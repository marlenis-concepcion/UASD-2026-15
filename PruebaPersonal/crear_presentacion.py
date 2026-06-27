from pathlib import Path
import sys

ROOT = Path(__file__).resolve().parent
sys.path.insert(0, str(ROOT / ".vendor"))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


OUT = ROOT / "Taller_Rediseno_Academico_IA.pptx"

NAVY = RGBColor(22, 67, 145)
BLUE = RGBColor(31, 78, 160)
CYAN = RGBColor(74, 210, 232)
PALE = RGBColor(227, 247, 251)
INK = RGBColor(25, 42, 67)
MUTED = RGBColor(88, 105, 127)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(244, 248, 252)
GREEN = RGBColor(53, 166, 116)
ORANGE = RGBColor(242, 161, 63)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


def rect(slide, x, y, w, h, fill, radius=False, line=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line if line else fill
    return shp


def circle(slide, x, y, d, fill, line=None, width=1.5):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line if line else fill
    shp.line.width = Pt(width)
    return shp


def text(slide, value, x, y, w, h, size=24, color=INK, bold=False,
         align=PP_ALIGN.LEFT, font="Aptos", valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = valign
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = value
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def bullets(slide, items, x, y, w, h, size=20, color=INK, gap=8):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for i, item in enumerate(items):
        p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(gap)
        p.text = "•  " + item
    return box


def header(slide, title, number, subtitle=None):
    rect(slide, 0, 0, 13.333, 0.12, CYAN)
    text(slide, f"{number:02d}", 0.55, 0.34, 0.7, 0.45, 14, BLUE, True)
    text(slide, title, 1.25, 0.28, 11.2, 0.62, 28, INK, True)
    if subtitle:
        text(slide, subtitle, 1.25, 0.88, 11.0, 0.42, 13, MUTED)
    text(slide, "Taller · Rediseño académico con IA", 9.65, 7.12, 3.05, 0.22, 9, MUTED, align=PP_ALIGN.RIGHT)


def card(slide, title_value, body, x, y, w, h, accent=CYAN, index=None):
    rect(slide, x, y, w, h, WHITE, True, RGBColor(220, 230, 239))
    rect(slide, x, y, 0.10, h, accent, True, accent)
    if index is not None:
        circle(slide, x + 0.28, y + 0.28, 0.55, accent)
        text(slide, str(index), x + 0.28, y + 0.29, 0.55, 0.52, 15, WHITE, True,
             PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        tx = x + 1.0
        tw = w - 1.25
    else:
        tx = x + 0.38
        tw = w - 0.65
    text(slide, title_value, tx, y + 0.28, tw, 0.38, 17, INK, True)
    text(slide, body, tx, y + 0.78, tw, h - 1.0, 12.5, MUTED)


# 1. Portada reconstruida a partir del afiche suministrado
slide = prs.slides.add_slide(blank)
rect(slide, 0, 0, 13.333, 7.5, LIGHT)
rect(slide, 0, 0, 7.8, 7.5, NAVY)
rect(slide, 0, 0, 7.8, 0.86, WHITE)
circle(slide, 0.55, 0.12, 0.62, PALE, CYAN, 2)
text(slide, "FIA", 0.55, 0.22, 0.62, 0.35, 15, BLUE, True, PP_ALIGN.CENTER)
text(slide, "Facultad de Ciencias", 1.35, 0.15, 4.4, 0.28, 15, INK, True)
text(slide, "Escuela de Informática", 1.35, 0.45, 4.4, 0.24, 13, INK, True)
rect(slide, 5.65, 0.16, 1.78, 0.48, CYAN, True)
text(slide, "27/JUN/2026", 5.73, 0.25, 1.61, 0.25, 11, INK, True, PP_ALIGN.CENTER)
text(slide, "TALLER", 0.65, 1.32, 4.8, 0.36, 18, CYAN, True)
text(slide, "Rediseño Académico\ncon IA", 0.65, 1.78, 6.25, 1.36, 34, WHITE, True)
text(slide, "Actualizar asignaturas según competencias\ny mercado laboral", 0.67, 3.22, 6.1, 0.82, 21, WHITE)
rect(slide, 0.65, 4.34, 1.75, 0.46, CYAN, True)
text(slide, "GRATIS", 0.76, 4.42, 1.52, 0.27, 16, NAVY, True, PP_ALIGN.CENTER)
text(slide, "MODALIDAD\nVIRTUAL", 0.68, 5.15, 1.72, 0.65, 13, WHITE, True)
text(slide, "HORARIO\nSábado · 08:00 AM – 12:00 PM", 2.58, 5.15, 3.75, 0.65, 13, WHITE, True)
text(slide, "DURACIÓN\n1 semana", 0.68, 6.03, 1.72, 0.60, 13, WHITE, True)
text(slide, "PROFESORA\nMarlenis Judith Concepción Cuevas", 2.58, 6.03, 4.42, 0.62, 13, WHITE, True)

# Decorative AI circles on right
for x, y, d, c in [(8.55, 0.85, 3.45, PALE), (9.25, 1.55, 2.05, WHITE), (10.05, 3.6, 2.45, WHITE)]:
    circle(slide, x, y, d, c, CYAN, 3)
text(slide, "IA", 9.25, 2.05, 2.05, 0.75, 34, BLUE, True, PP_ALIGN.CENTER)
text(slide, "DOCENCIA", 10.05, 4.35, 2.45, 0.42, 15, NAVY, True, PP_ALIGN.CENTER)
for r in range(4):
    for c in range(4):
        circle(slide, 8.25 + c * 0.34, 5.55 + r * 0.34, 0.10, CYAN)
text(slide, "Competencias · actividades · evaluación", 8.15, 6.92, 4.55, 0.28, 12, MUTED, True, PP_ALIGN.CENTER)


# 2. Propósito
slide = prs.slides.add_slide(blank)
rect(slide, 0, 0, 13.333, 7.5, LIGHT)
header(slide, "La IA redefine el trabajo docente", 2, "Una propuesta aplicable a un grupo multidisciplinario")
text(slide, "Más tiempo para enseñar.", 0.9, 1.62, 5.6, 0.52, 28, BLUE, True)
text(slide, "La IA reduce el trabajo repetitivo y amplía la capacidad para investigar, diseñar, evaluar y acompañar el aprendizaje.",
     0.9, 2.28, 5.55, 1.35, 19, INK)
card(slide, "Antes", "Horas para rediseñar actividades, revisar entregas y repetir correcciones.", 7.0, 1.55, 5.2, 1.25, CYAN, 1)
card(slide, "Con agentes", "Cada tarea docente recibe asistencia especializada y evidencia verificable.", 7.0, 3.0, 5.2, 1.25, BLUE, 2)
card(slide, "Nuevo rol", "El docente interpreta, acompaña, valida y toma la decisión final.", 7.0, 4.45, 5.2, 1.25, GREEN, 3)
rect(slide, 0.9, 4.25, 5.5, 1.25, PALE, True)
text(slide, "Idea clave", 1.2, 4.55, 1.3, 0.3, 14, BLUE, True)
text(slide, "La automatización devuelve al docente tiempo para acompañar al estudiante.", 1.2, 4.92, 4.75, 0.48, 17, INK, True)


# 3. Lo realizado
slide = prs.slides.add_slide(blank)
rect(slide, 0, 0, 13.333, 7.5, LIGHT)
header(slide, "¿Qué problema estamos resolviendo?", 3, "Demostración: NotebookLM con programas de asignaturas de programación")
steps = [
    ("1", "Actualizar", "Comparar programas, fuentes recientes y necesidades profesionales."),
    ("2", "Rediseñar", "Convertir competencias en actividades pertinentes y diferenciadas."),
    ("3", "Evaluar", "Revisar evidencias sin perder de vista criterios y rúbricas."),
    ("4", "Retroalimentar", "Explicar qué logró cada estudiante y cómo puede mejorar."),
]
for i, (n, ttl, body) in enumerate(steps):
    y = 1.52 + i * 1.25
    circle(slide, 0.95, y, 0.62, CYAN if i < 3 else GREEN)
    text(slide, n, 0.95, y + 0.04, 0.62, 0.5, 17, WHITE, True, PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    text(slide, ttl, 1.82, y - 0.02, 3.05, 0.34, 17, INK, True)
    text(slide, body, 1.82, y + 0.38, 4.55, 0.54, 13, MUTED)
    if i < 3:
        rect(slide, 1.23, y + 0.62, 0.06, 0.61, RGBColor(199, 223, 235))
rect(slide, 7.15, 1.55, 5.15, 4.95, NAVY, True)
text(slide, "NUEVO FLUJO", 7.55, 1.92, 4.35, 0.28, 12, CYAN, True)
text(slide, "Docente +\nagentes de IA", 7.55, 2.38, 4.35, 0.92, 27, WHITE, True)
text(slide, "Una tarea para cada agente", 7.55, 3.62, 4.0, 0.55, 17, WHITE, True)
for j, option in enumerate(["Investigar y actualizar", "Diseñar actividades", "Evaluar con rúbricas", "Retroalimentar y mejorar"]):
    fill = RGBColor(48, 91, 158) if j != 3 else GREEN
    rect(slide, 7.55, 4.32 + j * 0.46, 4.15, 0.35, fill, True)
    text(slide, option, 7.72, 4.38 + j * 0.46, 3.75, 0.19, 11, WHITE, j == 2)


# 4. Productos
slide = prs.slides.add_slide(blank)
rect(slide, 0, 0, 13.333, 7.5, LIGHT)
header(slide, "Agentes para las tareas docentes", 4, "Funciones adaptables a cualquier disciplina y contexto educativo")
products = [
    ("?", "Diagnóstico", "Detecta saberes previos y brechas", BLUE),
    ("↗", "Investigación", "Busca y verifica información actual", CYAN),
    ("✎", "Diseño", "Crea actividades y recursos alineados", ORANGE),
    ("✓", "Evaluación", "Contrasta evidencias con la rúbrica", GREEN),
    ("A+", "Calificación", "Propone resultados justificados", NAVY),
    ("AI", "Feedback", "Orienta la mejora de cada estudiante", RGBColor(118, 91, 166)),
]
for i, (icon, ttl, body, color) in enumerate(products):
    col, row = i % 3, i // 3
    x, y = 0.85 + col * 4.15, 1.55 + row * 2.35
    rect(slide, x, y, 3.72, 1.88, WHITE, True, RGBColor(220, 230, 239))
    circle(slide, x + 0.30, y + 0.32, 0.72, color)
    text(slide, icon, x + 0.30, y + 0.41, 0.72, 0.38, 14, WHITE, True, PP_ALIGN.CENTER)
    text(slide, ttl, x + 1.22, y + 0.32, 2.1, 0.35, 18, INK, True)
    text(slide, body, x + 1.22, y + 0.83, 2.1, 0.66, 13, MUTED)


# 5. Sistema de agentes
slide = prs.slides.add_slide(blank)
rect(slide, 0, 0, 13.333, 7.5, LIGHT)
header(slide, "Así funciona el equipo de agentes", 5, "Cada especialista recibe evidencia y entrega un producto verificable")
agents = [
    ("Diagnosticar", CYAN), ("Investigar", BLUE), ("Actualizar", NAVY),
    ("Diseñar", RGBColor(118, 91, 166)), ("Calificar", ORANGE), ("Retroalimentar", GREEN)
]
for i, (label, color) in enumerate(agents):
    x = 0.42 + i * 2.11
    circle(slide, x + 0.53, 2.0, 1.0, color)
    text(slide, str(i + 1), x + 0.53, 2.16, 1.0, 0.48, 20, WHITE, True, PP_ALIGN.CENTER)
    text(slide, label, x, 3.18, 2.05, 0.40, 14, INK, True, PP_ALIGN.CENTER)
    if i < len(agents) - 1:
        text(slide, "→", x + 1.75, 2.26, 0.42, 0.35, 21, MUTED, True, PP_ALIGN.CENTER)
rect(slide, 1.15, 4.35, 11.0, 1.20, PALE, True)
text(slide, "AGENTE COORDINADOR", 1.55, 4.62, 2.65, 0.30, 13, BLUE, True)
text(slide, "Conserva el contexto, verifica la alineación y solicita aprobación humana.", 4.05, 4.55, 7.55, 0.52, 18, INK, True)
text(slide, "El docente conserva la decisión final.", 4.05, 5.06, 7.2, 0.30, 13, MUTED)


# 6. Marco metodológico
slide = prs.slides.add_slide(blank)
rect(slide, 0, 0, 13.333, 7.5, LIGHT)
header(slide, "Los agentes obedecen a la pedagogía", 6, "La velocidad solo tiene valor cuando existe intención educativa")
frameworks = [
    ("Bloom", "Resultados observables: recordar, comprender, aplicar, analizar, evaluar y crear.", BLUE),
    ("Aula invertida", "Antes: explorar · Durante: aplicar · Después: consolidar y mejorar.", CYAN),
    ("Agile docente", "Ciclos cortos, producto verificable, revisión, retrospectiva y ajuste.", ORANGE),
    ("IA educativa", "Propósito explícito, verificación, transparencia y supervisión humana.", GREEN),
]
for i, (ttl, body, color) in enumerate(frameworks):
    col, row = i % 2, i // 2
    x, y = 0.95 + col * 6.15, 1.55 + row * 2.35
    card(slide, ttl, body, x, y, 5.35, 1.85, color, i + 1)


# 7. Diagnóstico
slide = prs.slides.add_slide(blank)
rect(slide, 0, 0, 13.333, 7.5, LIGHT)
header(slide, "El diagnóstico activa el primer agente", 7, "Diagnosticar para adaptar; no para castigar")
text(slide, "¿Qué necesitamos descubrir?", 0.85, 1.55, 4.85, 0.42, 22, INK, True)
bullets(slide, [
    "Conocimientos y habilidades previas",
    "Conceptos erróneos o vacíos",
    "Nivel inicial según Bloom",
    "Necesidades de apoyo y agrupamiento",
], 0.9, 2.12, 5.05, 2.75, 17, INK, 12)
rect(slide, 6.35, 1.50, 5.95, 4.60, WHITE, True, RGBColor(220, 230, 239))
text(slide, "RESULTADO DEL DIAGNÓSTICO", 6.78, 1.90, 5.10, 0.28, 12, BLUE, True)
for i, (label, pct, color) in enumerate([
    ("Recordar", 78, CYAN), ("Comprender", 62, BLUE), ("Aplicar", 45, ORANGE), ("Analizar", 28, GREEN)
]):
    y = 2.55 + i * 0.77
    text(slide, label, 6.78, y, 1.18, 0.25, 12, MUTED)
    rect(slide, 8.10, y + 0.02, 3.28, 0.22, RGBColor(229, 235, 241), True)
    rect(slide, 8.10, y + 0.02, 3.28 * pct / 100, 0.22, color, True)
    text(slide, f"{pct}%", 11.48, y - 0.01, 0.46, 0.25, 11, INK, True, PP_ALIGN.RIGHT)
text(slide, "Decisión docente: reforzar aplicación antes del primer reto.", 6.78, 5.52, 4.96, 0.35, 14, INK, True)
rect(slide, 0.9, 5.42, 4.92, 0.78, PALE, True)
text(slide, "No se utiliza como calificación sumativa.", 1.22, 5.67, 4.28, 0.28, 14, BLUE, True, PP_ALIGN.CENTER)


# 8. Aula invertida
slide = prs.slides.add_slide(blank)
rect(slide, 0, 0, 13.333, 7.5, LIGHT)
header(slide, "La IA acompaña todo el ciclo", 8, "El aula invertida libera el encuentro para aplicar, dialogar y crear")
phases = [
    ("ANTES", "Explorar", "Video breve, lectura guiada, tarjetas y diagnóstico.", CYAN),
    ("DURANTE", "Aplicar", "Casos, retos, colaboración y acompañamiento docente.", BLUE),
    ("DESPUÉS", "Mejorar", "Producto, reflexión, retroalimentación y nueva versión.", GREEN),
]
for i, (phase, verb, body, color) in enumerate(phases):
    x = 0.72 + i * 4.22
    rect(slide, x, 1.65, 3.72, 3.92, WHITE, True, RGBColor(220, 230, 239))
    rect(slide, x, 1.65, 3.72, 0.62, color, True)
    text(slide, phase, x + 0.25, 1.82, 3.22, 0.25, 13, WHITE, True, PP_ALIGN.CENTER)
    text(slide, verb, x + 0.35, 2.65, 3.02, 0.48, 23, INK, True, PP_ALIGN.CENTER)
    text(slide, body, x + 0.45, 3.45, 2.82, 1.12, 15, MUTED, align=PP_ALIGN.CENTER)
    if i < 2:
        text(slide, "→", x + 3.72, 3.25, 0.50, 0.40, 22, BLUE, True, PP_ALIGN.CENTER)
text(slide, "El estudiante llega preparado para trabajar; el docente observa, pregunta y retroalimenta.",
     1.3, 6.18, 10.75, 0.45, 17, BLUE, True, PP_ALIGN.CENTER)


# 9. Alineación
slide = prs.slides.add_slide(blank)
rect(slide, 0, 0, 13.333, 7.5, LIGHT)
header(slide, "Evaluar sin sofocarse", 9, "La IA organiza la evidencia; el docente interpreta y decide")
labels = ["Diagnóstico", "Competencia", "Resultado", "Bloom", "Actividad", "Evidencia", "Rúbrica", "Feedback"]
colors = [CYAN, BLUE, NAVY, RGBColor(118, 91, 166), ORANGE, RGBColor(209, 118, 82), GREEN, RGBColor(54, 145, 137)]
for i, (label, color) in enumerate(zip(labels, colors)):
    row, col = i // 4, i % 4
    x, y = 0.62 + col * 3.12, 1.62 + row * 2.05
    rect(slide, x, y, 2.48, 1.15, color, True)
    text(slide, label, x + 0.15, y + 0.38, 2.18, 0.34, 15, WHITE, True, PP_ALIGN.CENTER)
    if col < 3:
        text(slide, "→", x + 2.53, y + 0.40, 0.42, 0.32, 18, MUTED, True, PP_ALIGN.CENTER)
    elif row == 0:
        text(slide, "↓", x + 1.03, y + 1.23, 0.42, 0.32, 18, MUTED, True, PP_ALIGN.CENTER)
text(slide, "Regla de control", 0.85, 5.85, 2.05, 0.30, 14, BLUE, True)
text(slide, "Si un criterio no se cumplió, el agente identifica la evidencia y propone cómo mejorarlo.",
     2.88, 5.77, 9.25, 0.52, 19, INK, True)


# 10. Cierre
slide = prs.slides.add_slide(blank)
rect(slide, 0, 0, 13.333, 7.5, NAVY)
rect(slide, 0, 0, 13.333, 0.14, CYAN)
text(slide, "No se trata de hacer más.\nSe trata de enseñar mejor.", 0.95, 1.15, 7.25, 1.38, 33, WHITE, True)
text(slide, "La IA realiza el trabajo repetitivo.\nEl docente aporta propósito, vínculo y criterio.", 0.98, 2.92, 6.95, 1.05, 21, PALE)
rect(slide, 0.98, 4.48, 5.78, 0.08, CYAN)
text(slide, "AGENTES ESPECIALIZADOS  ·  HUMAN IN THE LOOP", 0.98, 4.82, 7.6, 0.32, 13, CYAN, True)
circle(slide, 9.13, 1.42, 2.85, PALE, CYAN, 3)
text(slide, "IA", 9.13, 2.12, 2.85, 0.75, 40, BLUE, True, PP_ALIGN.CENTER)
for r in range(4):
    for c in range(4):
        circle(slide, 9.35 + c * 0.46, 5.12 + r * 0.38, 0.12, CYAN)
text(slide, "Marlenis Judith Concepción Cuevas", 0.98, 6.68, 5.5, 0.28, 12, WHITE, True)
text(slide, "Facultad de Ciencias · Escuela de Informática", 7.55, 6.68, 4.8, 0.28, 12, WHITE, align=PP_ALIGN.RIGHT)


prs.core_properties.title = "Taller: Rediseño académico con IA"
prs.core_properties.subject = "Ejercicio con NotebookLM, agentes docentes y evaluación"
prs.core_properties.author = "Marlenis Judith Concepción Cuevas"
prs.core_properties.keywords = "IA, docencia, Bloom, aula invertida, Agile, NotebookLM"

# Notas del presentador. Se muestran en la vista del moderador de PowerPoint.
speaker_notes = [
    """Buenos días. Mi nombre es Marlenis Judith Concepción Cuevas y les doy la bienvenida a este taller sobre rediseño académico con inteligencia artificial. Hoy no vamos a mirar la IA solamente como una herramienta para escribir textos o preparar una presentación. Vamos a pensarla como una oportunidad para reorganizar el trabajo docente completo.

Tenemos un grupo multidisciplinario. Esto es una fortaleza, porque la necesidad de simplificar el trabajo cotidiano no pertenece a una sola carrera. Una persona puede enseñar programación, otra salud, derecho, ingeniería, idiomas, ciencias sociales o administración; todas deben investigar, planificar, crear actividades, evaluar evidencias y retroalimentar. Cambian los contenidos y los criterios profesionales, pero muchas tareas del flujo docente son comunes.

Nuestro propósito será partir de una asignatura real, revisar sus contenidos, relacionarla con las competencias profesionales y producir actividades, evaluaciones y retroalimentaciones más pertinentes. Trabajaremos con la Taxonomía de Bloom revisada, aula invertida y ciclos Agile para docentes.

También presentaré NotebookLM utilizando varios programas de asignaturas de programación. Veremos cómo una colección de fuentes institucionales puede convertirse en preguntas, cuestionarios, presentaciones y materiales de apoyo. Aunque el ejemplo sea de programación, el procedimiento puede transferirse a cualquier disciplina sustituyendo las fuentes y los criterios correspondientes.

Mi intención es que cada persona salga de aquí con al menos una idea concreta para simplificar su trabajo del día a día. Finalmente, conectaremos la demostración con una arquitectura de agentes especializados. La idea central es sencilla: la IA realiza y organiza buena parte del trabajo repetitivo, mientras el docente conserva el propósito pedagógico, el vínculo humano y la decisión final.""",

    """Quiero comenzar con una afirmación: la inteligencia artificial está redefiniendo el trabajo docente. No se trata simplemente de agregar otra aplicación a nuestra lista de herramientas. Se trata de revisar cómo distribuimos nuestro tiempo y cuáles tareas requieren realmente nuestro juicio profesional.

Durante años hemos dedicado muchas horas a comenzar actividades desde una página en blanco, adaptar instrucciones, elaborar preguntas similares, revisar entregas una por una y repetir comentarios. Con IA, buena parte de ese trabajo puede prepararse, clasificarse o compararse de manera asistida.

Eso no significa sacar al docente del proceso. Significa moverlo hacia donde aporta más valor: definir qué se debe aprender, observar cómo aprende el estudiante, interpretar dificultades, acompañar decisiones y validar la evaluación. A este principio lo llamamos **human in the loop**, o ser humano dentro del ciclo. La IA ejecuta una tarea, organiza datos o propone una respuesta; el docente revisa la evidencia, corrige cuando sea necesario y aprueba el resultado antes de utilizarlo.

Los agentes que propondremos no tomarán el control de la asignatura. Cada uno realizará una función delimitada y entregará evidencia para que el docente revise. El verdadero beneficio no es solamente producir más rápido; es recuperar tiempo para enseñar, conversar y atender las diferencias del grupo. Al escuchar los ejemplos, quiero que cada participante piense: ¿cuál tarea repetitiva de mi día a día podría delegar parcialmente a un agente sin delegar mi responsabilidad?""",

    """Aquí planteamos el problema que queremos resolver. Actualizar una asignatura exige revisar programas, fuentes nuevas y necesidades del mercado laboral. Luego debemos transformar esos hallazgos en actividades coherentes, revisar las evidencias y explicar a cada estudiante cómo mejorar. Cuando el grupo es numeroso, este ciclo puede convertirse en una carga difícil de sostener.

En este momento haré la demostración de NotebookLM. He cargado varios programas relacionados con asignaturas de programación, entre ellos documentos de INF-510 a INF-515, incluyendo asignaturas y laboratorios, además de planes académicos. Es importante señalar que NotebookLM responde apoyándose en las fuentes seleccionadas. Por eso primero debemos revisar cuáles documentos cargamos y cuáles tenemos activados.

[Mostrar NotebookLM]. Señalaré el panel de fuentes, el chat y el área Studio. A partir de esos programas podemos preguntar qué contenidos se repiten, cuáles competencias aparecen, qué temas podrían estar desactualizados y cómo se relacionan las asignaturas. También podemos producir un cuestionario sobre programación, una presentación, tarjetas didácticas o una guía. La herramienta genera un borrador fundamentado; el docente verifica y decide si es pedagógicamente adecuado.""",

    """En lugar de pedirle todo a una sola IA con un único mensaje, proponemos agentes especializados. El agente de diagnóstico identifica conocimientos previos y brechas. El investigador busca información actual y verifica sus fuentes. El diseñador convierte resultados de aprendizaje en actividades. El evaluador compara evidencias con una rúbrica. El agente de calificación propone una valoración justificada y el de retroalimentación explica cómo mejorar.

Estos agentes son útiles para un grupo multidisciplinario porque la arquitectura se conserva aunque cambie la materia. En programación, un agente puede analizar código y criterios de funcionamiento. En salud, puede organizar un caso clínico educativo sin utilizar datos personales. En derecho, puede contrastar una argumentación con los criterios de la actividad. En idiomas, puede clasificar errores y proponer práctica diferenciada. En administración, puede revisar un estudio de caso o un proyecto. En todos los ejemplos, las fuentes, la rúbrica y la validación humana deben pertenecer a la disciplina.

Esta separación es importante porque cada tarea requiere instrucciones, fuentes y controles distintos. Un agente que investiga no debe inventar criterios de evaluación. Un agente que califica no debe cambiar la rúbrica después de ver el trabajo del estudiante. Y un agente de retroalimentación no debe limitarse a felicitar o señalar errores; debe ofrecer acciones concretas.

En la demostración de NotebookLM vimos productos como preguntas, cuestionarios y presentaciones. Aquí damos un paso adicional: convertimos esos productos en partes de un flujo docente. Cada salida debe tener una finalidad pedagógica y un responsable de validarla. Así evitamos acumular materiales bonitos que no contribuyen realmente al aprendizaje.""",

    """Este es el flujo general del equipo de agentes. Primero se diagnostica. Después se investiga qué debe actualizarse. Con la evidencia recopilada se actualizan los contenidos. Solo entonces se diseñan actividades, instrumentos y productos. Cuando el estudiante entrega su evidencia, el agente evaluador la contrasta con la rúbrica y el agente de retroalimentación prepara orientaciones para la siguiente versión.

Sobre todos ellos se encuentra el agente coordinador. Su función no es enseñar ni calificar, sino conservar el contexto, comprobar que cada etapa recibió los productos anteriores y detener el proceso cuando falta información. También identifica los puntos donde debe intervenir el docente.

Por ejemplo, si la rúbrica pide argumentación, pero la actividad solo solicita marcar respuestas, el coordinador debe señalar que existe una desalineación. Si un estudiante no cumplió un criterio, el sistema no debe improvisar un castigo ni cambiar la regla. Debe mostrar qué evidencia falta, proponer una retroalimentación y remitir la decisión definitiva al docente.

Aquí opera nuevamente el principio human in the loop. No significa que el docente tenga que rehacer manualmente todo lo producido por la IA. Significa que existen puntos de control proporcionales al riesgo: aprobar las competencias, validar las fuentes, revisar la rúbrica, examinar casos dudosos y confirmar la calificación. Los agentes realizan el trabajo operativo; el docente supervisa y responde por la decisión. Cada agente asiste; ninguno sustituye la responsabilidad profesional.""",

    """La rapidez de la IA solo es valiosa cuando está gobernada por una metodología. Nuestro primer marco es la Taxonomía de Bloom revisada. Los resultados deben utilizar acciones observables: recordar, comprender, aplicar, analizar, evaluar y crear. No basta con decir que el estudiante conocerá un tema; debemos determinar qué hará para demostrarlo.

El segundo marco es el aula invertida. Antes del encuentro, el estudiante explora fuentes, responde un diagnóstico y prepara preguntas. Durante la clase aplicamos, analizamos y resolvemos problemas. Después de la clase consolidamos, producimos una evidencia y mejoramos a partir de la retroalimentación.

El tercer marco es Agile para docentes: trabajamos en ciclos cortos, con una lista priorizada, un producto verificable, una revisión y una retrospectiva. Finalmente, la IA educativa exige transparencia, verificación de fuentes, protección de datos y supervisión humana. Los agentes no diseñan en el vacío. Sus decisiones deben obedecer a estos principios y mantener la alineación entre lo que prometemos enseñar y lo que realmente evaluamos.""",

    """Toda unidad comienza con una prueba diagnóstica. Su finalidad no es castigar ni asignar una nota sumativa. Buscamos conocer el punto de partida: qué recuerda el estudiante, qué comprende, qué puede aplicar y cuáles conceptos erróneos podrían dificultar el nuevo aprendizaje.

El agente de diagnóstico puede ayudarnos a crear preguntas en varios niveles de Bloom, organizar las respuestas y detectar patrones. En el ejemplo de la gráfica, el grupo tiene mejores resultados en recordar y comprender, pero presenta mayor dificultad al aplicar y analizar. Esa información cambia nuestra planificación: quizá debamos dedicar menos tiempo a repetir definiciones y más tiempo a casos guiados.

Aquí también podemos utilizar NotebookLM para generar un primer banco de preguntas basadas en los programas y materiales seleccionados. Sin embargo, el docente debe comprobar el nivel cognitivo, la claridad y la pertinencia de cada pregunta. Los resultados del diagnóstico alimentan el primer ciclo Agile: permiten priorizar contenidos, formar grupos, preparar apoyos y ajustar el ritmo. Al final compararemos una evidencia equivalente para observar el progreso real.""",

    """La IA puede acompañar las tres etapas del aula invertida. Antes de la clase, puede preparar una lectura guiada, un video breve, tarjetas, preguntas de anticipación y el diagnóstico. El estudiante llega con una primera aproximación al contenido y el docente recibe información sobre las dificultades del grupo.

Durante la clase no queremos usar el tiempo para repetir lo que el estudiante pudo revisar previamente. Queremos resolver casos, programar, analizar errores, comparar soluciones y tomar decisiones. La IA puede proponer variantes de un problema o brindar apoyos diferenciados, pero el encuentro debe favorecer la discusión, la colaboración y el acompañamiento docente.

Después de la clase, el estudiante entrega un producto, explica su proceso y recibe retroalimentación. El agente evaluador organiza la evidencia según la rúbrica y el agente de retroalimentación propone próximos pasos. El estudiante mejora su trabajo y el docente observa qué debe cambiar en el siguiente ciclo. De esta forma, la tecnología no reduce la interacción humana; libera tiempo del encuentro para hacerla más profunda.""",

    """Una de las tareas que más agota al docente es corregir trabajos cuando algunos estudiantes no cumplen la rúbrica. El problema no es solamente colocar una puntuación. Debemos identificar qué criterio se cumplió, cuál evidencia lo demuestra, qué falta y cómo puede mejorar el estudiante.

Nuestro flujo conecta diagnóstico, competencia, resultado de aprendizaje, nivel de Bloom, actividad, evidencia, rúbrica y retroalimentación. El agente evaluador compara la entrega con cada criterio aprobado. No puede penalizar algo que nunca fue solicitado ni modificar los criterios después de recibir el trabajo. Cuando no encuentra evidencia suficiente, debe decirlo claramente y remitir el caso al docente.

El agente puede preparar una propuesta como esta: criterio, evidencia encontrada, nivel alcanzado, puntuación sugerida y justificación. Luego el agente de retroalimentación convierte ese análisis en un mensaje respetuoso y accionable. El docente revisa los casos dudosos y valida la calificación. Así dejamos de repetir comentarios manualmente, pero mantenemos una evaluación justa, trazable y humana.""",

    """Para cerrar, no estamos proponiendo que el docente haga más tareas con una herramienta nueva. Estamos proponiendo una manera distinta de trabajar. Los agentes especializados pueden investigar, organizar información, producir variantes, comparar evidencias y preparar retroalimentaciones. El docente aporta lo que ninguna automatización debe decidir por sí sola: el propósito de la asignatura, la lectura del contexto, la relación con el estudiante y el juicio profesional.

Como somos un grupo multidisciplinario, no espero que todos salgamos con el mismo agente. Una persona puede comenzar con un agente para crear actividades; otra, con uno para revisar una rúbrica; otra, con un agente para organizar retroalimentaciones frecuentes. Lo importante es identificar una tarea real que consume tiempo, definir qué información necesita la IA, establecer qué producto debe entregar y decidir dónde intervendrá el ser humano para validarlo.

La invitación es comenzar con una asignatura y una unidad, no automatizarlo todo de inmediato. Cargamos las fuentes confiables, aplicamos un diagnóstico, seleccionamos una tarea repetitiva y diseñamos un agente con límites claros. Revisamos sus resultados, recogemos evidencia y mejoramos el siguiente ciclo.

Quiero que nos quedemos con esta idea: la IA puede hacer por nosotros gran parte del trabajo operativo y repetitivo, pero trabaja dentro de límites definidos por el docente. Nosotros validamos los resultados mediante human in the loop. La meta es que cada participante salga con una idea aplicable para simplificar su día a día y recuperar tiempo para acompañar mejor a sus estudiantes.

Nuestro próximo paso será tomar una asignatura real y construir esta cadena de agentes desde la investigación hasta la retroalimentación. Muchas gracias. Ahora podemos conversar sobre cuál tarea docente desean transformar primero.""",
]

if len(speaker_notes) != len(prs.slides):
    raise ValueError("Debe existir una nota del presentador por cada diapositiva")

for slide, note in zip(prs.slides, speaker_notes):
    slide.notes_slide.notes_text_frame.text = note

prs.save(OUT)

# Copia de respaldo del guion para leerlo fuera de PowerPoint.
note_titles = [
    "Portada e introducción",
    "La IA redefine el trabajo docente",
    "Problema y demostración de NotebookLM",
    "Agentes para las tareas docentes",
    "Funcionamiento del equipo de agentes",
    "Marco metodológico",
    "Prueba diagnóstica",
    "Aula invertida e IA",
    "Evaluación, rúbrica y retroalimentación",
    "Cierre",
]
guide = ["# Guion de presentación", ""]
for number, (title_value, note) in enumerate(zip(note_titles, speaker_notes), start=1):
    guide.extend([f"## Diapositiva {number}: {title_value}", "", note, ""])
(ROOT / "GUION_PRESENTACION.md").write_text("\n".join(guide), encoding="utf-8")

print(OUT)
