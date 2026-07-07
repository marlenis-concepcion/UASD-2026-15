from pathlib import Path
import sys

ROOT = Path(__file__).resolve().parent
VENDOR = ROOT.parent / "TallerIA-27-06-2026" / "PruebaPersonal" / ".vendor"
sys.path.insert(0, str(VENDOR))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

OUT = ROOT / "Testing4All_Transformacion_QA_IA_Empleabilidad_2026.pptx"
LOGO = ROOT / "assets" / "testing4all_logo.png"
UASD = ROOT / "assets" / "logo_uasd.png"
REPO_QR = ROOT / "assets" / "qr_repositorio_recursos.png"
COMIC = ROOT / "assets" / "comic_nuevo_profesional_presionado.png"
COMIC_EXPECTATIONS = ROOT / "assets" / "comic_expectativas_hombre_mujer_presion.png"
ALLEGORY_SPRITE = ROOT / "assets" / "munequitos_cimientos_herramientas_personas.png"
CAREER_SPRITE = ROOT / "assets" / "munequitos_carreras_qa_movimiento_eleccion.png"
QR_BRAIN_DRAIN = ROOT / "assets" / "qr_fuga_cerebros_diariolibre.png"
QR_INFORMALITY = ROOT / "assets" / "qr_informalidad_banco_central.png"
QR_LINKEDIN = ROOT / "assets" / "qr_linkedin_marlenis.png"
QR_CIC = ROOT / "assets" / "qr_congreso_xxi_cic_2026.png"
AUTOMATION_SPRITE = ROOT / "assets" / "munequitos_automatizacion_valor_diseno_mantenimiento.png"
PRIORITY_SPRITE = ROOT / "assets" / "munequitos_priorizar_reconocer_preparar_cambiar.png"
PRESSURE_SPRITE = ROOT / "assets" / "munequitos_presion_desorden_recursos.png"
LEGAL_CIRCLE = ROOT / "assets" / "circulo_leyes_republica_dominicana.png"
RESOURCES_REPO = "https://github.com/marlenis-concepcion/TallerIA-Docente-27-06-2026"

NAVY = RGBColor(5, 13, 42)
PURPLE = RGBColor(93, 45, 145)
VIOLET = RGBColor(160, 77, 224)
CYAN = RGBColor(47, 210, 224)
ORANGE = RGBColor(255, 145, 55)
GREEN = RGBColor(68, 198, 135)
WHITE = RGBColor(255, 255, 255)
INK = RGBColor(25, 31, 50)
MUTED = RGBColor(94, 103, 126)
LIGHT = RGBColor(247, 247, 252)
LINE = RGBColor(222, 224, 235)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


def rect(slide, x, y, w, h, fill, radius=False, line=None):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line or fill
    return shp


def txt(slide, value, x, y, w, h, size=22, color=INK, bold=False,
        align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, italic=False, url=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.clear(); tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(.04)
    tf.margin_top = tf.margin_bottom = Inches(.03)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = value
    r.font.name = "Aptos"; r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    if url:
        r.hyperlink.address = url
    return box


def notes(slide, value):
    slide.notes_slide.notes_text_frame.text = value.strip()


def brand(slide, dark=False):
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(11.72), Inches(.22), width=Inches(1.05))
    txt(slide, "TESTING4ALL · 2026", .48, 7.12, 2.4, .18, 8,
        WHITE if dark else MUTED, True)


def base(title, number, subtitle=None):
    slide = prs.slides.add_slide(blank)
    rect(slide, 0, 0, 13.333, 7.5, LIGHT)
    rect(slide, 0, 0, 13.333, .10, VIOLET)
    txt(slide, f"{number:02d}", .48, .35, .5, .3, 11, PURPLE, True)
    # Tipografía preparada para proyección y lectura a larga distancia.
    txt(slide, title, 1.05, .23, 10.3, .66, 32, INK, True)
    if subtitle:
        txt(slide, subtitle, 1.05, .91, 10.8, .42, 18, MUTED)
    brand(slide)
    txt(slide, f"{number:02d}", 12.36, 7.10, .45, .2, 8, MUTED, True, PP_ALIGN.RIGHT)
    return slide


def bullets(slide, items, x, y, w, h, size=22, color=INK, accent=CYAN):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.clear(); tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(.02)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•  {item}"; p.font.name = "Aptos"; p.font.size = Pt(size)
        p.font.color.rgb = color; p.space_after = Pt(12); p.line_spacing = 1.05
    rect(slide, x-.32, y+.02, .10, min(h, 4.8), accent, True)
    return box


def cards(slide, data, top=1.55, cols=3):
    gap = .25; left = .62; total = 12.1
    w = (total-gap*(cols-1))/cols
    rows = (len(data)+cols-1)//cols
    h = 4.85 if rows == 1 else 2.28
    for i, (title, body, color) in enumerate(data):
        row, col = divmod(i, cols); x = left+col*(w+gap); y = top+row*(h+.25)
        rect(slide, x, y, w, h, WHITE, True, LINE)
        rect(slide, x, y, .10, h, color, True, color)
        # Las composiciones 2×2 funcionan como tablas/resúmenes densos y se
        # conservan intactas. Las tarjetas de una fila reciben letra de sala.
        title_size = 20 if rows == 1 else 15.5
        body_size = 18 if rows == 1 else 12.5
        txt(slide, title, x+.28, y+.24, w-.52, .52, title_size, INK, True)
        txt(slide, body, x+.28, y+.82, w-.52, h-1.0, body_size, MUTED)


def book_idea(slide, value, y=6.40):
    """Paráfrasis breve visible; la referencia completa aparece en la lámina 26."""
    rect(slide, .82, y, 11.70, .58, RGBColor(238, 232, 249), True,
         RGBColor(218, 204, 240))
    txt(slide, value, 1.02, y+.08, 11.30, .42, 15.5, PURPLE, True,
        PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)


def resource(slide, y, label, description, url, color=PURPLE):
    rect(slide, .70, y, 11.92, .70, WHITE, True, LINE)
    rect(slide, .70, y, .10, .70, color, True, color)
    txt(slide, label, .98, y+.12, 3.0, .28, 13, color, True, url=url)
    txt(slide, description, 4.05, y+.12, 7.95, .34, 11.5, MUTED, url=url)
    txt(slide, "↗", 12.08, y+.14, .28, .28, 14, color, True, url=url)


def repo_qr_badge(slide, x=11.35, y=.78, size=.95, dark=False):
    """Acceso alternativo cuando PowerPoint bloquea los hipervínculos."""
    if REPO_QR.exists():
        rect(slide, x-.05, y-.05, size+.10, size+.10, WHITE, True, WHITE)
        slide.shapes.add_picture(str(REPO_QR), Inches(x), Inches(y), width=Inches(size))
    txt(slide, "RECURSOS", x-.12, y+size+.04, size+.24, .18, 7.5,
        CYAN if dark else PURPLE, True, PP_ALIGN.CENTER)


def mini_pyramid(slide, x, y, w, title, levels, colors):
    """Pirámide vectorial de cuatro niveles, de base a cima."""
    rect(slide, x, y, w, 2.35, WHITE, True, LINE)
    txt(slide, title.upper(), x+.18, y+.15, w-.36, .28, 12.5, INK, True, PP_ALIGN.CENTER)
    base_y = y + 1.93
    widths = [w-.38, w-.92, w-1.46, w-2.00]
    for i, (label, color) in enumerate(zip(levels, colors)):
        level_w = widths[i]
        level_x = x + (w-level_w)/2
        level_y = base_y - i*.42
        rect(slide, level_x, level_y, level_w, .36, color, True, color)
        txt(slide, label, level_x+.05, level_y+.06, level_w-.10, .20,
            8.2, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)


def emoji_badge(slide, symbol, x, y, color, d=.52):
    """Icono emotivo contenido dentro de una tarjeta."""
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.color.rgb = color
    box = slide.shapes.add_textbox(Inches(x), Inches(y+.01), Inches(d), Inches(d-.02))
    tf = box.text_frame; tf.clear(); tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = symbol
    r.font.name = "Apple Color Emoji"; r.font.size = Pt(17)
    r.font.color.rgb = WHITE
    return box


def emoji_scene(slide, symbols, x, y, w, color):
    """Pequeña escena alegórica centrada dentro de una tarjeta."""
    bg = rect(slide, x, y, w, 1.18, RGBColor(246, 242, 252), True, color)
    bg.line.width = Pt(1.3)
    box = slide.shapes.add_textbox(Inches(x+.10), Inches(y+.10), Inches(w-.20), Inches(.82))
    tf = box.text_frame; tf.clear(); tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = symbols
    r.font.name = "Apple Color Emoji"; r.font.size = Pt(34)
    return box


def sprite_third(slide, third, x, y, w=2.55, h=2.75, image_path=None):
    """Inserta un tercio del sprite como imagen real, sin depender de emojis."""
    image_path = image_path or ALLEGORY_SPRITE
    if not image_path.exists():
        return None
    pic = slide.shapes.add_picture(str(image_path), Inches(x), Inches(y),
                                   width=Inches(w), height=Inches(h))
    if third == 0:
        pic.crop_right = 2/3
    elif third == 1:
        pic.crop_left = 1/3; pic.crop_right = 1/3
    else:
        pic.crop_left = 2/3
    return pic


def photo_half(slide, half, x, y, w, h, image_path):
    """Inserta una mitad de una fotografía doble dentro de una tarjeta."""
    if not image_path.exists():
        return None
    pic = slide.shapes.add_picture(str(image_path), Inches(x), Inches(y),
                                   width=Inches(w), height=Inches(h))
    if half == 0:
        pic.crop_right = .5
    else:
        pic.crop_left = .5
    return pic


def pill(slide, value, x, y, w, color):
    rect(slide, x, y, w, .38, color, True, color)
    txt(slide, value, x+.06, y+.08, w-.12, .18, 9, WHITE, True, PP_ALIGN.CENTER)


# 1 · Portada
s = prs.slides.add_slide(blank)
rect(s, 0, 0, 13.333, 7.5, NAVY)
rect(s, 8.75, 0, 4.583, 7.5, PURPLE)
rect(s, 8.93, .28, 4.10, 6.95, NAVY, True, VIOLET)
if LOGO.exists(): s.shapes.add_picture(str(LOGO), Inches(9.48), Inches(.72), width=Inches(3.0))
if QR_LINKEDIN.exists():
    rect(s, 10.15, 4.18, 1.55, 1.55, WHITE, True, WHITE)
    s.shapes.add_picture(str(QR_LINKEDIN), Inches(10.25), Inches(4.28), width=Inches(1.35))
txt(s, "CONECTA CONMIGO", 9.55, 5.88, 2.75, .28, 11, CYAN, True, PP_ALIGN.CENTER)
txt(s, "Escanea para agregarme en LinkedIn", 9.35, 6.18, 3.15, .42, 9.5, WHITE, False, PP_ALIGN.CENTER,
    url="https://www.linkedin.com/in/marlenis-judith-c-55f1117a3/")
txt(s, "SOFTWARE QA ENGINEERING", .72, .75, 6.8, .28, 13, CYAN, True)
txt(s, "IA, automatización\ny empleabilidad", .72, 1.35, 7.45, 1.55, 34, WHITE, True)
txt(s, "Cómo evolucionar en un mercado que cambia más rápido que el currículo", .75, 3.20, 7.15, .85, 20, RGBColor(222,225,245), True)
rect(s, .75, 4.50, 6.80, .06, VIOLET)
txt(s, "Mtra. Marlenis Judith Concepción Cuevas", .75, 4.88, 7.0, .32, 15, WHITE, True)
txt(s, "Senior SDET · Docente e investigadora · UASD", .75, 5.30, 6.8, .28, 12, CYAN)
txt(s, "linkedin.com/in/marlenis-judith-c-55f1117a3/", .75, 5.76, 6.9, .30, 11, WHITE,
    url="https://www.linkedin.com/in/marlenis-judith-c-55f1117a3/")
if UASD.exists(): s.shapes.add_picture(str(UASD), Inches(.76), Inches(6.40), width=Inches(.68))
notes(s, "Abra con la pregunta: ¿qué habilidad de tu puesto cambió durante el último año? Evite una introducción biográfica larga.")

# 2
s = base("Lo que conecta estas conversaciones", 2, "Agentes, IA, automatización y observabilidad dependen de bases humanas y técnicas")
cards(s, [
    ("Primero los cimientos", "No hay agentes confiables ni buena automatización sin conceptos básicos, contexto y diseño.", VIOLET),
    ("Después la herramienta", "IA, observabilidad y automatización amplifican lo que el equipo ya entiende y practica.", CYAN),
    ("Siempre las personas", "La calidad se construye conversando, preguntando, negociando riesgo y creando confianza.", ORANGE),
])
sprite_third(s, 0, 1.46, 3.98, 2.15, 2.32)
sprite_third(s, 1, 5.58, 3.98, 2.15, 2.32)
sprite_third(s, 2, 9.69, 3.98, 2.15, 2.32)
book_idea(s, "El comportamiento laboral surge de la persona, el grupo y el sistema (Robbins y Judge, 2017).")
notes(s, "Conecte con los temas de las demás charlas: agentes autónomos, automatización, observabilidad, IA y confianza. La tecnología cambia, pero todas las propuestas necesitan fundamentos y relaciones humanas saludables.")

# 3
s = base("El currículo corre detrás", 3, "Una idea retomada del XXI Congreso Científico Internacional")
bullets(s, [
    "El mercado y las herramientas cambian más rápido que los programas académicos.",
    "Varias asignaturas pueden repetir contenidos mientras faltan competencias emergentes.",
    "Actualizar una regulación o un plan de estudios requiere acuerdos, procesos y tiempo.",
    "La universidad aporta fundamentos; la vigencia profesional exige aprendizaje continuo."
], 1.0, 1.60, 7.3, 4.8, 18)
rect(s, 8.78, 1.60, 3.75, 4.60, NAVY, True)
txt(s, "La brecha", 9.18, 2.05, 2.9, .35, 14, CYAN, True)
txt(s, "No es culpa de una sola persona.\n\nPero sí es tu responsabilidad cerrarla.", 9.18, 2.72, 2.85, 2.8, 22, WHITE, True)
rect(s, 1.02, 5.54, 7.10, 1.13, RGBColor(238, 232, 249), True,
     RGBColor(218, 204, 240))
txt(s, "Ninguna institución posee por sí sola todo el conocimiento necesario: este se encuentra disperso, incompleto y, a veces, contradictorio entre muchas personas.",
    1.24, 5.61, 6.66, .68, 16, PURPLE, True, PP_ALIGN.CENTER,
    MSO_ANCHOR.MIDDLE)
txt(s, "Idea basada en Friedrich A. Hayek, Individualism and Economic Order (1948).",
    1.24, 6.31, 6.66, .25, 11.5, MUTED, False, PP_ALIGN.CENTER,
    MSO_ANCHOR.MIDDLE, italic=True,
    url="https://press.uchicago.edu/ucp/books/book/chicago/I/bo3615705.html")
if QR_CIC.exists():
    rect(s, 11.47, 6.08, .98, .98, WHITE, True, WHITE)
    s.shapes.add_picture(str(QR_CIC), Inches(11.54), Inches(6.15), width=Inches(.84))
txt(s, "MATERIALES DEL XXI CIC\nEscanea para abrir la carpeta", 8.84, 6.24, 2.46, .54,
    10, PURPLE, True, PP_ALIGN.CENTER,
    url="https://github.com/marlenis-concepcion/TallerIA-Docente-27-06-2026/tree/main/Congreso-XXI-CIC-2026")
notes(s, "No presente la academia como enemiga. Explique la diferencia de velocidad entre gobernanza institucional y cambio tecnológico. La frase es una paráfrasis de la tesis del conocimiento disperso desarrollada en el ensayo The use of knowledge in society, luego recopilado en el libro. Referencia del libro: Hayek, F. A. (1948). Individualism and economic order. University of Chicago Press. Referencia del artículo original: Hayek, F. A. (1945). The use of knowledge in society. The American Economic Review, 35(4), 519–530. https://www.aeaweb.org/aer/top20/35.4.519-530.pdf")

# 4
s = prs.slides.add_slide(blank)
rect(s, 0, 0, 13.333, 7.5, NAVY)
txt(s, "NO PUEDES ESPERAR", .75, .82, 7.0, .38, 16, CYAN, True)
txt(s, "a que las regulaciones, los programas de clase y las instituciones se pongan de acuerdo.", .75, 1.50, 11.45, 2.0, 31, WHITE, True)
rect(s, .75, 4.02, 11.20, .07, VIOLET)
txt(s, "Debes construir tu propio plan.", .75, 4.55, 10.8, .78, 33, RGBColor(224,175,255), True)
txt(s, "Aprender · practicar · demostrar · ajustar", .78, 5.65, 9.0, .35, 17, CYAN, True)
book_idea(s, "El cambio genera incertidumbre y puede provocar resistencia (Robbins y Judge, 2017).", 6.30)
brand(s, True)
notes(s, "Haga una pausa. Esta es la tesis central de la charla. No espere permiso para aprender, pero tampoco confunda autonomía con improvisación.")

# 5
s = base("Tu plan empieza con una vacante real", 5, "Aprende desde la demanda, no desde el ruido")
cards(s, [
    ("1 · Observar", "Selecciona 10–15 vacantes del rol que deseas.", CYAN),
    ("2 · Mapear", "Separa requisitos frecuentes, deseables y diferenciadores.", VIOLET),
    ("3 · Priorizar", "Elige una brecha técnica y una humana por ciclo.", ORANGE),
    ("4 · Practicar", "Construye algo pequeño que funcione y pueda explicarse.", GREEN),
    ("5 · Mostrar", "Publica evidencia: repositorio, reporte, demo o reflexión.", CYAN),
    ("6 · Ajustar", "Revisa el mapa cada 30 días; conserva fundamentos.", VIOLET),
], cols=3)
notes(s, "Recalque que una ruta propia no es coleccionar cursos. Debe salir de evidencia del mercado y terminar en evidencia del aprendizaje.")

# 6
s = base("No quieras volar sin dominar lo básico", 6, "La velocidad sin fundamento solo produce errores más rápidos")
cards(s, [
    ("Testing", "Riesgo, casos de prueba, particiones, límites, defectos y trazabilidad.", CYAN),
    ("Web y API", "Cliente-servidor, HTTP, contratos, estados, datos y errores.", VIOLET),
    ("Código", "Lógica, estructuras, un lenguaje, Git, depuración y diseño legible.", ORANGE),
    ("Producto", "Usuario, necesidad, negocio, impacto, prioridad y calidad esperada.", GREEN),
    ("Automatización", "Solo después: UI/API, CI, reportes, estrategia y mantenimiento.", CYAN),
    ("IA aplicada", "Con fundamentos: contexto, límites, privacidad y verificación humana.", VIOLET),
], cols=3)
notes(s, "Use la metáfora de los cimientos. Antes de automatizar una prueba hay que saber diseñarla; antes de pedir código a la IA hay que poder leerlo, cuestionarlo y probarlo.")

# 7
s = base("Software tiene muchas carreras", 7, "QA es una especialidad completa, no una sala de espera")
cards(s, [
    ("Profundizar en QA", "Funcional, API, mobile, automatización, SDET, performance, seguridad, accesibilidad o liderazgo.", CYAN),
    ("Moverse si quiere", "Desarrollo, producto, análisis de negocio, datos, UX, DevOps, seguridad o gestión.", VIOLET),
    ("Elegir y reelegir", "No existe una ruta superior ni definitiva. Puedes cambiar cuando tus fortalezas, intereses o realidad también cambien.", ORANGE),
])
sprite_third(s, 0, 1.46, 3.98, 2.15, 2.32, CAREER_SPRITE)
sprite_third(s, 1, 5.58, 3.98, 2.15, 2.32, CAREER_SPRITE)
sprite_third(s, 2, 9.69, 3.98, 2.15, 2.32, CAREER_SPRITE)
book_idea(s, "La satisfacción laboral influye en las actitudes y decisiones profesionales (Robbins y Judge, 2017).")
notes(s, "Explique que el ciclo de desarrollo reúne muchas especialidades. Nadie tiene que abandonar QA para crecer. Puede avanzar en profundidad, amplitud, influencia o liderazgo dentro de calidad; un cambio lateral es una opción, no una obligación.")

# 8
s = base("Así se veía una vacante antes del auge de IA", 8,
         "Ejemplo ilustrativo de expectativas frecuentes hasta 2023; no representa a todas las empresas")
rect(s, .68, 1.42, 7.78, 5.35, WHITE, True, LINE)
txt(s, "QA AUTOMATION ENGINEER", 1.00, 1.72, 6.85, .42, 21, INK, True)
txt(s, "Producto digital · Equipo ágil · Modalidad híbrida", 1.00, 2.18, 6.60, .30, 11.5, MUTED)
pill(s, "Selenium", 1.00, 2.70, 1.15, CYAN)
pill(s, "Java", 2.27, 2.70, .82, PURPLE)
pill(s, "API", 3.21, 2.70, .72, ORANGE)
pill(s, "SQL", 4.05, 2.70, .72, GREEN)
pill(s, "Git + CI", 4.89, 2.70, 1.05, VIOLET)
bullets(s, [
    "Analizar requisitos y diseñar casos de prueba.",
    "Ejecutar pruebas funcionales, API y regresión.",
    "Automatizar escenarios estables y mantener la suite.",
    "Reportar defectos con evidencia y colaborar con desarrollo.",
    "Conocer programación, control de versiones y pipelines."
], 1.10, 3.35, 6.72, 2.75, 15.2, INK, CYAN)
rect(s, 8.78, 1.54, 3.82, 4.95, NAVY, True)
txt(s, "EL CENTRO", 9.18, 1.96, 2.95, .30, 12, CYAN, True)
txt(s, "Diseñar\nEjecutar\nAutomatizar\nDocumentar\nColaborar", 9.18, 2.58, 2.95, 2.70, 22, WHITE, True)
txt(s, "La técnica importaba, pero también la comunicación y el conocimiento del producto.", 9.18, 5.55, 2.92, .62, 11, RGBColor(220,225,245))
notes(s, "Esta vacante es una síntesis ilustrativa. Antes de 2024 ya existían IA, machine learning y automatización avanzada, pero muchas vacantes generales de QA se concentraban en fundamentos, API, UI, SQL, Git, CI, diseño y colaboración.")

# 9
s = base("Así se reorganiza una vacante con IA", 9,
         "Los fundamentos permanecen; se añaden uso asistido, evaluación, seguridad y responsabilidad")
rect(s, .68, 1.42, 7.78, 5.35, WHITE, True, LINE)
txt(s, "QUALITY ENGINEER · AI-ASSISTED", 1.00, 1.72, 6.85, .42, 21, INK, True)
txt(s, "Producto con IA · Automatización moderna · Equipos multidisciplinarios", 1.00, 2.18, 6.90, .30, 11.5, MUTED)
pill(s, "Playwright", 1.00, 2.70, 1.28, CYAN)
pill(s, "API + CI", 2.40, 2.70, 1.08, PURPLE)
pill(s, "GenAI", 3.60, 2.70, .92, ORANGE)
pill(s, "LLM eval", 4.64, 2.70, 1.05, GREEN)
pill(s, "AI security", 5.81, 2.70, 1.24, VIOLET)
bullets(s, [
    "Mantener diseño de pruebas, API, datos, código, Git y CI.",
    "Usar asistentes de IA sin entregarles el juicio profesional.",
    "Verificar código, casos, resúmenes y respuestas generadas.",
    "Evaluar alucinaciones, sesgos, seguridad, privacidad y guardrails.",
    "Comunicar límites, evidencia y riesgo a producto y negocio."
], 1.10, 3.35, 6.72, 2.75, 15.2, INK, CYAN)
rect(s, 8.78, 1.54, 3.82, 4.95, PURPLE, True)
txt(s, "LO NUEVO", 9.18, 1.96, 2.95, .30, 12, CYAN, True)
txt(s, "Contexto\nVerificación\nPrivacidad\nEvaluación de IA\nAprendizaje continuo", 9.18, 2.58, 2.95, 2.70, 21, WHITE, True)
txt(s, "No busques solamente el título del puesto: lee responsabilidades, riesgos y evidencia esperada.", 9.18, 5.55, 2.92, .62, 11, RGBColor(235,225,245))
notes(s, "Desde 2024 muchas vacantes conservaron las responsabilidades anteriores y añadieron alfabetización en IA, verificación de resultados, privacidad, seguridad, evaluación de LLM y colaboración con producto y datos. Los títulos cambian entre empresas; compare funciones, no solo nombres.")

# 10
s = base("La calidad también es una relación humana", 10, "Un defecto se reporta a una persona; una decisión se toma con otras personas")
cards(s, [
    ("Escuchar antes de responder", "Comprender intención, contexto y restricciones; no preparar la defensa mientras la otra persona habla.", CYAN),
    ("Preguntar sin atacar", "Cuestionar requisitos y supuestos con curiosidad, respeto y propósito compartido.", VIOLET),
    ("Reportar sin culpar", "Describir evidencia, impacto y condiciones; el defecto no define a quien escribió el código.", ORANGE),
    ("Recibir feedback", "Separar la observación profesional del valor personal y convertirla en una acción.", GREEN),
    ("Manejar desacuerdos", "Hablar de riesgo, datos y alternativas; no convertir una diferencia técnica en conflicto personal.", CYAN),
    ("Construir confianza", "Decir la verdad, cumplir acuerdos, reconocer límites, pedir ayuda y comunicar temprano.", VIOLET),
], cols=3)
book_idea(s, "La comunicación reduce incertidumbre y facilita la coordinación (Robbins y Judge, 2017).")
notes(s, "Pida ejemplos reales: cómo reportar un defecto sin humillar, cómo disentir con desarrollo y cómo reconocer que uno se equivocó. Las habilidades humanas se demuestran en conductas.")

# 11
s = base("La IA amplifica; no reemplaza el juicio ni la relación", 11)
bullets(s, [
    "Úsala para explorar casos, generar borradores, explicar código y acelerar investigación.",
    "Verifica requisitos, datos, resultados y afirmaciones antes de confiar.",
    "No compartas secretos, datos personales ni propiedad intelectual sin autorización.",
    "No delegues conversaciones difíciles a una herramienta: pregunta, escucha y acuerda con el equipo."
], 1.0, 1.60, 7.35, 4.8, 18)
rect(s, 8.72, 1.62, 3.82, 4.55, PURPLE, True)
txt(s, "Fórmula útil", 9.10, 2.05, 3.0, .35, 14, CYAN, True)
txt(s, "IA\n+ contexto\n+ criterio\n+ verificación", 9.10, 2.65, 3.0, 2.75, 24, WHITE, True)
book_idea(s, "Las emociones influyen en cómo interpretamos y respondemos al trabajo (Robbins y Judge, 2017).")
notes(s, "La habilidad importante no es obtener una respuesta, sino producir una decisión verificable y poder conversar sobre ella con el equipo.")

# 12
s = base("Automatizar no es grabar clics", 12, "Es decidir qué merece repetirse y mantenerse")
cards(s, [
    ("Valor", "Prioriza recorridos críticos, riesgos recurrentes y feedback rápido.", CYAN),
    ("Diseño", "Pruebas legibles, datos controlados, selectores estables y capas claras.", VIOLET),
    ("Mantenimiento", "Una prueba frágil consume confianza, tiempo y atención del equipo.", ORANGE),
])
sprite_third(s, 0, 1.46, 3.98, 2.15, 2.32, AUTOMATION_SPRITE)
sprite_third(s, 1, 5.58, 3.98, 2.15, 2.32, AUTOMATION_SPRITE)
sprite_third(s, 2, 9.69, 3.98, 2.15, 2.32, AUTOMATION_SPRITE)
notes(s, "Conecte la automatización con pensamiento de sistemas. Una suite grande no equivale automáticamente a mayor calidad.")

# 13
s = base("Certificar no sustituye practicar", 13, "La teoría da lenguaje y estructura; la práctica desarrolla criterio")
cards(s, [
    ("Teoría", "Explica por qué: principios, vocabulario, técnicas, riesgos y modelos mentales.", CYAN),
    ("Práctica", "Demuestra cómo: diseña, ejecuta, falla, depura, comunica y vuelve a intentar.", VIOLET),
    ("Certificación", "Es valiosa cuando organiza el aprendizaje y se conecta con experiencia aplicable.", ORANGE),
    ("Evidencia", "Repositorio, reporte, caso explicado, conversación técnica y reflexión sobre lo aprendido.", GREEN),
], cols=2)
book_idea(s, "La motivación conecta esfuerzo, desempeño y recompensas percibidas (Robbins y Judge, 2017).")
notes(s, "No desacredite las certificaciones. Explique su lugar correcto: validan un cuerpo de conocimiento, pero la competencia aparece al aplicar, equivocarse, corregir y explicar decisiones. Teoría y práctica van de la mano.")

# 14
s = base("Práctica para hacer a tu ritmo", 14, "Una ruta de casa: ejecutar, comprender, mejorar y explicar")
rect(s, .72, 1.50, 11.88, 4.85, NAVY, True)
txt(s, "vamcodeAutomationPractice", 1.12, 1.92, 6.8, .50, 24, WHITE, True,
    url="https://github.com/marlenis-concepcion/vamcodeAutomationPractice")
bullets(s, [
    "Lee el README y ejecuta la solución existente.",
    "Agrega o mejora un escenario con valor claro.",
    "Refactoriza solo cuando puedas explicar el beneficio.",
    "Conserva commits, reporte y una reflexión breve."
], 1.42, 2.75, 7.25, 3.1, 17, WHITE, CYAN)
txt(s, "ABRIR REPOSITORIO  ↗", 9.12, 3.17, 2.75, .60, 15, CYAN, True, PP_ALIGN.CENTER,
    MSO_ANCHOR.MIDDLE, url="https://github.com/marlenis-concepcion/vamcodeAutomationPractice")
txt(s, "No busques perfección.\nBusca una mejora demostrable.", 9.05, 4.18, 2.95, 1.0, 17, WHITE, True, PP_ALIGN.CENTER)
notes(s, "Explique que pueden avanzar por etapas y pedir ayuda mostrando el error, el intento y la evidencia disponible.")

# 15
s = base("Aprender testing gratis", 15, "Recursos para comenzar hoy")
resource(s, 1.42, "Test Automation University", "Cursos gratuitos de testing y automatización.", "https://testautomationu.applitools.com/", CYAN)
resource(s, 2.28, "Ruta codeless", "Trayecto gratuito para explorar automatización sin código.", "https://testautomationu.applitools.com/learningpaths.html?id=codeless-path", VIOLET)
resource(s, 3.14, "Swagger Petstore", "API pública para practicar solicitudes, respuestas y contratos.", "https://petstore.swagger.io/#/pet/addPet", ORANGE)
resource(s, 4.00, "Playwright", "Documentación oficial para automatización web moderna.", "https://playwright.dev/", GREEN)
resource(s, 4.86, "Playwright en GitHub", "Código, ejemplos, issues y evolución del proyecto.", "https://github.com/microsoft/playwright", CYAN)
resource(s, 5.72, "Katalon", "Plataforma con opciones de aprendizaje y automatización.", "https://katalon.com/", VIOLET)
notes(s, "No recomiende consumirlos todos. Cada persona debe elegir un recurso alineado con su brecha actual.")

# 16
s = base("Aprender IA y datos gratis", 16, "La IA también exige fundamentos")
resource(s, 1.55, "Claude Courses", "Cursos introductorios y prácticos sobre IA.", "https://claude.com/resources/courses", VIOLET)
resource(s, 2.55, "DeepLearning.AI", "Cursos de IA y datos; filtra según tu ruta.", "https://www.deeplearning.ai/courses?q=data", CYAN)
resource(s, 3.55, "Google Trends", "Compara interés relativo; úsalo como señal, no como sentencia.", "https://trends.google.com/explore?date=today%201-y&geo=Worldwide&q=%2Fm%2F0mkz%2C%2Fm%2F05z1_%2C%2Fm%2F02p97%2C%2Fm%2F012l1vxv%2C%2Fm%2F0bbxf89%2C%2Fg%2F11c6w0ddw9%2C%2Fm%2F0n50hxv%2C%2Fm%2F0_h5pbr", ORANGE)
resource(s, 4.55, "TallerIA-Docente-27-06-2026", "Repositorio central: aquí encontrarás los enlaces y materiales.", RESOURCES_REPO, GREEN)
rect(s, .72, 5.65, 11.88, .72, RGBColor(235,229,248), True, RGBColor(220,205,240))
txt(s, "Regla: teoría → práctica → feedback humano → evidencia → mejora.", 1.05, 5.83, 11.2, .32, 16, PURPLE, True, PP_ALIGN.CENTER)
notes(s, "Presente los recursos como puntos de partida. La habilidad crece al aplicar y explicar, no al acumular certificados.")

# 17
s = base("Tu ciclo de 30 días", 17, "Pequeño, sostenible y verificable")
cards(s, [
    ("Semana 1", "Analiza vacantes y define una brecha técnica y una humana.", CYAN),
    ("Semana 2", "Estudia lo mínimo necesario y reproduce un ejemplo.", VIOLET),
    ("Semana 3", "Construye una mejora propia; pide feedback específico.", ORANGE),
    ("Semana 4", "Publica evidencia, explica decisiones y actualiza tu mapa.", GREEN),
], cols=2)
book_idea(s, "Las metas orientan mejor cuando son claras y cuentan con recursos para ejecutarse (Robbins y Judge, 2017).")
notes(s, "Invite a escribir ahora la habilidad, la práctica y la evidencia que cada participante completará durante los próximos 30 días.")

# 18
s = base("Priorizarte también es tomar decisiones", 18,
         "No conviertas una situación que puedes cambiar en una culpa permanente hacia los demás")
cards(s, [
    ("Reconoce la realidad", "Si el salario, el ambiente o el crecimiento no te funcionan, nómbralo con honestidad y evalúa tus opciones.", CYAN),
    ("Prepara el movimiento", "Negocia, aprende, ahorra, busca y crea evidencia. Si no puedes salir hoy, construye una transición responsable.", VIOLET),
    ("Cambia de rumbo", "Puedes moverte de posición, especialidad, empresa o carrera. El mercado cambia; tú también tienes derecho a evolucionar.", ORANGE),
])
sprite_third(s, 0, 1.46, 3.98, 2.15, 2.32, PRIORITY_SPRITE)
sprite_third(s, 1, 5.58, 3.98, 2.15, 2.32, PRIORITY_SPRITE)
sprite_third(s, 2, 9.69, 3.98, 2.15, 2.32, PRIORITY_SPRITE)
book_idea(s, "La percepción influye directamente en la toma de decisiones (Robbins y Judge, 2017).")
notes(s, "Explique la diferencia entre responsabilidad y culpabilización. Si una persona no está conforme con salario, posición o crecimiento, puede negociar o preparar un cambio. No todo el mundo puede renunciar inmediatamente por responsabilidades económicas, familiares o migratorias; aun así puede diseñar una transición. Adaptarse no significa mentir, incumplir o abandonar clientes sin comunicación. Tanto la organización como el profesional pueden cambiar objetivos y rumbo, respetando acuerdos y consecuencias.")

# 19
s = base("Aprende a leer la organización y el país", 19,
         "Manejar lo técnico también exige comprender el lenguaje donde se toman decisiones")
cards(s, [
    ("Lenguaje corporativo", "Valor, cliente, capacidad, coste, riesgo, prioridad, indicador, dependencia, retorno y nivel de servicio.", CYAN),
    ("Contexto nacional", "Regulación, políticas públicas, economía, inversión, informalidad, idioma, sectores y concentración geográfica del empleo.", VIOLET),
    ("Decisión profesional", "Contrasta aspiraciones con vacantes, salarios, modalidad, barreras de entrada y alternativas locales o remotas.", ORANGE),
])
if LEGAL_CIRCLE.exists():
    s.shapes.add_picture(str(LEGAL_CIRCLE), Inches(5.43), Inches(4.02), width=Inches(2.48), height=Inches(2.24))
txt(s, "LEY\n30-26", 6.41, 4.10, .52, .34, 8.2, PURPLE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
txt(s, "16-92\nTRABAJO", 5.48, 4.86, .62, .36, 7.2, CYAN, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
txt(s, "87-01\nSISALRIL", 7.27, 4.86, .62, .36, 7.2, ORANGE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
txt(s, "11-92\nTRIBUTARIA", 6.37, 5.74, .66, .34, 7.0, NAVY, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
if QR_BRAIN_DRAIN.exists():
    s.shapes.add_picture(str(QR_BRAIN_DRAIN), Inches(1.00), Inches(4.78), width=Inches(1.10))
txt(s, "FUGA DE CEREBROS\nEscanea el reportaje", 2.20, 5.03, 2.10, .62,
    12, PURPLE, True, PP_ALIGN.CENTER,
    url="https://www.diariolibre.com/economia/empleo/2021/12/20/rd-registra-fuga-de-cerebros-en-la-industria-de-software/1542664")
if QR_INFORMALITY.exists():
    s.shapes.add_picture(str(QR_INFORMALITY), Inches(11.20), Inches(4.78), width=Inches(1.10))
txt(s, "INFORMALIDAD LABORAL\nEscanea el informe", 8.88, 5.03, 2.15, .62,
    12, PURPLE, True, PP_ALIGN.CENTER,
    url="https://cdn.bancentral.gov.do/documents/publicaciones-economicas/boletin-trimestral-del-mercado-laboral/documents/Boletin_Trimestral_Mercado_Laboral_jul-sep_2025.pdf")
rect(s, .72, 6.42, 11.88, .46, RGBColor(235,229,248), True, RGBColor(220,205,240))
txt(s, "RD, jul.–sep. 2025: 5.15 millones de ocupados; 54.6 % de informalidad. Crecer no elimina todas las brechas.",
    .92, 6.52, 11.50, .24, 10.5, PURPLE, True, PP_ALIGN.CENTER,
    url="https://cdn.bancentral.gov.do/documents/publicaciones-economicas/boletin-trimestral-del-mercado-laboral/documents/Boletin_Trimestral_Mercado_Laboral_jul-sep_2025.pdf")
txt(s, "Referencia: L. Cuatrecasas Arbós, Organización de la producción y dirección de operaciones, Díaz de Santos, 2011, ISBN 9788479789978.",
    .72, 6.91, 11.90, .18, 7.3, MUTED, italic=True,
    url="https://www.editdiazdesantos.com/libros/9788479789978/Organizacion-de-la-produccion-y-direccion-de-operaciones.html")
notes(s, "La política se aborda como contexto institucional, no como afiliación partidaria: regulación laboral, educación, inversión, estabilidad y políticas públicas afectan las oportunidades. Según el Banco Central, en julio-septiembre de 2025 había 5,149,829 ocupados y la informalidad laboral era 54.6 %. Mencione el libro de Cuatrecasas para ampliar el vocabulario de procesos, valor, calidad, costes, capacidad, logística y eficiencia. QA debe traducir defectos técnicos a impacto organizacional.")

# 20
s = prs.slides.add_slide(blank)
if COMIC.exists():
    s.shapes.add_picture(str(COMIC), 0, 0, width=prs.slide_width, height=prs.slide_height)
txt(s, "PRIMER EMPLEO: NIVEL PRINCIPIANTE", .52, .22, 8.60, .42, 22, WHITE, True)
txt(s, "Expectativas: modo leyenda", .54, .68, 6.50, .32, 15, CYAN, True)
for label, x, color in [
    ("ESTADO", .64, PURPLE),
    ("EDUCACIÓN", 3.02, VIOLET),
    ("EMPRESA", 8.73, ORANGE),
    ("SOCIEDAD / FAMILIA", 10.65, PURPLE),
]:
    rect(s, x, 1.25, 1.95 if label != "SOCIEDAD / FAMILIA" else 2.18, .42, color, True, WHITE)
    txt(s, label, x+.05, 1.33, 1.85 if label != "SOCIEDAD / FAMILIA" else 2.08, .22,
        9.5, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
rect(s, 5.16, 6.68, 3.10, .48, NAVY, True, CYAN)
txt(s, "NUEVO PROFESIONAL: «¿RESPIRAR CUENTA?»", 5.27, 6.79, 2.88, .22,
    9.2, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
brand(s, True)
notes(s, "Use la viñeta como pausa cómica. El nuevo profesional recibe simultáneamente normas, títulos, indicadores y expectativas familiares. Pregunte al público cuál de esas voces pesa más al iniciar. La respuesta no es rechazar toda expectativa, sino priorizar, negociar, construir fundamentos y conservar integridad.")

# 21
s = prs.slides.add_slide(blank)
if COMIC_EXPECTATIONS.exists():
    s.shapes.add_picture(str(COMIC_EXPECTATIONS), 0, 0, width=prs.slide_width, height=prs.slide_height)
txt(s, "DOS CAMINADORAS, DEMASIADAS EXPECTATIVAS", .52, .20, 9.80, .43, 22, WHITE, True)
txt(s, "La sociedad entrega la carga… y después pregunta por qué te cansas.", .54, .68, 8.80, .30, 14, CYAN, True)
rect(s, .62, 1.18, 4.70, .56, PURPLE, True, WHITE)
txt(s, "SOBRE ÉL: «PROVEE, RESUELVE Y NUNCA TE CANSES»", .76, 1.31, 4.42, .25,
    10.8, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
rect(s, 7.18, 1.18, 5.45, .56, VIOLET, True, WHITE)
txt(s, "SOBRE ELLA: «PRODUCE, CUIDA, ORGANIZA Y SONRÍE»", 7.32, 1.31, 5.17, .25,
    10.8, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
rect(s, 4.43, 6.77, 4.48, .42, NAVY, True, ORANGE)
txt(s, "GERENCIA: «¿Y SI SUBIMOS LA VELOCIDAD?»", 4.55, 6.86, 4.24, .20,
    9.6, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
brand(s, True)
notes(s, "Aclare que son presiones sociales frecuentes, no destinos universales: al hombre se le puede exigir proveer, resolver y ocultar cansancio; a la mujer, producir, cuidar, organizar y mantener una imagen impecable. Ambos pierden cuando la sociedad distribuye roles rígidos. El personaje calmado representa planificación y recursos; el gerente caótico representa presión fabricada por mala gestión.")

# 22
s = base("Trabajar bajo presión no debe ser el estado normal", 22,
         "Una urgencia ocasional es real; vivir en urgencia es una señal de gestión")
cards(s, [
    ("Presión ≠ excelencia", "Responder a una emergencia puede ser valioso. Convertir el agotamiento permanente en medalla destruye calidad y personas.", CYAN),
    ("Desorden ≠ agilidad", "Prioridades que cambian sin criterio, retrabajo y dependencia de héroes suelen revelar planificación incorrecta.", VIOLET),
    ("Meta sin recursos", "Alcance, tiempo, personas, información y herramientas deben ser compatibles. Lo imposible no se vuelve posible por insistir.", ORANGE),
])
sprite_third(s, 0, 1.46, 3.98, 2.15, 2.32, PRESSURE_SPRITE)
sprite_third(s, 1, 5.58, 3.98, 2.15, 2.32, PRESSURE_SPRITE)
sprite_third(s, 2, 9.69, 3.98, 2.15, 2.32, PRESSURE_SPRITE)
rect(s, .72, 6.46, 11.88, .50, RGBColor(235,229,248), True, RGBColor(220,205,240))
txt(s, "El estrés aumenta cuando las exigencias superan los recursos percibidos (Robbins y Judge, 2017).",
    .94, 6.51, 11.44, .40, 15.5, PURPLE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
notes(s, "Distinga presión ocasional de presión crónica. La primera puede requerir coordinación especial; la segunda normaliza deuda, errores y desgaste. No culpe al equipo por una meta inalcanzable con los recursos actuales. Haga visibles alcance, capacidad, dependencias, riesgos y decisiones de prioridad.")

# 23
s = base("El mundo laboral no sigue el guion universitario", 23,
         "Observa patrones del mercado; investiga el contexto sin convertir percepciones en hechos")
cards(s, [
    ("Mira el mercado", "Rastreadores conocidos como Layoffs.fyi y TrueUp muestran empresas y sectores donde se han reportado recortes de personal.", CYAN),
    ("Busca contexto cultural", "La cultura orienta cómo las personas interpretan el trabajo (Robbins y Judge, 2017). Las reseñas pueden ofrecer señales para investigar.", VIOLET),
    ("Contrasta las fuentes", "Un rastreador muestra el evento; noticias y comunicados aportan contexto. Una reseña anónima expresa una experiencia, no un hecho probado.", ORANGE),
    ("Prepárate profesionalmente", "Documenta resultados, lee el ambiente, conversa temprano, cuida tu red y construye alternativas antes de necesitarlas.", GREEN),
], cols=2)
txt(s, "Fuentes para explorar tendencias: Layoffs.fyi · TrueUp Tech Layoffs Tracker · Glassdoor Employee Reviews.",
    .74, 6.62, 11.85, .28, 7.5, MUTED, italic=True, align=PP_ALIGN.CENTER,
    url="https://layoffs.fyi/")
notes(s, "No mencione una empresa ni convierta esta lámina en una denuncia particular. Layoffs.fyi y TrueUp permiten observar recortes reportados y patrones del mercado tecnológico. Después, las noticias, los comunicados corporativos y las reseñas de Glassdoor pueden ayudar a comprender el contexto. Aclare siempre que una reseña anónima representa una percepción individual: sirve como señal para investigar, no como prueba definitiva. La lección es aprender a leer el mercado y la cultura organizacional, comunicar, documentar resultados y preparar alternativas con anticipación.")

# 24
s = base("Tu enfoque sí tiene respaldo en administración", 24,
         "Operaciones explica el sistema; administración explica cómo dirigir personas dentro de él")
cards(s, [
    ("Cuatrecasas · Operaciones",
     "“Es en los procesos de producción donde la empresa genera su mayor o menor valor añadido.”\n\nQA debe traducir defectos a valor, riesgo, coste y capacidad.", CYAN),
    ("Varela · Compensación",
     "“La buena voluntad del empleado de hacer esas contribuciones depende de la forma como perciba su contribución en relación con la recompensa que recibe.”\n\nPriorizarse y pedir reciprocidad no es deslealtad.", VIOLET),
    ("Robbins y Coulter · Administración",
     "“Aprender a leer la cultura de una organización.”\n\nLa empleabilidad también requiere lenguaje corporativo, negociación y criterio político-organizacional.", ORANGE),
    ("Robbins y Judge · Comportamiento organizacional",
     "“Cambio organizacional y manejo del estrés.”\n\nEl desempeño también depende de valores, emociones, motivación, comunicación, conflicto, cultura y relaciones humanas.", GREEN),
], cols=2)
txt(s, "Fuentes: Díaz de Santos; Pearson, Administración; Pearson, Comportamiento organizacional; Varela Juárez, Administración de la compensación.",
    .72, 6.63, 11.88, .28, 7.5, MUTED, italic=True, align=PP_ALIGN.CENTER,
    url="https://www.pearsonenespanol.com/mexico/tienda-online/administracion-1ed-robbins-mx-ebook")
notes(s, "Presente estas citas como respaldo, no como autoridad incuestionable. Cuatrecasas sitúa el valor en los procesos; Varela conecta contribución y recompensa; Robbins y Coulter incluyen leer la cultura organizacional como habilidad profesional; Robbins y Judge estudian satisfacción laboral, emociones, valores, motivación, comunicación, liderazgo, poder, política, conflicto, cultura, cambio y estrés. Esto confirma que hablar de personas, reciprocidad y presión laboral explica el sistema humano donde la tecnología produce valor.")

# 25
s = base("Cuatro pirámides, cuatro prioridades", 25,
         "El sistema funciona mejor cuando cada grupo entiende lo que necesita el otro")
mini_pyramid(s, .55, 1.45, 6.00, "Estado",
             ["Instituciones y reglas", "Infraestructura y seguridad", "Empleo formal y productividad", "Bienestar sostenible"],
             [PURPLE, VIOLET, CYAN, GREEN])
mini_pyramid(s, 6.78, 1.45, 6.00, "Educación",
             ["Acceso y alfabetización", "Fundamentos disciplinares", "Práctica, investigación y criterio", "Impacto social"],
             [PURPLE, VIOLET, CYAN, GREEN])
mini_pyramid(s, .55, 4.05, 6.00, "Empresa",
             ["Viabilidad y cumplimiento", "Procesos, personas y calidad", "Valor para cliente y mercado", "Continuidad e innovación"],
             [PURPLE, VIOLET, ORANGE, GREEN])
mini_pyramid(s, 6.78, 4.05, 6.00, "Profesional",
             ["Salud, valores e integridad", "Fundamentos y lenguaje", "Práctica, relaciones y evidencia", "Autonomía y propósito"],
             [PURPLE, VIOLET, CYAN, GREEN])
txt(s, "El poder y la política forman parte de la vida organizacional (Robbins y Judge, 2017).",
    .78, 6.53, 11.78, .42, 15.5, PURPLE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
notes(s, "Lea cada pirámide desde la base. El Estado busca reglas, capacidad económica y bienestar; la educación convierte acceso y fundamentos en impacto; la empresa necesita viabilidad antes de innovar; y el profesional necesita salud e integridad antes de sostener autonomía. Relacione la lámina con dirección de operaciones: el valor surge de coordinar personas, procesos, recursos, calidad y objetivos, no de optimizar una sola parte.")

# 26
s = base("Comportamiento organizacional: del libro a la vida", 26,
         "Robbins y Judge conectan persona, grupo y sistema organizacional")
cards(s, [
    ("La persona",
     "“Las emociones y los estados de ánimo.”\n“Percepción y toma de decisiones individual.”\n\nLo que sientes y percibes influye en cómo decides.", CYAN),
    ("El grupo",
     "“Comunicación.”\n“Poder y política.”\n“Conflicto y negociación.”\n\nTrabajar también exige leer relaciones e intereses.", VIOLET),
    ("La organización",
     "“Cultura organizacional.”\n“Políticas y prácticas de recursos humanos.”\n\nLa empresa es un sistema, no solo un puesto técnico.", ORANGE),
    ("Cambio y bienestar",
     "“Cambio organizacional y manejo del estrés.”\n\nAdaptarse no significa normalizar desgaste, caos o falta de planificación.", GREEN),
], cols=2)
txt(s, "Referencia: Robbins, S. P. y Judge, T. A. · Comportamiento organizacional, 17.ª ed. · Pearson · ISBN 9786073239851.",
    .72, 6.62, 11.88, .28, 8, MUTED, italic=True, align=PP_ALIGN.CENTER,
    url="https://www.pearsonenespanol.com/mexico/educacion-superior/robbins/robbins_comportamiento_organizacional_17e_contenido")
notes(s, "Estas frases corresponden a títulos de capítulos de Comportamiento organizacional. Úselas como mapa: individuo, grupo y sistema. Pregunte cuál dimensión suelen ignorar los profesionales técnicos. El objetivo no es diagnosticar personas, sino observar cómo emociones, percepción, comunicación, poder, cultura, recursos humanos, cambio y estrés influyen en el trabajo real.")

# 27
s = prs.slides.add_slide(blank)
rect(s, 0, 0, 13.333, 7.5, NAVY)
if LOGO.exists(): s.shapes.add_picture(str(LOGO), Inches(9.82), Inches(.40), width=Inches(2.15))
txt(s, "LUEGO DE\nESTOS TEMAS…", .75, .85, 8.2, 1.45, 34, WHITE, True)
txt(s, "¿Qué aplicarás mañana en tu día a día?", .78, 2.72, 7.85, .72, 24, RGBColor(221,215,242), True)
rect(s, .78, 4.12, 7.30, .07, VIOLET)
txt(s, "Comunicar mejor · leer el ambiente · cuidar tus límites · preparar tu siguiente paso", .78, 4.48, 8.25, 1.0, 19, CYAN, True)
repo_qr_badge(s, 10.05, 2.72, 1.62, True)
txt(s, "TallerIA-Docente-27-06-2026", 8.70, 4.70, 4.30, .36, 13, WHITE, True, PP_ALIGN.CENTER,
    url=RESOURCES_REPO)
txt(s, "github.com/marlenis-concepcion/\nTallerIA-Docente-27-06-2026", 8.70, 5.14, 4.30, .62, 10.5, CYAN, False, PP_ALIGN.CENTER,
    url=RESOURCES_REPO)
txt(s, "Marlenis Judith Concepción Cuevas · LinkedIn ↗", 8.70, 6.02, 4.30, .32, 10.5, WHITE, True, PP_ALIGN.CENTER,
    url="https://www.linkedin.com/in/marlenis-judith-c-55f1117a3/")
brand(s, True)
notes(s, "Haga la pregunta directamente al público y permita algunas respuestas: Luego de exponerte estos temas, ¿qué aplicarás en tu día a día? Pida una conducta observable, no una intención genérica. Cierre indicando que el código QR abre el repositorio TallerIA-Docente-27-06-2026 con enlaces y materiales.")

prs.core_properties.title = "La transformación del Software QA Engineering: IA, automatización y empleabilidad"
prs.core_properties.subject = "Testing4All 2026"
prs.core_properties.author = "Marlenis Judith Concepción Cuevas"
prs.core_properties.keywords = "QA, testing, IA, automatización, empleabilidad, habilidades blandas"
prs.save(OUT)
print(f"Creada: {OUT}")
print(f"Diapositivas: {len(prs.slides)}")
