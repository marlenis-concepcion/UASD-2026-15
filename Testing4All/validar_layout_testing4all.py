from pathlib import Path
import sys

ROOT = Path(__file__).resolve().parent
VENDOR = ROOT.parent / "TallerIA-27-06-2026" / "PruebaPersonal" / ".vendor"
sys.path.insert(0, str(VENDOR))

from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE


PPTX = ROOT / "Testing4All_Transformacion_QA_IA_Empleabilidad_2026_LETRA_GRANDE_V5.pptx"


def inside(inner, outer, tolerance=2):
    return (
        inner.left >= outer.left - tolerance
        and inner.top >= outer.top - tolerance
        and inner.left + inner.width <= outer.left + outer.width + tolerance
        and inner.top + inner.height <= outer.top + outer.height + tolerance
    )


prs = Presentation(PPTX)
issues = []
text_boxes = 0

for slide_number, slide in enumerate(prs.slides, 1):
    for shape in slide.shapes:
        if shape.width <= 0 or shape.height <= 0:
            issues.append(f"Lámina {slide_number}: objeto sin dimensiones válidas")
        if (
            shape.left < 0
            or shape.top < 0
            or shape.left + shape.width > prs.slide_width
            or shape.top + shape.height > prs.slide_height
        ):
            issues.append(f"Lámina {slide_number}: objeto fuera de los límites de la diapositiva")

        if not shape.has_text_frame or not shape.text.strip():
            continue

        text_boxes += 1
        if shape.text_frame.auto_size != MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE:
            issues.append(f"Lámina {slide_number}: texto sin protección contra desbordamiento")

        # Comprueba el marco visual más pequeño que contiene al texto.
        containers = [
            candidate
            for candidate in slide.shapes
            if candidate is not shape
            and not candidate.has_text_frame
            and inside(shape, candidate)
        ]
        if containers:
            container = min(containers, key=lambda item: item.width * item.height)
            if not inside(shape, container):
                issues.append(f"Lámina {slide_number}: texto fuera de su marco visual")

print(f"Diapositivas revisadas: {len(prs.slides)}")
print(f"Cuadros de texto revisados: {text_boxes}")
if issues:
    print(f"Incidencias: {len(issues)}")
    for issue in issues:
        print(f"- {issue}")
    raise SystemExit(1)
print("Resultado: todos los objetos están dentro de la lámina y los textos tienen ajuste al marco.")
