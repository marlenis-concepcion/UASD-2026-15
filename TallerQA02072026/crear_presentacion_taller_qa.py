from pathlib import Path
import re
import sys

ROOT = Path(__file__).resolve().parent
VENDOR = ROOT.parent / "TallerIA-27-06-2026" / "PruebaPersonal" / ".vendor"
sys.path.insert(0, str(VENDOR))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


OUT = ROOT / "Taller_QA_IA_Automatizacion_Empleabilidad_02072026.pptx"
UASD_LOGO = ROOT / "logo_uasd.png"
RAFFLE_FORM_URL = (
    "https://docs.google.com/forms/d/e/"
    "1FAIpQLSfkSlA1WPgVU4sj6K3Ol3QkT7Dnbv3puESRw7DDC_3QESRljA/viewform?usp=header"
)
EVENT_INSTAGRAM_URL = "https://www.instagram.com/p/DaB3AGmACVT/"

# Paleta inspirada en tecnología, QA y UASD.
NAVY = RGBColor(7, 19, 42)
NAVY_2 = RGBColor(12, 35, 65)
BLUE = RGBColor(25, 92, 170)
CYAN = RGBColor(49, 202, 219)
TEAL = RGBColor(27, 163, 156)
GREEN = RGBColor(51, 173, 118)
ORANGE = RGBColor(244, 160, 58)
RED = RGBColor(225, 82, 92)
PURPLE = RGBColor(124, 92, 188)
INK = RGBColor(22, 38, 60)
MUTED = RGBColor(88, 106, 128)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(244, 248, 252)
PALE = RGBColor(224, 245, 249)
LINE = RGBColor(218, 228, 238)
DARK_CARD = RGBColor(18, 42, 75)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


def rect(slide, x, y, w, h, fill, radius=False, line=None, transparency=0):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.fill.transparency = transparency
    shp.line.color.rgb = line if line else fill
    return shp


def circle(slide, x, y, d, fill, line=None, width=1.2):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line if line else fill
    shp.line.width = Pt(width)
    return shp


def line(slide, x1, y1, x2, y2, color=LINE, width=1.5):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x1), Inches(y1), Inches(max(x2 - x1, 0.02)), Inches(max(y2 - y1, 0.02))
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.color.rgb = color
    return shp


def text(slide, value, x, y, w, h, size=22, color=INK, bold=False,
         align=PP_ALIGN.LEFT, font="Aptos", valign=MSO_ANCHOR.TOP,
         margin=0.04, italic=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = valign
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = value
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def bullets(slide, items, x, y, w, h, size=18, color=INK, gap=8,
            bullet_color=None, check=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.04)
    frame.margin_right = Inches(0.04)
    for i, item in enumerate(items):
        p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        p.text = ("✓  " if check else "•  ") + item
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = bullet_color if bullet_color else color
        p.space_after = Pt(gap)
        p.line_spacing = 1.05
    return box


def add_notes(slide, note):
    frame = slide.notes_slide.notes_text_frame
    if frame is not None:
        frame.text = note.strip()


def footer(slide, number, label="Taller QA · 2 de julio de 2026"):
    text(slide, label, 0.52, 7.12, 5.7, 0.20, 8.5, MUTED)
    text(slide, f"{number:02d}", 12.20, 7.08, 0.55, 0.25, 9, MUTED, True, PP_ALIGN.RIGHT)


def base(slide, number, title_value, subtitle=None, source=None):
    rect(slide, 0, 0, 13.333, 7.5, LIGHT)
    rect(slide, 0, 0, 13.333, 0.10, CYAN)
    text(slide, f"{number:02d}", 0.55, 0.30, 0.62, 0.38, 12, BLUE, True)
    text(slide, title_value, 1.20, 0.24, 11.35, 0.56, 26, INK, True)
    if subtitle:
        text(slide, subtitle, 1.20, 0.83, 11.15, 0.42, 12.5, MUTED)
    if source:
        text(slide, source, 0.58, 6.82, 11.6, 0.20, 7.5, MUTED, italic=True)
    footer(slide, number)


def add_cover(title_value, subtitle, meta, note):
    slide = prs.slides.add_slide(blank)
    rect(slide, 0, 0, 13.333, 7.5, NAVY)
    rect(slide, 0, 0, 13.333, 0.12, CYAN)
    rect(slide, 8.15, 0, 5.18, 7.5, NAVY_2)
    for i, (x, y, d, c) in enumerate([
        (8.70, 0.75, 2.85, CYAN), (10.25, 2.25, 2.25, BLUE),
        (8.45, 4.45, 1.72, TEAL), (10.70, 5.15, 1.25, ORANGE)
    ]):
        circle(slide, x, y, d, c, c)
        text(slide, ["QA", "IA", "CI", "CV"][i], x, y + d * 0.34, d, 0.48,
             22 if d > 2 else 16, WHITE, True, PP_ALIGN.CENTER)
    text(slide, "CHARLA–TALLER", 0.72, 0.76, 4.0, 0.35, 15, CYAN, True)
    text(slide, title_value, 0.72, 1.28, 7.05, 2.05, 31, WHITE, True)
    text(slide, subtitle, 0.75, 3.62, 6.85, 0.88, 18, PALE)
    rect(slide, 0.74, 4.85, 6.45, 0.06, CYAN)
    text(slide, meta, 0.74, 5.20, 6.60, 1.25, 13.5, WHITE)
    text(slide, "Mtra. Marlenis Judith Concepción Cuevas", 0.74, 6.32, 7.0, 0.28, 11.5, CYAN, True)
    text(slide, "Senior Software Development Engineer in Test (SDET) · Docente e investigadora",
         0.74, 6.61, 7.12, 0.24, 9.1, WHITE, True)
    text(slide, "Escuela de Informática · Facultad de Ciencias · Universidad Autónoma de Santo Domingo",
         0.74, 6.88, 7.16, 0.22, 8.2, PALE)
    if UASD_LOGO.exists():
        slide.shapes.add_picture(str(UASD_LOGO), Inches(11.92), Inches(0.30), width=Inches(1.05))
    add_notes(slide, note)


def add_section(number, roman, title_value, subtitle, note, accent=CYAN):
    slide = prs.slides.add_slide(blank)
    rect(slide, 0, 0, 13.333, 7.5, NAVY)
    rect(slide, 0, 0, 0.18, 7.5, accent)
    text(slide, "SECCIÓN", 0.85, 0.82, 1.8, 0.28, 12, accent, True)
    text(slide, roman, 0.78, 1.28, 2.15, 1.15, 54, WHITE, True)
    rect(slide, 3.15, 1.12, 0.06, 4.9, accent)
    text(slide, title_value, 3.65, 1.40, 8.65, 1.25, 31, WHITE, True)
    text(slide, subtitle, 3.67, 3.05, 7.90, 1.15, 18, PALE)
    text(slide, "Preguntar · practicar · demostrar", 3.67, 5.42, 6.0, 0.32, 13, accent, True)
    footer(slide, number, "Taller QA · Ruta de aprendizaje y empleabilidad")
    add_notes(slide, note)


def add_bullet_slide(number, title_value, subtitle, items, note, accent=BLUE,
                     callout=None, source=None):
    slide = prs.slides.add_slide(blank)
    base(slide, number, title_value, subtitle, source)
    rect(slide, 0.78, 1.50, 0.12, 4.62, accent, True)
    bullets(slide, items, 1.18, 1.52, 7.20, 4.78, 17.2, INK, 10)
    if callout:
        rect(slide, 8.78, 1.58, 3.72, 4.36, NAVY, True)
        text(slide, callout[0].upper(), 9.16, 1.97, 2.95, 0.32, 12, CYAN, True)
        text(slide, callout[1], 9.16, 2.58, 2.92, 2.70, 20, WHITE, True)
        if len(callout) > 2:
            text(slide, callout[2], 9.16, 5.24, 2.92, 0.42, 11, PALE)
    add_notes(slide, note)


def add_cards_slide(number, title_value, subtitle, cards, note, source=None):
    slide = prs.slides.add_slide(blank)
    base(slide, number, title_value, subtitle, source)
    count = len(cards)
    if count == 2:
        cols, rows = 2, 1
    elif count <= 4:
        cols, rows = 2, 2
    else:
        cols, rows = 3, 2
    gap_x = 0.30
    left = 0.78
    usable_w = 11.78
    card_w = (usable_w - gap_x * (cols - 1)) / cols
    top = 1.48
    card_h = 4.80 if rows == 1 else 2.26
    gap_y = 0.30
    for i, card in enumerate(cards):
        col, row = i % cols, i // cols
        x = left + col * (card_w + gap_x)
        y = top + row * (card_h + gap_y)
        accent = card[2] if len(card) > 2 else CYAN
        rect(slide, x, y, card_w, card_h, WHITE, True, LINE)
        rect(slide, x, y, 0.11, card_h, accent, True, accent)
        circle(slide, x + 0.30, y + 0.28, 0.55, accent)
        text(slide, str(i + 1), x + 0.30, y + 0.34, 0.55, 0.30, 12, WHITE, True, PP_ALIGN.CENTER)
        title_size = 15.5 if rows == 2 else 17
        text(slide, card[0], x + 1.00, y + 0.24, card_w - 1.30, 0.54, title_size, INK, True)
        text(slide, card[1], x + 1.00, y + 0.84, card_w - 1.32, card_h - 1.10,
             12.8 if rows == 2 else 16.5, MUTED)
    add_notes(slide, note)


def add_process_slide(number, title_value, subtitle, steps, note, source=None):
    slide = prs.slides.add_slide(blank)
    base(slide, number, title_value, subtitle, source)
    n = len(steps)
    start_x = 0.58
    gap = 0.16
    width = (12.18 - gap * (n - 1)) / n
    y = 2.05
    for i, step in enumerate(steps):
        x = start_x + i * (width + gap)
        accent = step[2] if len(step) > 2 else [CYAN, BLUE, PURPLE, ORANGE, GREEN][i % 5]
        circle(slide, x + width / 2 - 0.32, 1.45, 0.64, accent)
        text(slide, str(i + 1), x + width / 2 - 0.32, 1.56, 0.64, 0.30,
             13, WHITE, True, PP_ALIGN.CENTER)
        rect(slide, x, y, width, 3.38, WHITE, True, LINE)
        text(slide, step[0], x + 0.22, y + 0.35, width - 0.44, 0.48,
             15.5, INK, True, PP_ALIGN.CENTER)
        text(slide, step[1], x + 0.25, y + 1.02, width - 0.50, 1.78,
             11.8, MUTED, align=PP_ALIGN.CENTER)
        if i < n - 1:
            text(slide, "→", x + width - 0.04, 3.30, gap + 0.10, 0.35,
                 17, BLUE, True, PP_ALIGN.CENTER)
    add_notes(slide, note)


def add_split_slide(number, title_value, subtitle, left_title, left_items,
                    right_title, right_items, note, source=None):
    slide = prs.slides.add_slide(blank)
    base(slide, number, title_value, subtitle, source)
    rect(slide, 0.72, 1.50, 5.92, 4.82, WHITE, True, LINE)
    rect(slide, 6.84, 1.50, 5.76, 4.82, NAVY, True)
    text(slide, left_title.upper(), 1.10, 1.84, 5.10, 0.35, 13, BLUE, True)
    bullets(slide, left_items, 1.10, 2.38, 5.02, 3.35, 16, INK, 9)
    text(slide, right_title.upper(), 7.24, 1.84, 4.90, 0.35, 13, CYAN, True)
    bullets(slide, right_items, 7.24, 2.38, 4.82, 3.35, 16, WHITE, 9)
    add_notes(slide, note)


def add_matrix_slide(number, title_value, subtitle, headers, rows, note,
                     source=None, col_widths=None):
    slide = prs.slides.add_slide(blank)
    base(slide, number, title_value, subtitle, source)
    left = 0.60
    top = 1.48
    total_w = 12.12
    if col_widths is None:
        col_widths = [total_w / len(headers)] * len(headers)
    x = left
    for i, header_value in enumerate(headers):
        rect(slide, x, top, col_widths[i], 0.58, NAVY if i == 0 else BLUE)
        text(slide, header_value, x + 0.10, top + 0.15, col_widths[i] - 0.20,
             0.25, 11.5, WHITE, True, PP_ALIGN.CENTER)
        x += col_widths[i]
    row_h = min(0.74, 4.52 / max(len(rows), 1))
    for r_idx, row in enumerate(rows):
        x = left
        y = top + 0.58 + r_idx * row_h
        fill = WHITE if r_idx % 2 == 0 else RGBColor(235, 242, 248)
        for c_idx, value in enumerate(row):
            rect(slide, x, y, col_widths[c_idx], row_h, fill, False, LINE)
            text(slide, value, x + 0.10, y + 0.10, col_widths[c_idx] - 0.20,
                 row_h - 0.16, 10.5, INK, c_idx == 0,
                 PP_ALIGN.LEFT, valign=MSO_ANCHOR.MIDDLE)
            x += col_widths[c_idx]
    add_notes(slide, note)


def add_code_slide(number, title_value, subtitle, code, side_title, side_items,
                   note, source=None):
    slide = prs.slides.add_slide(blank)
    base(slide, number, title_value, subtitle, source)
    rect(slide, 0.65, 1.46, 7.65, 4.95, NAVY, True)
    circle(slide, 0.95, 1.72, 0.16, RED)
    circle(slide, 1.22, 1.72, 0.16, ORANGE)
    circle(slide, 1.49, 1.72, 0.16, GREEN)
    text(slide, code, 0.98, 2.15, 6.95, 3.88, 11.5, WHITE, False, font="Menlo")
    rect(slide, 8.62, 1.46, 3.98, 4.95, WHITE, True, LINE)
    text(slide, side_title.upper(), 9.02, 1.86, 3.16, 0.32, 12, BLUE, True)
    bullets(slide, side_items, 9.02, 2.42, 3.05, 3.45, 13.8, INK, 9, check=True)
    add_notes(slide, note)


def add_activity_slide(number, title_value, instruction, deliverable, time_value, note):
    slide = prs.slides.add_slide(blank)
    rect(slide, 0, 0, 13.333, 7.5, NAVY)
    rect(slide, 0, 0, 13.333, 0.12, ORANGE)
    text(slide, "ACTIVIDAD PRÁCTICA", 0.72, 0.70, 3.4, 0.32, 13, ORANGE, True)
    text(slide, title_value, 0.72, 1.28, 8.45, 0.82, 29, WHITE, True)
    rect(slide, 0.72, 2.50, 8.15, 2.52, DARK_CARD, True)
    text(slide, "INSTRUCCIÓN", 1.10, 2.87, 1.75, 0.30, 11, CYAN, True)
    text(slide, instruction, 1.10, 3.35, 7.32, 1.18, 18, WHITE)
    rect(slide, 9.30, 1.28, 3.05, 3.74, ORANGE, True)
    text(slide, time_value, 9.55, 1.90, 2.55, 0.72, 31, NAVY, True, PP_ALIGN.CENTER)
    text(slide, "ENTREGABLE", 9.55, 2.92, 2.55, 0.26, 11, NAVY, True, PP_ALIGN.CENTER)
    text(slide, deliverable, 9.63, 3.40, 2.40, 1.12, 15, NAVY, True, PP_ALIGN.CENTER)
    text(slide, "Trabaja con una vacante, requisito o caso real.", 1.02, 5.68, 10.65, 0.42,
         17, PALE, True, PP_ALIGN.CENTER)
    footer(slide, number, "Taller QA · Aplicación inmediata")
    add_notes(slide, note)


def add_stats_slide(number, title_value, subtitle, stats, note, source):
    slide = prs.slides.add_slide(blank)
    base(slide, number, title_value, subtitle, source)
    for i, stat in enumerate(stats):
        x = 0.70 + i * 4.18
        color = stat[2]
        rect(slide, x, 1.65, 3.66, 4.45, WHITE, True, LINE)
        circle(slide, x + 1.18, 2.06, 1.30, color)
        text(slide, stat[0], x + 0.82, 2.42, 2.02, 0.42, 24, WHITE, True, PP_ALIGN.CENTER)
        text(slide, stat[1], x + 0.42, 3.68, 2.82, 1.48, 16, INK, True, PP_ALIGN.CENTER)
    add_notes(slide, note)


# 01
add_cover(
    "La transformación del Software QA Engineering",
    "IA, automatización y empleabilidad en un mercado cambiante",
    "Modalidad virtual · 2 de julio de 2026\n6:30 p. m. – 9:30 p. m. · Duración: 3 horas\nEnfoque práctico, participativo y orientado al empleo",
    """[Tiempo sugerido: 3 minutos]
Buenas tardes y gracias por acompañarme. Hoy no vengo a decirles que la inteligencia artificial eliminará el trabajo de QA ni que todos deben convertirse mañana en programadores de automatización. Vengo a mostrarles cómo está cambiando el rol, qué capacidades siguen siendo fundamentales y cómo construir una ruta profesional basada en evidencias reales del mercado.

Durante estas tres horas vamos a conectar cuatro ideas: calidad, inteligencia artificial, automatización y empleabilidad. El resultado esperado es que cada participante identifique un puesto objetivo, reconozca sus brechas y pueda comenzar una guía de estudio autónoma basada en las vacantes a las que realmente desea aplicar."""
)

# 02
add_cards_slide(2, "¿Qué te llevarás de este taller?",
    "El objetivo no es aprender todas las herramientas; es saber decidir qué aprender y demostrarlo.",
    [
        ("Una dirección profesional", "Elegirás un rol objetivo dentro de QA y reconocerás qué exige el mercado para ese rol.", CYAN),
        ("Una ruta de aprendizaje", "Convertirás requisitos de vacantes reales en una guía de estudio organizada por prioridades.", BLUE),
        ("Evidencias de capacidad", "Sabrás qué proyectos, documentos y automatizaciones publicar para demostrar lo que sabes hacer.", GREEN),
        ("Un plan de acción", "Cerrarás con acciones para buscar empleo, prepararte para entrevistas y mantenerte actualizado.", ORANGE),
    ],
    """[Tiempo sugerido: 3 minutos]
Explique que este taller no pretende cubrir de manera profunda todas las áreas de QA. El propósito es entregar un mapa para tomar mejores decisiones. Muchas personas estudian cursos aislados y acumulan certificados, pero no construyen una dirección profesional ni evidencias visibles.

Al finalizar, cada participante deberá poder responder cuatro preguntas: ¿a qué puesto quiero aplicar?, ¿qué competencias aparecen repetidamente en esas vacantes?, ¿qué me falta por aprender? y ¿cómo demostraré que ya puedo realizar ese trabajo? Estas preguntas serán el hilo conductor de toda la sesión."""
)

# 03
add_process_slide(3, "Ruta de las tres horas", "Una sesión con explicación, práctica y decisiones personales.",
    [
        ("6:30–7:00", "Mercado y transformación del rol QA", CYAN),
        ("7:00–7:35", "IA aplicada al trabajo y al testing", BLUE),
        ("7:35–8:20", "Automatización y portafolio", PURPLE),
        ("8:20–8:30", "Pausa breve", ORANGE),
        ("8:30–9:30", "Empleabilidad, estudio autónomo y sorteo", GREEN),
    ],
    """[Tiempo sugerido: 2 minutos]
Presente la agenda y explique que habrá pequeñas actividades. Recomiende tener abierto un navegador, LinkedIn o el portal de empleo que utilice cada participante, una hoja de cálculo o documento y, si es posible, su CV actual.

Indique que la sesión comienza con el mercado porque una ruta de estudio debe responder a un destino. Luego se analizará cómo la IA y la automatización cambian el trabajo. Finalmente, se traducirá todo en acciones concretas de búsqueda de empleo y actualización continua."""
)

# 04 · Dinámica de apertura; corresponde a la diapositiva 3 de la versión compacta.
add_activity_slide(4, "Conozcámonos: ¿desde dónde participas?",
    "Responde por el chat o levanta la mano si deseas compartir: ¿desde dónde nos escuchas?, ¿cuál es tu nivel de seniority?, ¿cuál es tu profesión o área actual? y ¿qué esperas de este taller?",
    "Ubicación + seniority + profesión + expectativa", "5 min",
    """[Tiempo sugerido: 5 minutos]
Invite a los participantes a responder por el chat. También pregunte quién se anima a levantar la mano y presentarse brevemente. Cada intervención debe incluir cuatro elementos: desde qué ciudad, provincia o país participa; cuál considera que es su nivel de seniority; cuál es su profesión, ocupación o área tecnológica actual; y qué espera obtener de este taller.

Puede orientar el nivel de seniority con ejemplos sencillos: persona que está iniciando o explorando el área, junior, semisenior, senior, líder o especialista. Aclare que no se trata de examinar ni comparar a nadie, sino de conocer la diversidad del grupo y adaptar mejor la conversación.

Escuche las expectativas: búsqueda del primer empleo, transición hacia QA, automatización, uso de inteligencia artificial, actualización profesional, preparación para entrevistas, networking o crecimiento hacia liderazgo. Utilice las respuestas para destacar durante el taller los temas que más necesita el grupo.

Si hay muchas personas principiantes, explique cada término con ejemplos. Si hay participantes con experiencia, invítelos a compartir herramientas, aprendizajes y retos. Cierre diciendo que el punto de partida no limita el destino profesional; solamente ayuda a identificar el próximo paso realista."""
)

# 05
add_section(5, "I", "El mercado cambió; el valor del QA también",
    "La calidad deja de ser una fase final y se convierte en una capacidad transversal del producto.",
    """[Tiempo sugerido: 1 minuto]
Introduzca la primera sección. Explique que antes de hablar de herramientas debemos comprender qué está comprando una empresa cuando contrata a una persona de QA: reducción de riesgo, confianza para liberar cambios, claridad sobre el comportamiento del producto y capacidad de prevenir problemas.""", CYAN)

# 06
add_stats_slide(6, "Tres señales del mercado laboral",
    "Las herramientas cambian, pero la capacidad de aprender y razonar sigue ganando valor.",
    [
        ("39%", "de las competencias esenciales podrían cambiar hacia 2030 según empleadores consultados.", CYAN),
        ("#1", "el pensamiento analítico continúa como competencia central para los empleadores.", BLUE),
        ("IA", "la alfabetización en IA aparece entre las capacidades de crecimiento más rápido.", GREEN),
    ],
    """[Tiempo sugerido: 5 minutos]
Explique que estas cifras no significan que el 39 % de las personas perderá su trabajo. Significan que una parte importante de lo que hoy hacemos y de las herramientas que usamos cambiará. Por eso, memorizar una interfaz no es suficiente.

El pensamiento analítico se conecta directamente con QA: formular preguntas, detectar inconsistencias, priorizar riesgos y explicar evidencias. La alfabetización en IA tampoco significa desarrollar modelos complejos; significa comprender sus posibilidades, limitaciones y uso responsable dentro del trabajo.

Concluya: el profesional sostenible combina competencias técnicas, criterio humano y capacidad de aprendizaje.""",
    "Fuentes: World Economic Forum, Future of Jobs Report 2025; LinkedIn, Skills on the Rise 2025.")

# 07
add_split_slide(7, "QA no desaparece: cambia el tipo de trabajo",
    "Las tareas repetitivas se aceleran; el criterio y la responsabilidad aumentan.",
    "Trabajo que se automatiza o acelera", [
        "Generación inicial de casos y datos de prueba",
        "Ejecución repetitiva de regresiones estables",
        "Resumen de logs y documentación",
        "Creación de borradores de reportes",
    ],
    "Trabajo que gana valor", [
        "Análisis de riesgo y ambigüedades",
        "Pruebas exploratorias y pensamiento crítico",
        "Comprensión del negocio y del usuario",
        "Decisiones éticas, privacidad y calidad",
    ],
    """[Tiempo sugerido: 5 minutos]
Evite presentar la IA como sustitución automática del QA. Explique que se automatizan tareas, no responsabilidades completas. Una herramienta puede proponer cien casos, pero alguien debe decidir cuáles son relevantes, qué riesgo cubren, qué datos pueden utilizarse y qué evidencia es suficiente.

Use un ejemplo: si una IA genera pruebas para una transferencia bancaria, todavía necesitamos a una persona que conozca reglas, límites, permisos, fraude, experiencia del cliente y consecuencias de un error. El valor se desplaza desde ejecutar pasos mecánicos hacia diseñar una estrategia de calidad y tomar decisiones sustentadas."""
)

# Franja complementaria de la diapositiva 5 compacta: automatización fuera de TI.
qa_change_slide = prs.slides[-1]
rect(qa_change_slide, 0.92, 5.18, 11.48, 1.02, ORANGE, True, ORANGE)
text(qa_change_slide, "TAMBIÉN SE AUTOMATIZAN TAREAS EN PROFESIONES NO TECNOLÓGICAS",
     1.20, 5.32, 10.92, 0.24, 11.5, NAVY, True, PP_ALIGN.CENTER)
text(qa_change_slide,
     "Administración: agenda y documentos  ·  Finanzas: facturas y conciliaciones  ·  RR. HH.: filtros y onboarding\n"
     "Educación: borradores y cuestionarios  ·  Salud: citas y transcripción  ·  Ventas/mercadeo: mensajes y segmentación",
     1.15, 5.65, 11.02, 0.42, 10.2, NAVY, True, PP_ALIGN.CENTER)
qa_change_notes = qa_change_slide.notes_slide.notes_text_frame
if qa_change_notes is not None:
    qa_change_notes.text = qa_change_notes.text.rstrip() + """

La automatización no se limita al desarrollo de software o al aseguramiento de la calidad. En administración se automatizan la programación de citas, clasificación de documentos, transcripción de reuniones y primeros borradores de correos. En contabilidad y finanzas se aceleran la captura de facturas, conciliaciones, clasificación de gastos y alertas sobre transacciones inusuales.

En recursos humanos se apoyan el filtrado inicial de información, la coordinación de entrevistas, la preparación de documentos de incorporación y las respuestas a preguntas frecuentes. En educación pueden generarse borradores de planificación, bancos iniciales de preguntas, rúbricas y resúmenes de resultados, pero el docente conserva la responsabilidad de adaptar, evaluar y acompañar.

En salud se automatizan tareas administrativas como citas, recordatorios, transcripción y organización documental; las decisiones clínicas requieren profesionales autorizados y revisión rigurosa. En ventas, mercadeo y servicio al cliente se automatizan la segmentación, los borradores de mensajes, la clasificación de prospectos y las respuestas iniciales. En el ámbito jurídico también pueden compararse documentos y localizar cláusulas, pero la interpretación y el asesoramiento permanecen bajo responsabilidad profesional.

El patrón se repite: la tecnología acelera tareas predecibles, repetitivas y basadas en información. Ganan valor la supervisión, el criterio del dominio, la comunicación humana, la protección de datos, la atención de excepciones y la responsabilidad sobre la decisión final. No debemos preguntar solamente cuál profesión será automatizada, sino cuáles tareas cambiarán y qué competencias humanas y digitales necesitaremos para dirigir ese cambio."""

# 08
add_process_slide(8, "De probar al final a construir calidad",
    "El QA moderno participa en todo el ciclo de desarrollo.",
    [
        ("Descubrir", "Preguntar por usuarios, valor, reglas y riesgos.", CYAN),
        ("Definir", "Aclarar requisitos y criterios de aceptación.", BLUE),
        ("Construir", "Colaborar con desarrollo y prevenir defectos.", PURPLE),
        ("Validar", "Ejecutar pruebas funcionales y no funcionales.", ORANGE),
        ("Observar", "Revisar producción, métricas y aprendizaje.", GREEN),
    ],
    """[Tiempo sugerido: 5 minutos]
Recorra el ciclo de izquierda a derecha. Destaque que encontrar un defecto antes de programar suele ser más barato que encontrarlo en producción. El QA puede aportar desde el refinamiento de una historia de usuario preguntando por casos alternos, permisos, errores, límites y datos.

Explique el concepto de shift left sin convertirlo en un eslogan: participar antes. Añada shift right: aprender también de producción mediante monitoreo, métricas y comportamiento real. El QA moderno conecta prevención, validación y observabilidad."""
)

# 09
add_cards_slide(9, "El ecosistema profesional de QA",
    "No todas las vacantes piden lo mismo. Primero identifica la familia de roles que te interesa.",
    [
        ("QA manual / funcional", "Requisitos, casos, exploración, reporte de defectos, evidencia y negocio.", CYAN),
        ("API y datos", "HTTP, JSON, Postman, contratos, consultas SQL y consistencia de datos.", BLUE),
        ("QA automation", "Programación, framework, patrones, Git, CI y mantenimiento de pruebas.", PURPLE),
        ("Performance / seguridad", "Carga, observabilidad, riesgos, OWASP y herramientas especializadas.", ORANGE),
        ("SDET / quality engineering", "Arquitectura de pruebas, calidad del código, pipelines y estrategia.", GREEN),
        ("AI testing", "Datos, sesgo, comportamiento probabilístico, modelos, red teaming y monitoreo.", TEAL),
    ],
    """[Tiempo sugerido: 6 minutos]
Explique brevemente cada familia. El objetivo no es obligar a elegir una especialización definitiva, sino evitar una guía de estudio infinita. Una vacante junior manual puede priorizar documentación, criterios de aceptación, Jira y SQL básico. Una vacante de automation puede priorizar lenguaje de programación, Playwright o Selenium, API, Git y CI.

Mencione que los nombres cambian entre empresas. Por eso no basta con buscar un solo título. Se deben revisar funciones y requisitos. Recomiende guardar vacantes con títulos distintos pero responsabilidades parecidas."""
)

# 10
add_cards_slide(10, "¿A qué otras especialidades puedes evolucionar?",
    "La experiencia en QA desarrolla capacidades transferibles, pero cada ruta requiere preparación adicional.",
    [
        ("Análisis de negocio", "Requisitos, procesos, reglas, criterios de aceptación y comunicación con usuarios.", CYAN),
        ("Producto", "Product Owner, product analyst o product operations con foco en valor y experiencia.", BLUE),
        ("DevOps / release", "Pipelines, ambientes, observabilidad, despliegues, confiabilidad y mejora operativa.", PURPLE),
        ("Ciberseguridad", "Application security, pruebas de seguridad, riesgos, OWASP y DevSecOps.", ORANGE),
        ("Datos e IA", "Data quality, ETL testing, validación de modelos, sesgo, drift y sistemas generativos.", TEAL),
        ("Liderazgo de calidad", "QA Lead, Test Manager, Quality Coach, consultoría y estrategia organizacional.", GREEN),
    ],
    """[Tiempo sugerido: 4 minutos]
Explique que QA puede ser una puerta de entrada a múltiples trayectorias. La persona desarrolla pensamiento analítico, comprensión del negocio, documentación, comunicación y visión de riesgo. Estas capacidades son valiosas en análisis de negocio, producto, DevOps, seguridad, datos y liderazgo.

Sin embargo, evite prometer una transición automática. Para llegar a análisis de negocio habrá que profundizar en procesos y levantamiento de requisitos. Para DevOps se necesitan sistemas, nube, automatización e infraestructura. Para seguridad se requieren fundamentos específicos. Para producto se debe comprender estrategia, usuarios y métricas.

También existen rutas como desarrollo de software, ingeniería de confiabilidad, accesibilidad, UX research, soporte técnico avanzado, gestión de proyectos y documentación técnica. La decisión debe surgir de las tareas que la persona disfruta, sus fortalezas y las vacantes disponibles.""",
    "Rutas cercanas: negocio, producto, operaciones, seguridad, datos, IA y liderazgo.")

# 11
add_activity_slide(11, "Elige un rol objetivo, no diez cursos",
    "Selecciona un puesto que te gustaría alcanzar en los próximos 6 a 12 meses. Escribe el título, nivel, tipo de empresa y tres tareas que esperas realizar.",
    "Ficha de puesto objetivo", "5 min",
    """[Tiempo sugerido: 5 minutos]
Pida que cada participante escriba una ficha mínima. Ejemplo: QA Automation Junior, empresa de servicios digitales, modalidad remota o híbrida; tareas: automatizar flujos web, probar APIs y ejecutar pruebas en CI.

Explique que esta elección es provisional. Se puede ajustar después de analizar vacantes. Lo importante es comenzar con una hipótesis profesional concreta. Sin destino, cualquier curso parece urgente; con destino, se puede priorizar."""
)

# 12
add_section(12, "II", "Inteligencia artificial aplicada a QA",
    "Usar IA para ampliar cobertura y velocidad sin delegar el criterio profesional.",
    """[Tiempo sugerido: 1 minuto]
Introduzca la segunda sección. Diferencie dos campos: utilizar IA para apoyar actividades de testing y probar productos que incorporan IA. Ambos aparecen cada vez más en el desarrollo profesional.""", BLUE)

# 13
add_process_slide(13, "¿Dónde puede ayudar la IA?",
    "Un asistente acelera borradores; el QA valida relevancia, riesgo y exactitud.",
    [
        ("Requisitos", "Detectar ambigüedades, reglas faltantes y preguntas.", CYAN),
        ("Diseño", "Proponer escenarios, datos y casos límite.", BLUE),
        ("Ejecución", "Explicar errores, resumir logs y comparar resultados.", PURPLE),
        ("Reporte", "Mejorar claridad, pasos, impacto y evidencia.", ORANGE),
        ("Aprendizaje", "Crear ejercicios, explicar código y practicar entrevistas.", GREEN),
    ],
    """[Tiempo sugerido: 5 minutos]
Explique que el mejor uso inicial de IA ocurre en tareas con revisión clara. Por ejemplo, pedir preguntas sobre una historia de usuario, proponer datos de prueba o mejorar la redacción de un bug.

Evite copiar resultados directamente. Una salida extensa puede crear sensación de productividad sin aportar cobertura útil. Recomiende solicitar clasificación por riesgo y justificar por qué cada escenario importa. El QA debe verificar reglas de negocio, seguridad, privacidad y exactitud técnica.""",
    "Referencia: ISTQB Certified Tester – Testing with Generative AI, versión 1.1.")

# 14
add_code_slide(14, "Prompt útil para analizar requisitos",
    "Un buen prompt define contexto, objetivo, restricciones, criterios y formato de salida.",
    """ROL: Actúa como QA analyst de una aplicación bancaria.

CONTEXTO: Historia de usuario sobre transferencias entre cuentas.

OBJETIVO: Detecta ambigüedades, riesgos y criterios faltantes.

RESTRICCIONES: No inventes reglas. Marca cada supuesto.

SALIDA:
1. Preguntas al negocio
2. Riesgos priorizados
3. Escenarios positivos y negativos
4. Datos de prueba
5. Criterios de aceptación sugeridos""",
    "Antes de usar la salida", [
        "Contrasta con el requisito original",
        "Elimina supuestos no confirmados",
        "Prioriza según riesgo",
        "Protege datos sensibles",
        "Documenta tu decisión final",
    ],
    """[Tiempo sugerido: 6 minutos]
Lea el prompt por bloques. Explique que el rol aporta perspectiva, el contexto evita respuestas genéricas, el objetivo define la tarea, las restricciones reducen invenciones y el formato facilita revisar la salida.

Haga énfasis en la frase «no inventes reglas». Aun así, el modelo puede asumir información. Por eso se solicita marcar supuestos y luego se valida con negocio. Invite a adaptar este prompt a comercio electrónico, educación, salud o cualquier dominio."""
)

# 15
add_cards_slide(15, "La estructura de un prompt profesional",
    "Preguntar mejor no reemplaza saber de QA; obliga a expresar el razonamiento.",
    [
        ("Contexto", "Producto, usuario, requisito, ambiente y dominio.", CYAN),
        ("Tarea", "La acción precisa: analizar, diseñar, comparar, explicar o revisar.", BLUE),
        ("Criterios", "Riesgo, cobertura, accesibilidad, seguridad, calidad y prioridad.", PURPLE),
        ("Restricciones", "Qué no debe inventar, datos prohibidos y alcance.", ORANGE),
        ("Formato", "Tabla, checklist, Gherkin, JSON, código o reporte.", GREEN),
        ("Verificación", "Cómo revisarás la salida y qué evidencia aceptarás.", TEAL),
    ],
    """[Tiempo sugerido: 4 minutos]
Explique cada componente con rapidez. Destaque que la verificación debe planificarse desde el prompt. Si se pide código, se debe ejecutar. Si se piden casos, se deben contrastar con requisitos y riesgos. Si se pide una explicación, se deben comprobar conceptos en documentación oficial.

Una señal de madurez es poder explicar por qué se aceptó, modificó o rechazó una sugerencia de IA."""
)

# 16
add_split_slide(16, "Uso responsable: cinco límites que no se negocian",
    "La productividad nunca justifica comprometer información o calidad.",
    "Riesgos frecuentes", [
        "Alucinaciones y respuestas plausibles pero falsas",
        "Exposición de credenciales o datos personales",
        "Sesgos en ejemplos, datos o decisiones",
        "Código inseguro o difícil de mantener",
        "Dependencia sin aprendizaje real",
    ],
    "Controles profesionales", [
        "Anonimizar y minimizar los datos",
        "Verificar contra fuentes y ejecución",
        "Revisar seguridad y licencias",
        "Conservar trazabilidad de cambios",
        "Mantener decisión y responsabilidad humana",
    ],
    """[Tiempo sugerido: 5 minutos]
Advierta claramente: nunca pegar credenciales, datos de clientes, código privado o información institucional sensible en una herramienta sin autorización. Use ejemplos ficticios o anonimizados.

Explique que validar no significa leer rápidamente y decir «parece correcto». Validar puede implicar ejecutar, consultar documentación oficial, comparar con el requisito y pedir revisión de otra persona. La IA debe fortalecer el aprendizaje, no ocultar que no comprendemos lo producido.""",
    "Referencia: ISTQB CT-GenAI; OWASP Top 10 for LLM Applications 2025.")

# 17
add_cards_slide(17, "Probar sistemas con IA exige nuevas preguntas",
    "La salida puede ser probabilística; la calidad necesita criterios estadísticos, éticos y contextuales.",
    [
        ("Datos", "Representatividad, calidad, etiquetas, privacidad, sesgo y trazabilidad.", CYAN),
        ("Modelo", "Métricas, robustez, variabilidad, ataques, drift y comparación.", BLUE),
        ("Sistema", "Integración, latencia, seguridad, experiencia y fallback.", PURPLE),
        ("Generativa", "Alucinaciones, toxicidad, prompt injection, red teaming y uso indebido.", ORANGE),
    ],
    """[Tiempo sugerido: 5 minutos]
Explique que probar una calculadora y probar un modelo generativo no es lo mismo. Una salida puede variar aun con entradas parecidas. A veces no existe una única respuesta correcta. Se necesitan conjuntos de evaluación, umbrales, revisión humana y análisis estadístico.

Mencione que ISTQB actualizó en 2026 su syllabus de AI Testing con foco en datos, modelos, ciclo de machine learning, sistemas generativos y red teaming. Esto muestra una posible ruta de especialización para perfiles de QA.""",
    "Fuente: ISTQB Certified Tester AI Testing, syllabus 2.0, 2026.")

# 18
add_activity_slide(18, "Revisa una salida de IA como QA",
    "Toma un caso de prueba generado por IA. Identifica: un supuesto, un riesgo no cubierto, un dato sensible, un resultado esperado ambiguo y una mejora necesaria.",
    "Checklist de validación", "6 min",
    """[Tiempo sugerido: 6 minutos]
Puede mostrar un caso generado previamente o pedir a los participantes que utilicen uno propio. Dé dos minutos para revisión individual y luego solicite ejemplos.

Concluya que la habilidad no consiste en obtener una respuesta larga, sino en inspeccionarla. Esa capacidad de revisión es transferible a reportes, código, datos y documentación. El QA debe convertirse en un usuario exigente de la IA."""
)

# 19
add_section(19, "III", "Automatización con propósito",
    "Automatizar para obtener retroalimentación rápida y confiable, no para acumular scripts.",
    """[Tiempo sugerido: 1 minuto]
Introduzca la tercera sección. Explique que la automatización es una inversión: tiene costos de diseño, datos, ambientes, mantenimiento y análisis. Por eso debe priorizarse según valor y riesgo.""", PURPLE)

# 20
add_split_slide(20, "¿Qué conviene automatizar?",
    "La mejor candidata es una prueba valiosa, repetible, estable y suficientemente clara.",
    "Buenos candidatos", [
        "Regresiones críticas y frecuentes",
        "Validaciones de API y contratos",
        "Flujos estables con resultados claros",
        "Combinaciones de datos repetitivas",
        "Pruebas necesarias en cada integración",
    ],
    "Automatizar con cautela", [
        "Experiencia visual altamente subjetiva",
        "Funcionalidad que cambia cada semana",
        "Escenarios de una sola ejecución",
        "Pruebas sin criterio de éxito claro",
        "Casos que requieren juicio exploratorio",
    ],
    """[Tiempo sugerido: 5 minutos]
Presente una regla simple: automatizar no es grabar todos los pasos posibles. Es seleccionar controles que aporten retroalimentación repetible. Una prueba que cambia constantemente puede consumir más tiempo de mantenimiento que el valor que entrega.

Subraye que la exploración humana y la automatización se complementan. La automatización libera tiempo para investigar riesgos nuevos; no elimina la necesidad de pensar."""
)

# 21
add_process_slide(21, "Ruta técnica progresiva",
    "Aprende por capas; cada capa debe producir una evidencia visible.",
    [
        ("Manual", "Casos, bugs, exploración y riesgo.", CYAN),
        ("API", "HTTP, JSON, Postman y contratos.", BLUE),
        ("Datos", "SQL, consultas y validación.", TEAL),
        ("Código", "Un lenguaje, Git y buenas prácticas.", PURPLE),
        ("Automatización", "Framework, patrones, reportes y CI.", GREEN),
    ],
    """[Tiempo sugerido: 5 minutos]
Explique que no todas las personas necesitan esperar a dominar cada capa al cien por ciento. Sin embargo, avanzar sin fundamentos crea automatizaciones frágiles.

Recomiende elegir un solo lenguaje al inicio según las vacantes objetivo. Si predominan Playwright y TypeScript, esa combinación puede ser razonable. Si el entorno usa Java y Selenium, la ruta será distinta. La vacante y el ecosistema de la empresa deben orientar la elección."""
)

# 22
add_matrix_slide(22, "Stack práctico para un portafolio QA",
    "Las herramientas son ejemplos; selecciona las que se repiten en tus vacantes objetivo.",
    ["Capacidad", "Herramientas posibles", "Evidencia en portafolio"],
    [
        ("Gestión", "Jira, Trello, Qase", "Plan, casos y defectos"),
        ("API", "Postman, Swagger", "Colección con pruebas"),
        ("Datos", "SQL", "Consultas y validaciones"),
        ("UI web", "Playwright, Cypress, Selenium", "Suite automatizada"),
        ("Código", "JavaScript/TS, Java, Python, C#", "Repositorio legible"),
        ("CI", "GitHub Actions, GitLab CI", "Ejecución automática"),
        ("IA", "Asistentes autorizados", "Prompts y revisión documentada"),
    ],
    """[Tiempo sugerido: 5 minutos]
Recorra la tabla y enfatice la tercera columna. No basta con escribir «Postman» en el CV. La persona debe mostrar una colección, variables, aserciones, documentación y ejecución.

Tampoco conviene aprender tres frameworks de UI simultáneamente. Es mejor un proyecto sólido, mantenible y explicado que tres tutoriales copiados. La evidencia debe permitir que otra persona entienda el problema, ejecute las pruebas y observe el resultado.""",
    col_widths=[2.20, 4.00, 5.92])

# 23
add_code_slide(23, "Ejemplo mínimo con Playwright",
    "Una prueba profesional expresa intención, usa localizadores estables y produce evidencia.",
    """import { test, expect } from '@playwright/test';

test('usuario válido accede al panel', async ({ page }) => {
  await page.goto('/login');
  await page.getByLabel('Correo').fill('qa@example.com');
  await page.getByLabel('Contraseña').fill('ClaveSegura1!');
  await page.getByRole('button', { name: 'Entrar' }).click();

  await expect(page.getByRole('heading', {
    name: 'Panel principal'
  })).toBeVisible();
});""",
    "Qué observar", [
        "Nombre orientado al comportamiento",
        "Localizadores accesibles",
        "Resultado esperado explícito",
        "Datos separados del código real",
        "Reporte y trazas ante fallos",
    ],
    """[Tiempo sugerido: 6 minutos]
Explique el código sin asumir experiencia avanzada. La prueba abre la página, completa datos, ejecuta una acción y verifica un resultado observable. Destaque que en un proyecto real no se dejarían credenciales en el código; se usarían variables seguras y datos de prueba.

Mencione que Playwright ofrece auto-waiting, aserciones reintentables, trazas, ejecución paralela y soporte para varios navegadores. Aclare que una herramienta reduce ciertos problemas, pero no corrige un mal diseño de pruebas.""",
    "Fuente técnica: documentación oficial de Playwright.")

# 24
add_process_slide(24, "De tu computadora al pipeline",
    "La automatización gana valor cuando se ejecuta de manera repetible y visible.",
    [
        ("Commit", "El cambio se envía al repositorio.", CYAN),
        ("Build", "Se instala y prepara la aplicación.", BLUE),
        ("Tests", "Se ejecutan controles automáticos.", PURPLE),
        ("Report", "Se publican resultados y evidencias.", ORANGE),
        ("Decision", "El equipo evalúa si puede avanzar.", GREEN),
    ],
    """[Tiempo sugerido: 5 minutos]
Explique integración continua de manera sencilla. Cada cambio puede activar un flujo que instala dependencias, construye, ejecuta pruebas y publica resultados. Esto reduce la dependencia de una sola computadora y hace visible el estado del producto.

Para portafolio, GitHub Actions permite mostrar que la suite se ejecuta automáticamente. La insignia del pipeline no es el objetivo; el objetivo es demostrar una práctica profesional reproducible.""",
    "Fuente: GitHub Docs, Building and testing your code with GitHub Actions.")

# 25
add_cards_slide(25, "Proyecto de portafolio que sí cuenta una historia",
    "Un solo producto puede demostrar análisis, ejecución, automatización y comunicación.",
    [
        ("1. Contexto", "Producto seleccionado, usuarios, alcance, riesgos y estrategia de prueba.", CYAN),
        ("2. Evidencia manual", "Casos, checklist exploratorio, defectos y capturas organizadas.", BLUE),
        ("3. API y datos", "Colección de Postman, aserciones, ambientes y consultas SQL.", TEAL),
        ("4. Automatización", "Flujos críticos, estructura mantenible, reportes y manejo de datos.", PURPLE),
        ("5. Integración continua", "Pipeline que ejecuta pruebas y conserva resultados.", ORANGE),
        ("6. README", "Cómo ejecutar, decisiones, limitaciones, aprendizajes y próximos pasos.", GREEN),
    ],
    """[Tiempo sugerido: 6 minutos]
Explique que el portafolio debe contar una historia completa. Puede utilizarse una aplicación pública de práctica, un proyecto personal o un sistema académico autorizado. Nunca publicar información confidencial de una empresa.

El README es fundamental: debe explicar qué problema se probó, qué riesgos se priorizaron, cómo ejecutar el proyecto y qué decisiones tomó la persona. Esto permite evaluar criterio, no solo código.""",
    "Referencia práctica: GitHub Docs sobre perfil, README, repositorios fijados y GitHub Actions.")

# 26
add_activity_slide(26, "Diseña tu evidencia de automatización",
    "Define un proyecto pequeño: producto, tres riesgos, un flujo crítico, una API, una consulta SQL y una ejecución automática. No escribas código todavía; diseña la evidencia.",
    "Borrador de proyecto QA", "7 min",
    """[Tiempo sugerido: 7 minutos]
Dé cinco minutos para diseño y dos para compartir. Corrija proyectos demasiado grandes. Recomiende un alcance que pueda completarse en dos a cuatro semanas.

Pregunte: ¿qué capacidad demostrará cada elemento? Si una actividad no aporta evidencia relacionada con la vacante objetivo, puede posponerse. Concluya que el portafolio es una prueba de trabajo, no un depósito de archivos."""
)

# 27
add_section(27, "PAUSA", "Respira y regresa con una vacante abierta",
    "Al volver trabajaremos búsqueda de empleo, guía de estudio autónoma, CV, LinkedIn y entrevistas.",
    """[Tiempo sugerido: 10 minutos]
Indique la hora exacta de regreso. Pida a cada participante localizar durante la pausa al menos una vacante que le interese. Puede ser local, remota, de pasantía, junior o de un nivel que aspire a alcanzar.""", ORANGE)
pause_slide = prs.slides[-1]
for pause_shape in pause_slide.shapes:
    if not getattr(pause_shape, "has_text_frame", False) or pause_shape.text.strip() != "PAUSA":
        continue
    pause_shape.width = Inches(2.30)
    pause_shape.height = Inches(0.90)
    pause_run = pause_shape.text_frame.paragraphs[0].runs[0]
    pause_run.font.size = Pt(47)

# 28
add_section(28, "IV", "Empleabilidad y aprendizaje autónomo",
    "Estudiar lo que el mercado pide, practicarlo y convertirlo en evidencia profesional.",
    """[Tiempo sugerido: 1 minuto]
Introduzca el bloque central de empleabilidad. Explique que buscar empleo no comienza cuando el CV está listo. Comienza observando el mercado, seleccionando un objetivo y creando evidencia alineada con ese objetivo.""", GREEN)

# 29
add_split_slide(29, "Primero busca la vacante; después diseña el estudio",
    "La autonomía no es estudiar sin dirección. Es tomar decisiones informadas sobre qué aprender.",
    "Ruta poco efectiva", [
        "Comprar cursos por impulso",
        "Aprender muchas herramientas superficialmente",
        "Copiar proyectos sin comprenderlos",
        "Usar el mismo CV para todas las vacantes",
        "Esperar sentirse totalmente preparado",
    ],
    "Ruta orientada al mercado", [
        "Elegir un rol y nivel objetivo",
        "Analizar entre 10 y 20 vacantes",
        "Detectar competencias repetidas",
        "Priorizar brechas por frecuencia e impacto",
        "Crear evidencias y comenzar a aplicar",
    ],
    """[Tiempo sugerido: 5 minutos]
Esta es una de las ideas más importantes del taller. Explique que una guía de estudio no debe comenzar con una lista genérica de tecnologías. Debe comenzar con el trabajo que se desea realizar.

La persona busca vacantes reales, aunque todavía no cumpla todos los requisitos. Extrae tareas, herramientas, competencias y nivel. Luego agrupa patrones. Si Playwright aparece una vez y Selenium aparece en quince vacantes objetivo, esa diferencia es una señal. No es una orden absoluta, pero orienta la prioridad."""
)

# 30
add_process_slide(30, "Método VACANTE para estudiar con autonomía",
    "Un ciclo repetible para alinear aprendizaje, práctica y búsqueda de empleo.",
    [
        ("V", "Visualiza el puesto y nivel objetivo.", CYAN),
        ("A", "Analiza vacantes reales y tareas.", BLUE),
        ("C", "Clasifica competencias y frecuencia.", PURPLE),
        ("A", "Autoevalúa tus brechas con evidencia.", ORANGE),
        ("NTE", "Nivela, trabaja y evidencia mientras aplicas.", GREEN),
    ],
    """[Tiempo sugerido: 6 minutos]
Presente el método como una ayuda mnemotécnica. Visualizar significa definir el destino. Analizar significa leer funciones, no solo títulos. Clasificar implica separar fundamentos, herramientas, habilidades humanas, experiencia y requisitos deseables.

Autoevaluar no consiste en decir «creo que sé». Debe preguntarse: ¿puedo explicarlo?, ¿puedo hacerlo sin tutorial?, ¿puedo mostrar una evidencia? Finalmente, nivelar, trabajar y evidenciar significa estudiar en ciclos breves, construir proyectos y aplicar sin esperar perfección."""
)

# 31
add_matrix_slide(31, "Matriz de vacantes: convierte anuncios en datos",
    "Ejemplo de análisis. La prioridad se define por repetición, importancia y brecha personal.",
    ["Competencia", "Frecuencia", "Mi nivel", "Evidencia actual", "Prioridad"],
    [
        ("Casos y bugs", "9/10", "Medio", "Proyecto manual", "Alta"),
        ("API / Postman", "8/10", "Bajo", "Ninguna", "Muy alta"),
        ("SQL", "7/10", "Bajo", "Ejercicios", "Alta"),
        ("Git / GitHub", "7/10", "Medio", "2 repositorios", "Media"),
        ("Playwright", "5/10", "Bajo", "Tutorial", "Media"),
        ("Inglés técnico", "6/10", "Medio", "Lectura", "Media"),
    ],
    """[Tiempo sugerido: 7 minutos]
Explique las columnas. Frecuencia indica cuántas vacantes mencionan la competencia. Mi nivel debe basarse en una escala clara: no conozco, comprendo, practico con apoyo, trabajo de forma autónoma o puedo enseñar. Evidencia actual indica qué se puede mostrar.

La prioridad no depende solo de frecuencia. Una competencia puede aparecer pocas veces, pero ser obligatoria para el rol. También se considera el tiempo disponible y las dependencias: antes de automatizar APIs, quizá se necesite comprender HTTP y JavaScript.""",
    col_widths=[2.30, 1.55, 1.55, 3.75, 2.97])

# 32
add_activity_slide(32, "Construye tu matriz con una vacante real",
    "Extrae de la vacante: cinco tareas, cinco herramientas o conocimientos, dos habilidades humanas y cualquier requisito de experiencia. Marca qué puedes demostrar hoy.",
    "Primera fila de tu matriz", "8 min",
    """[Tiempo sugerido: 8 minutos]
Pida trabajar con la vacante localizada durante la pausa. Si alguien no tiene una, puede buscar «QA junior», «software tester», «QA analyst», «QA automation», «SDET» o equivalentes en inglés.

Advierta que muchas descripciones mezclan requisitos ideales. No es necesario cumplir el cien por ciento para aplicar. El análisis busca patrones entre varias vacantes, no obedecer literalmente una sola. Solicite que dos participantes compartan una competencia repetida y la evidencia que podrían crear."""
)

# 33
add_process_slide(33, "Cómo crear tu guía de estudio autónoma",
    "Cada tema debe tener propósito, práctica, evidencia y criterio de dominio.",
    [
        ("Prioriza", "Elige 2 o 3 brechas de alto impacto.", CYAN),
        ("Define", "Escribe qué podrás hacer al terminar.", BLUE),
        ("Estudia", "Usa documentación, curso o mentoría.", PURPLE),
        ("Practica", "Resuelve una tarea semejante al trabajo.", ORANGE),
        ("Demuestra", "Publica evidencia y solicita feedback.", GREEN),
    ],
    """[Tiempo sugerido: 7 minutos]
Explique que cada unidad de la guía debe responder: ¿por qué necesito aprender esto?, ¿qué tarea laboral podré realizar?, ¿qué recurso confiable usaré?, ¿qué ejercicio completaré?, ¿qué evidencia quedará? y ¿cómo sabré que lo domino?

Ejemplo: no escribir solo «aprender Postman». Escribir: «crear una colección para una API, usar variables, autenticar, validar códigos y cuerpo, ejecutar en línea de comandos y publicar documentación». Eso convierte un tema difuso en desempeño observable."""
)

# 34
add_matrix_slide(34, "Plantilla de guía de estudio basada en vacantes",
    "La guía se actualiza cuando cambian las vacantes, tu nivel o el puesto objetivo.",
    ["Brecha", "Resultado esperado", "Práctica", "Evidencia", "Fecha"],
    [
        ("Postman", "Probar API con aserciones", "Colección de 12 requests", "Repo + reporte", "Semana 1"),
        ("SQL", "Validar datos y joins", "20 consultas sobre BD demo", "Archivo SQL comentado", "Semana 2"),
        ("Playwright", "Automatizar flujo crítico", "Login + compra + error", "Suite + reporte", "Semanas 3–4"),
        ("CI", "Ejecutar pruebas en cada push", "Workflow de Actions", "Pipeline visible", "Semana 4"),
    ],
    """[Tiempo sugerido: 6 minutos]
Recorra el ejemplo. Destaque que el resultado esperado utiliza un verbo observable. La práctica se parece a una tarea real. La evidencia puede compartirse y revisarse. La fecha crea compromiso.

Recomiende limitar la guía a cuatro semanas y revisarla al finalizar. Una lista de seis meses suele quedar obsoleta y desmotiva. El aprendizaje autónomo funciona mejor con ciclos breves y revisión periódica.""",
    col_widths=[1.45, 2.55, 3.10, 3.05, 1.97])

# 35
add_cards_slide(35, "Técnicas para aprender sin depender de un curso",
    "Ser autónomo no significa estar solo: significa dirigir el proceso y buscar ayuda de forma intencional.",
    [
        ("Documentación primero", "Lee la guía oficial, reproduce un ejemplo y explica con tus palabras.", CYAN),
        ("Práctica deliberada", "Aísla una habilidad, repite con variaciones y registra errores.", BLUE),
        ("Aprender construyendo", "Cada tema debe mejorar un proyecto conectado con el puesto objetivo.", PURPLE),
        ("Recuperación activa", "Cierra el recurso y explica o resuelve sin copiar.", ORANGE),
        ("Feedback temprano", "Comparte código, casos y reportes antes de considerarlos terminados.", GREEN),
        ("Bitácora de aprendizaje", "Registra qué intentaste, qué falló, qué cambió y qué sigue.", TEAL),
    ],
    """[Tiempo sugerido: 7 minutos]
Explique cada técnica con un ejemplo. Documentación primero evita depender de tutoriales desactualizados. Práctica deliberada significa trabajar una dificultad concreta. Recuperación activa permite comprobar comprensión real.

La bitácora puede ser un archivo Markdown semanal. Debe registrar decisiones, errores y soluciones. También sirve para preparar entrevistas porque ayuda a recordar historias reales de aprendizaje y resolución de problemas.

Insista: la IA puede explicar y proponer ejercicios, pero el estudiante debe intentar, ejecutar, equivocarse y verificar."""
)

# 36
add_process_slide(36, "Sprint semanal de aprendizaje",
    "Una semana pequeña, medible y repetible vale más que un plan perfecto que nunca comienza.",
    [
        ("Lunes", "Revisar vacantes y elegir objetivo semanal.", CYAN),
        ("Martes", "Estudiar concepto y reproducir ejemplo.", BLUE),
        ("Miércoles", "Resolver ejercicio sin copiar.", PURPLE),
        ("Jueves", "Integrar al proyecto y pedir feedback.", ORANGE),
        ("Viernes", "Publicar evidencia, reflexionar y aplicar.", GREEN),
    ],
    """[Tiempo sugerido: 5 minutos]
Proponga bloques realistas de 45 a 90 minutos. La constancia importa más que sesiones maratónicas. El viernes no debe ser únicamente para seguir estudiando: también debe utilizarse para hacer visible el trabajo y enviar aplicaciones.

Recomiende conservar una lista de preguntas y errores. Si una dificultad consume demasiado tiempo, buscar documentación, comunidad, mentoría o ayuda de IA, pero después explicar la solución con palabras propias."""
)

# 37
add_bullet_slide(37, "Cómo encontrar oportunidades de manera sistemática",
    "Buscar empleo también es un proceso que se mide y mejora.",
    [
        "Diversifica: LinkedIn, Indeed, Glassdoor, Get on Board, Wellfound, Dice, RD Trabaja y Concursa.",
        "Consulta «Trabaja con nosotros» en portales institucionales de empresas públicas y privadas.",
        "Define un plan A y un plan B; aplica a perfiles relacionados con competencias transferibles.",
        "Alinea tu LinkedIn con el puesto objetivo y adapta el CV a cada tipo de vacante.",
        "Usa Glassdoor para contrastar referencias positivas y negativas; busca patrones, no opiniones aisladas.",
        "Mide aplicaciones y cultiva networking, proyectos, servicios o un emprendimiento propio.",
    ],
    """[Tiempo sugerido: 6 minutos]
Explique que una búsqueda ordenada reduce ansiedad y permite aprender. No dependan de un solo portal. LinkedIn Jobs permite descubrir vacantes y conectar con reclutadores, colegas y empresas. Indeed amplía la búsqueda general. Glassdoor combina ofertas con referencias compartidas por empleados. Para tecnología y startups pueden explorar Get on Board, Wellfound y Dice. Para oportunidades formales en República Dominicana están RD Trabaja, del Ministerio de Trabajo, y Concursa, del Ministerio de Administración Pública. También deben visitar las secciones «Carreras», «Empleos», «Vacantes» o «Trabaja con nosotros» de empresas privadas e instituciones públicas.

Una vacante encontrada en un agregador debe verificarse en el portal institucional de la empresa. Confirmen que el dominio, el correo y el proceso sean legítimos. Desconfíen de solicitudes de dinero, compras obligatorias, ofertas sin funciones claras o comunicaciones que no puedan relacionarse con una organización real.

Actualicen LinkedIn en función de las vacantes que desean alcanzar. El titular, el resumen, las competencias, la experiencia y los proyectos deben contar una historia coherente con el puesto objetivo. Si buscan QA Automation, la automatización, API, SQL, Git, CI y las evidencias relacionadas deben ser fáciles de localizar. No agreguen habilidades que no puedan explicar o demostrar.

El CV también debe adaptarse. No significa inventar una versión diferente de la persona, sino priorizar las experiencias, resultados, palabras clave y evidencias pertinentes para ese puesto. Un CV para QA funcional puede destacar análisis, casos, defectos y negocio; uno para automation debe dar mayor visibilidad a programación, frameworks, API, Git y pipelines. Mantengan una versión base y creen variantes honestas para las familias de roles a las que aplicarán.

Glassdoor puede ayudar a investigar cultura, entrevistas, salarios y experiencias positivas y negativas. No tomen una reseña aislada como verdad absoluta. Comparen patrones, fechas, ubicación, equipo y tipo de puesto. Una empresa grande puede ofrecer experiencias distintas entre departamentos. Contrasten esa información con su página oficial, LinkedIn, noticias y conversaciones con personas que conozcan la organización.

Construyan un plan A y un plan B. El plan A puede ser el puesto principal que desean, por ejemplo QA Automation Junior. El plan B puede ser una ruta cercana donde sus competencias también generen valor: QA funcional, soporte de aplicaciones, análisis de negocio, datos, implementación, desarrollo junior o coordinación tecnológica. Aplicar a perfiles relacionados no es abandonar la meta; puede ser una puerta de entrada y una fuente de experiencia.

También conviene desarrollar una alternativa propia: un pequeño servicio, proyecto, producto digital, tutoría, consultoría inicial, automatización para negocios, creación de contenido especializado o emprendimiento. No esperen pasivamente que el mercado les conceda la oportunidad. Un proyecto personal puede producir aprendizaje, contactos, portafolio y, con el tiempo, ingresos. Debe manejarse con ética y dentro de las competencias reales de cada persona.

El mercado no preguntará cuándo estarán completamente preparados. Las vacantes aparecen cuando una empresa tiene una necesidad, no cuando el candidato termina todos sus cursos. Si todavía no están listos para demostrar una competencia esencial, otra persona responderá a esa oportunidad. Por eso deben construir una preparación mínima defendible, aplicar mientras continúan aprendiendo y utilizar cada respuesta, entrevista o rechazo como información.

La empleabilidad combina habilidades técnicas y blandas. Las técnicas permiten realizar el trabajo; las blandas permiten comprender requisitos, escuchar, preguntar, colaborar, manejar desacuerdos, explicar riesgos, cumplir compromisos y aprender. Durante una entrevista, la madurez se observa en cómo reaccionan ante una pregunta desconocida, una corrección, un cambio de requisito o una limitación.

Negociar no consiste simplemente en pedir más dinero o aceptar menos para asegurar el puesto. Una negociación madura considera responsabilidades, rango salarial, beneficios, horario, modalidad, aprendizaje, estabilidad, expectativas de desempeño y posibilidades de crecimiento. La actitud profesional incluye escuchar, hacer preguntas, explicar el valor que se puede aportar y reconocer con honestidad lo que todavía debe desarrollarse.

Utilicen una hoja de seguimiento con fecha, empresa, enlace oficial, puesto, requisitos, versión del CV, contacto, estado y próximo paso. Midan aplicaciones, respuestas, entrevistas, pruebas y ofertas. Participen en eventos y comunidades porque muchas oportunidades llegan cuando otras personas conocen su trabajo, criterio y confiabilidad.

Si nadie responde, revisar CV, palabras clave y tipo de vacante. Si se obtienen entrevistas pero no ofertas, practicar comunicación y parte técnica. Si se llega a pruebas técnicas, analizar las brechas observadas. La búsqueda produce datos para mejorar.""",
    GREEN,
    ("PLAN A + PLAN B", "Posición objetivo → perfil alternativo → proyecto o emprendimiento", "Prepárate, aplica, negocia y ajusta."))

# 38
add_split_slide(38, "CV adaptado al puesto, no biografía completa",
    "En pocos segundos debe quedar claro qué rol buscas, qué sabes hacer y qué evidencia posees.",
    "Incluye", [
        "Titular alineado con el puesto",
        "Resumen breve con valor y nivel",
        "Herramientas que realmente utilizas",
        "Proyectos y logros con evidencia",
        "Enlaces a LinkedIn y GitHub",
    ],
    "Evita", [
        "Lista interminable de cursos",
        "Herramientas sin contexto",
        "Responsabilidades copiadas",
        "Diseño difícil de leer",
        "Datos irrelevantes o sensibles",
    ],
    """[Tiempo sugerido: 6 minutos]
Explique la diferencia entre responsabilidad y logro. «Ejecuté pruebas» es débil. «Diseñé y ejecuté 45 casos para tres flujos críticos, documenté ocho defectos y automaticé la regresión principal» aporta contexto, aunque sea un proyecto académico correctamente identificado.

No recomendar falsificar experiencia. Los proyectos personales deben presentarse como proyectos. El CV se adapta usando el lenguaje de la vacante sin copiar afirmaciones que no se puedan defender."""
)

# 39
add_cards_slide(39, "LinkedIn y GitHub: tu evidencia pública",
    "El perfil profesional debe facilitar que una persona entienda tu dirección y encuentre tu mejor trabajo.",
    [
        ("Titular de LinkedIn", "Rol objetivo + capacidades principales + tipo de valor que aportas.", CYAN),
        ("Acerca de", "Historia breve: hacia dónde vas, qué practicas y qué buscas.", BLUE),
        ("Actividad", "Publica aprendizajes, proyectos, análisis y participación en comunidad.", TEAL),
        ("Perfil de GitHub", "README claro, biografía profesional y repositorios fijados.", PURPLE),
        ("Repositorios", "Instrucciones, decisiones, evidencias, reportes y pipeline.", ORANGE),
        ("Consistencia", "Mismo nombre profesional, enlaces funcionales y contacto actualizado.", GREEN),
    ],
    """[Tiempo sugerido: 6 minutos]
Explique que publicar no significa convertirse en creador de contenido diario. Puede compartirse una lección concreta, un error resuelto o una mejora del proyecto.

GitHub permite crear un README de perfil y fijar repositorios. Recomiende seleccionar entre tres y seis trabajos relevantes, no mostrar veinte tutoriales incompletos. Cada repositorio debe indicar propósito, herramientas, instalación, ejecución, resultados y próximos pasos.""",
    "Fuente práctica: documentación oficial de GitHub sobre perfil, README y repositorios fijados.")

# 40
add_cards_slide(40, "Cómo prepararte para la entrevista",
    "Las empresas evalúan conocimiento, razonamiento, comunicación y capacidad de aprendizaje.",
    [
        ("Fundamentos", "Explica tipos de prueba, riesgo, severidad, prioridad y ciclo de defectos.", CYAN),
        ("Casos reales", "Prepara historias sobre un error, una decisión, un conflicto y un aprendizaje.", BLUE),
        ("Práctica técnica", "Diseña casos, consulta SQL, prueba API o explica automatización.", PURPLE),
        ("Pensamiento visible", "Explica preguntas, supuestos, prioridades y criterios.", ORANGE),
        ("Preguntas a la empresa", "Consulta proceso, equipo, calidad, herramientas, mentoría y expectativas.", GREEN),
        ("Seguimiento", "Agradece, registra feedback y actualiza tu guía de estudio.", TEAL),
    ],
    """[Tiempo sugerido: 6 minutos]
Recomiende practicar en voz alta. Una persona puede conocer la respuesta y no lograr comunicarla bajo presión. Para historias de experiencia, puede utilizarse la estructura situación, tarea, acción y resultado.

En preguntas técnicas, no intentar adivinar inmediatamente. Aclarar contexto, identificar riesgos y explicar el proceso. Si no se conoce una herramienta, decir qué concepto relacionado se domina y cómo se aprendería. La honestidad acompañada de razonamiento es mejor que inventar."""
)

# 41
add_process_slide(41, "Sistema para mantenerte actualizado",
    "No persigas cada novedad. Construye un radar conectado con tu rol y revísalo periódicamente.",
    [
        ("Semanal", "Leer releases, una vacante y una fuente técnica.", CYAN),
        ("Mensual", "Actualizar matriz, proyecto, CV y perfil.", BLUE),
        ("Trimestral", "Revisar 20 vacantes y ajustar la ruta.", PURPLE),
        ("Comunidad", "Eventos, meetups, QA Dominicana y colaboración.", ORANGE),
        ("Evidencia", "Publicar una mejora demostrable y reflexionar.", GREEN),
    ],
    """[Tiempo sugerido: 6 minutos]
Explique que mantenerse actualizado no significa leer noticias todo el día. Se necesita un sistema con frecuencia y filtros. Las fuentes principales deben incluir documentación oficial, notas de versiones, comunidades profesionales y vacantes.

Cada tres meses, repetir el análisis de mercado. Si cambian las herramientas más frecuentes, evaluar si la ruta debe ajustarse. No abandonar fundamentos por una moda. La actualización sostenible combina señales del mercado con profundidad técnica.""",
    "Fuentes sugeridas: documentación oficial de herramientas, ISTQB, GitHub, OWASP, comunidades y vacantes objetivo.")

# 42
add_activity_slide(42, "Tu plan de 30 días",
    "Define: un puesto objetivo, diez vacantes para analizar, dos brechas prioritarias, un proyecto de evidencia, cinco aplicaciones semanales y una revisión cada viernes.",
    "Plan personal de una página", "8 min",
    """[Tiempo sugerido: 8 minutos]
Dé cinco minutos para escribir y tres para compartir compromisos. Recalque que el plan debe caber en una página y tener acciones verificables.

Ejemplo: analizar diez vacantes de QA Automation Junior, estudiar Postman y GitHub Actions, publicar una colección y un pipeline, adaptar el CV y enviar cinco aplicaciones por semana. Al final de cada semana revisar qué se completó y qué señal entregó el mercado."""
)

# 43
add_cards_slide(43, "Recursos recomendados para continuar",
    "Prioriza fuentes oficiales y práctica verificable.",
    [
        ("Fundamentos", "ISTQB CTFL 4.0 y glosario oficial.", CYAN),
        ("IA para testing", "ISTQB CT-GenAI y sus riesgos, prompts y validación.", BLUE),
        ("Testing de IA", "ISTQB CT-AI 2.0: datos, modelos, GenAI y red teaming.", TEAL),
        ("Automatización", "Documentación de Playwright, Cypress o Selenium según vacantes.", PURPLE),
        ("CI y portafolio", "GitHub Docs: perfil, README, Actions y seguridad.", ORANGE),
        ("Seguridad de IA", "OWASP Top 10 for LLM Applications 2025.", GREEN),
    ],
    """[Tiempo sugerido: 3 minutos]
Explique que estos recursos no deben estudiarse todos al mismo tiempo. Cada participante seleccionará los que respondan a su puesto objetivo y brechas prioritarias.

Recomiende guardar enlaces oficiales y comprobar fechas o versiones. Un video puede servir para comenzar, pero la documentación oficial ayuda a verificar comportamientos y cambios."""
)

# 44 · Fundamentos de programación para QA Automation
add_cards_slide(44, "POO y principios para automatización mantenible",
    "El objetivo no es escribir más código, sino crear pruebas claras, reutilizables y fáciles de cambiar.",
    [
        ("POO", "Encapsular datos y comportamientos en objetos como LoginPage, ApiClient o TestUser.", CYAN),
        ("DRY", "No repetir pasos, localizadores o datos; extraer únicamente lo que realmente se reutiliza.", BLUE),
        ("KISS", "Preferir una solución sencilla y legible frente a una arquitectura innecesariamente compleja.", GREEN),
        ("SOLID", "Separar responsabilidades, extender sin romper y depender de abstracciones controlables.", PURPLE),
        ("Código limpio", "Nombres expresivos, funciones pequeñas, errores claros y estructura consistente.", ORANGE),
        ("Pruebas independientes", "Cada prueba prepara sus datos, valida un propósito y deja el ambiente controlado.", TEAL),
    ],
    """[Tiempo sugerido: 7 minutos]
Texto para leer:
Cuando una persona comienza automatización puede concentrarse únicamente en lograr que el script funcione. Sin embargo, una suite profesional también debe mantenerse. Aquí aparecen la programación orientada a objetos, conocida como POO, y principios como DRY, KISS y SOLID.

POO permite representar responsabilidades mediante objetos. Una clase LoginPage puede encapsular los localizadores y acciones de la pantalla de acceso. DRY significa «Don't Repeat Yourself»: no debemos copiar el mismo inicio de sesión en veinte archivos. KISS significa «Keep It Simple»: una solución sencilla y entendible suele ser mejor que una estructura sofisticada que nadie puede mantener.

SOLID reúne cinco principios. El más útil para comenzar es responsabilidad única: una clase no debe encargarse al mismo tiempo de navegar, crear datos, consultar la base y producir reportes. También es importante depender de componentes sustituibles para poder simular servicios y cambiar herramientas con menos impacto.

Estos principios no se aplican mecánicamente. Extraer todo también puede crear complejidad. La pregunta correcta es: ¿esta decisión hace que la prueba sea más clara, estable y fácil de modificar?""",
    "Conceptos aplicados a frameworks de automatización y código de pruebas.")

# 45 · Ejemplos de principios y patrones
add_matrix_slide(45, "Ejemplos: DRY, KISS, SOLID y patrones de diseño",
    "Los patrones son soluciones conocidas; deben utilizarse cuando resuelven un problema real.",
    ["Concepto", "Ejemplo aplicado a QA", "Beneficio"],
    [
        ("POO", "LoginPage agrupa campos y acciones de acceso", "Encapsulación y lectura"),
        ("DRY", "Una fixture reutiliza preparación y limpieza", "Menos duplicación"),
        ("KISS", "Una prueba valida un comportamiento principal", "Diagnóstico más simple"),
        ("SOLID", "ApiClient, TestData y Reporter separados", "Cambio con menor impacto"),
        ("Page Object", "Cada pantalla expone acciones del usuario", "Localizadores centralizados"),
        ("Builder", "TestUserBuilder crea usuarios con variaciones", "Datos legibles y flexibles"),
        ("Factory", "BrowserFactory selecciona navegador o ambiente", "Creación controlada"),
        ("Strategy", "PaymentStrategy cambia la forma de pago", "Comportamientos intercambiables"),
    ],
    """[Tiempo sugerido: 7 minutos]
Texto para leer:
Veamos ejemplos concretos. Page Object Model, o modelo de objetos de página, centraliza localizadores y acciones de una pantalla. La prueba puede decir loginPage.iniciarSesion en lugar de repetir detalles técnicos. El patrón Builder permite construir datos de prueba de forma expresiva: un usuario válido, uno bloqueado o uno sin dirección. Factory centraliza la creación de navegadores, clientes o ambientes. Strategy permite intercambiar un comportamiento, como distintos métodos de autenticación o pago.

DRY evita repetición, pero no significa crear una función genérica para todo. Si dos pasos se parecen pero representan reglas distintas, unirlos puede ocultar intención. KISS recuerda que la prueba debe poder leerse y diagnosticarse. SOLID ayuda a separar cambios: modificar el reporte no debería romper la navegación.

Otro patrón útil es Adapter, que envuelve una herramienta externa para evitar que todo el framework dependa directamente de ella. También puede utilizarse Facade para ofrecer una interfaz sencilla a un flujo complejo. La recomendación es comenzar con Page Object, fixtures y Builder, y añadir otros patrones solamente cuando el proyecto los necesite.""",
    "Ejemplos orientativos; la arquitectura debe responder al tamaño y riesgo del proyecto.",
    col_widths=[1.80, 6.50, 3.82])

# 46 · Evolución de posiciones tecnológicas
add_split_slide(46, "Vacantes tecnológicas: 2015-2024 y desde 2024 con IA",
    "Los puestos anteriores continúan, pero ahora incorporan IA, verificación y nuevas responsabilidades.",
    "Posiciones frecuentes 2015–2024", [
        "Software developer, frontend, backend y full-stack",
        "QA analyst, automation engineer y SDET",
        "Business analyst y systems analyst",
        "Data analyst, BI developer y data engineer",
        "UX/UI designer y UX researcher",
        "DevOps, cloud y cybersecurity analyst",
        "Project manager, Scrum, soporte y technical writer",
    ],
    "Posiciones y capacidades 2024–actualidad", [
        "AI-assisted developer y GenAI engineer",
        "AI-augmented QA, AI tester y LLM evaluator",
        "AI business analyst y AI product manager",
        "Data/ML/AI engineer y model validator",
        "AI UX y conversation designer",
        "MLOps, LLMOps y AI security",
        "Responsible AI y AI agent builder",
    ],
    """[Tiempo sugerido: 7 minutos]
Texto para leer:
En esta comparación no estamos diciendo que las posiciones existentes entre 2015 y 2024 desaparecieron en 2024. Software developer, QA, análisis de negocio, datos, UX, DevOps, ciberseguridad, gestión, soporte y documentación siguen siendo necesarias. El cambio es que muchas vacantes ahora esperan que estos profesionales sepan utilizar herramientas de inteligencia artificial, evaluar sus resultados y proteger la información.

Entre 2015 y 2024, de un desarrollador se esperaba que diseñara, programara, depurara y documentara soluciones utilizando lenguajes, frameworks, control de versiones y patrones. De QA se esperaba analizar requisitos, diseñar y ejecutar pruebas, reportar defectos y automatizar regresiones. Un business analyst levantaba procesos, entrevistaba usuarios y redactaba requisitos. Un data analyst preparaba datos, escribía consultas y construía reportes. UX investigaba usuarios y diseñaba prototipos. DevOps trabajaba ambientes, pipelines y nube. Seguridad analizaba riesgos e incidentes. Gestión preparaba planes y reportes; soporte clasificaba casos y documentación técnica producía guías.

Desde 2024 estas responsabilidades no desaparecen, pero se añaden nuevas expectativas. Se espera alfabetización en IA, capacidad para entregar contexto y restricciones, verificar respuestas, reconocer alucinaciones y sesgos, proteger datos, revisar código generado e integrar asistentes dentro de un flujo de trabajo. También aparecen posiciones como GenAI application engineer, AI product manager, AI test engineer, LLM evaluator, AI red teamer, MLOps o LLMOps, Responsible AI specialist, AI automation specialist y AI agent builder. Los títulos todavía cambian entre empresas; por eso debemos leer las funciones.

Ahora compartiré tiempos ilustrativos para comprender el cambio. No son garantías ni métricas universales. Varían según experiencia, complejidad, calidad de los datos, herramienta y revisión requerida.

Un resumen inicial de reunión o requisitos podía tomar entre dos y cuatro horas; con apoyo de IA, el primer borrador puede reducirse a treinta o noventa minutos, pero una persona debe validar decisiones y acuerdos. Diseñar un conjunto inicial de casos de prueba podía tomar dos o tres horas; la IA puede proponer un borrador en treinta o sesenta minutos, mientras QA dedica tiempo a riesgo, cobertura y resultados esperados. Crear código repetitivo o pruebas unitarias iniciales podía tomar entre dos y cuatro horas; un asistente puede producir una base en treinta o noventa minutos, seguida de ejecución, revisión y refactorización.

Una exploración inicial de datos, consultas y narrativa podía requerir entre tres y seis horas; con IA puede tomar de una a tres horas, siempre verificando cálculos y privacidad. Un primer borrador de documentación o reporte podía tomar entre dos y cinco horas; con IA puede reducirse a treinta o noventa minutos. Clasificar un caso de soporte y redactar una respuesta podía tomar de diez a veinte minutos; un borrador asistido puede tomar de dos a cinco minutos, pero el agente humano confirma exactitud y tono. Sintetizar entrevistas de UX podía tomar entre tres y cinco horas; la IA puede apoyar una primera clasificación en una o dos horas, sin reemplazar la interpretación del investigador.

La idea no es trabajar más rápido para aceptar resultados sin revisar. El tiempo recuperado debe utilizarse para investigar excepciones, comprobar fuentes, conversar con usuarios, mejorar seguridad y tomar decisiones. La IA reduce parte del tiempo de producción inicial; la responsabilidad profesional permanece.""",
    "Tiempos orientativos para discusión, no mediciones universales. Fuentes de contexto: WEF 2025, LinkedIn 2025 e ISTQB 2026.")

# 44
slide = prs.slides.add_slide(blank)
base(slide, 44, "Fuentes principales", "Referencias utilizadas para actualizar el contenido del taller.")
sources = [
    "World Economic Forum. Future of Jobs Report 2025. https://www.weforum.org/publications/the-future-of-jobs-report-2025/",
    "LinkedIn. Skills on the Rise in 2025. https://www.linkedin.com/business/talent/blog/learning-and-development/skills-on-the-rise",
    "ISTQB. Certified Tester AI Testing, syllabus 2.0. https://istqb.org/",
    "ISTQB. Certified Tester – Testing with Generative AI. https://istqb.org/certifications/gen-ai/",
    "Playwright. Documentación oficial. https://playwright.dev/docs/intro",
    "GitHub Docs. Profile, README y Actions. https://docs.github.com/",
    "OWASP. Top 10 for LLM Applications 2025. https://owasp.org/www-project-top-10-for-large-language-model-applications/",
]
bullets(slide, sources, 0.88, 1.52, 11.72, 4.92, 13.6, INK, 10)
add_notes(slide, """[Tiempo sugerido: 2 minutos]
Indique que estas fuentes se incluyen para que los participantes puedan verificar y profundizar. No es necesario leer cada documento completo de inmediato. La recomendación es seleccionar la fuente relacionada con la próxima competencia de su guía de estudio.""")

# 45
slide = prs.slides.add_slide(blank)
rect(slide, 0, 0, 13.333, 7.5, NAVY)
rect(slide, 0, 0, 13.333, 0.12, ORANGE)
text(slide, "BENEFICIO ESPECIAL", 0.68, 0.50, 3.55, 0.30, 12, ORANGE, True)
text(slide, "Sorteo de 6 boletas · Testing 4 All 2026", 0.68, 0.92, 9.75, 0.62, 27, WHITE, True)
circle(slide, 11.10, 0.48, 1.35, ORANGE)
text(slide, "6", 11.10, 0.72, 1.35, 0.46, 28, NAVY, True, PP_ALIGN.CENTER)
text(slide, "GANADORES", 10.80, 1.88, 1.95, 0.25, 10, ORANGE, True, PP_ALIGN.CENTER)

rect(slide, 0.68, 1.82, 12.00, 0.72, DARK_CARD, True)
text(slide, "11 de julio de 2026 · 9:00 a. m.–5:00 p. m. · QA Dominicana · Universidad APEC",
     0.98, 2.03, 11.40, 0.30, 14.5, PALE, True, PP_ALIGN.CENTER)

steps_raffle = [
    ("1", "Completar el form", "Registrar correctamente los datos solicitados durante la charla."),
    ("2", "Disponibilidad", "No debes residir en Santo Domingo, pero sí poder llegar al evento."),
    ("3", "Compromiso", "Confirmar asistencia para aprender, conectar y retirar la boleta."),
    ("4", "Sorteo y DM", "Se elegirán seis nombres al azar y se confirmarán por mensaje privado."),
]
for i, (step_no, step_title, step_body) in enumerate(steps_raffle):
    x = 0.68 + i * 3.05
    rect(slide, x, 2.90, 2.75, 2.18, DARK_CARD, True, RGBColor(49, 77, 108))
    circle(slide, x + 0.22, 3.16, 0.56, ORANGE if i < 3 else GREEN)
    text(slide, step_no, x + 0.22, 3.26, 0.56, 0.28, 13, NAVY, True, PP_ALIGN.CENTER)
    text(slide, step_title, x + 0.88, 3.08, 1.68, 0.58, 12.8, WHITE, True,
         valign=MSO_ANCHOR.MIDDLE)
    text(slide, step_body, x + 0.28, 3.82, 2.19, 0.92, 11.3, PALE, align=PP_ALIGN.CENTER)

rect(slide, 1.45, 5.38, 10.45, 0.72, ORANGE, True)
text(slide, "6 boletas · Formulario · Disponibilidad para asistir · Confirmación y retiro personal",
     1.72, 5.59, 9.92, 0.30, 14.3, NAVY, True, PP_ALIGN.CENTER)
raffle_cta = rect(slide, 1.45, 6.18, 10.45, 0.44, CYAN, True, CYAN)
raffle_cta.click_action.hyperlink.address = RAFFLE_FORM_URL
raffle_cta_text = text(slide, "ABRIR FORMULARIO DEL SORTEO · 6 BOLETAS  ↗",
                       1.72, 6.28, 9.92, 0.24, 11.8, NAVY, True, PP_ALIGN.CENTER)
raffle_cta_text.click_action.hyperlink.address = RAFFLE_FORM_URL
instagram = text(slide, "VER AGENDA OFICIAL EN INSTAGRAM · @qa_dominicana  ↗",
                 2.30, 6.72, 8.72, 0.22, 10.2, CYAN, True, PP_ALIGN.CENTER)
instagram.click_action.hyperlink.address = EVENT_INSTAGRAM_URL
footer(slide, 45, "Taller QA · Sorteo especial")
add_notes(slide, """[Tiempo sugerido: 5 minutos]
Texto para leer:
Como beneficio especial de esta charla–taller, donaré seis boletas para el evento Testing 4 All 2026. Por tanto, seleccionaremos seis personas ganadoras mediante una dinámica sencilla y transparente.

Estas son las bases del sorteo. Primero, durante la charla compartiré el enlace de un formulario de Google Forms. Las personas interesadas deberán completarlo correctamente con la información solicitada. La lista utilizada para el sorteo se obtendrá únicamente de las respuestas recibidas en ese formulario.

Segundo, no es obligatorio residir en Santo Domingo. Pueden participar personas de otras localidades, siempre que tengan la disponibilidad y la facilidad para trasladarse a la Universidad APEC el día del evento. Lo importante es asegurar que realmente podrán asistir, porque la boleta será entregada personalmente en la entrada y no se enviará de manera digital.

Tercero, al completar el formulario la persona declara que tiene la intención real de asistir y se compromete a retirar personalmente la boleta en la entrada. Les pido participar solamente si pueden cumplir este compromiso, para evitar que una entrada destinada a la comunidad quede sin utilizar.

La finalidad no es únicamente entrar al evento. Queremos que las personas ganadoras aprovechen la jornada completa para actualizar sus conocimientos, conocer tendencias de QA, automatización e inteligencia artificial, conversar con representantes de empresas, ampliar su red profesional y relacionarse con colegas del área. El networking puede generar mentorías, colaboraciones, referencias profesionales y futuras oportunidades de empleo.

Cuando cierre el formulario, utilizaré la lista de respuestas válidas para realizar una selección aleatoria de seis nombres. Después contactaré a cada persona seleccionada mediante un mensaje directo, o DM, para confirmar su identidad, disponibilidad para trasladarse, participación y compromiso de asistencia. La selección se considerará confirmada cuando la persona responda el mensaje privado.

Testing 4 All 2026 se realizará el 11 de julio de 2026, de nueve de la mañana a cinco de la tarde, en la Universidad APEC, y es organizado por QA Dominicana. En la diapositiva encontrarán un enlace directo a la publicación de Instagram donde pueden consultar la agenda oficial del evento.

Yo estaré en la entrada el día del evento y entregaré personalmente las seis boletas a las personas ganadoras que hayan sido confirmadas por DM. Recuerden: formulario obligatorio, disponibilidad real para asistir, compromiso de aprovechar la actividad y retiro personal de la boleta en la entrada.

El botón naranja de esta diapositiva abre directamente el formulario oficial del sorteo. Debajo también encontrarán el acceso a la agenda oficial publicada en Instagram.""")

# 49 · Preguntas del formulario para el sorteo
add_bullet_slide(49, "Preguntas para el formulario del sorteo",
    "Solicita solamente los datos necesarios para validar la participación y contactar a las personas ganadoras.",
    [
        "Nombre completo y correo electrónico",
        "Usuario de Instagram o LinkedIn para confirmar por DM",
        "Ciudad o provincia desde donde se trasladará",
        "Área de estudio, trabajo o interés profesional",
        "Disponibilidad para asistir el 11 de julio, de 9:00 a. m. a 5:00 p. m.",
        "Compromiso de retirar personalmente la boleta en la entrada",
        "Aceptación de las bases y autorización de contacto para el sorteo",
    ],
    """[Tiempo sugerido: 4 minutos]
Texto para leer:
Este será el contenido recomendado para el formulario de Google. Todas las preguntas de identificación, disponibilidad y aceptación deben configurarse como obligatorias. Evitaremos solicitar documentos de identidad, dirección exacta u otros datos que no sean necesarios para realizar el sorteo.

Pregunta uno, respuesta corta y obligatoria: ¿Cuál es tu nombre completo?

Pregunta dos, respuesta corta y obligatoria: ¿Cuál es tu correo electrónico?

Pregunta tres, respuesta corta y obligatoria: Escribe tu usuario de Instagram o la dirección de tu perfil de LinkedIn donde podamos contactarte por mensaje privado. Debe indicarse que este dato se utilizará únicamente para confirmar el resultado del sorteo.

Pregunta cuatro, respuesta corta y obligatoria: ¿Desde cuál ciudad o provincia te trasladarías al evento? Esta pregunta no excluye a quienes viven fuera de Santo Domingo; permite confirmar que la persona ha considerado su traslado.

Pregunta cinco, opción múltiple y obligatoria: ¿Cuál opción describe mejor tu situación actual? Las opciones pueden ser: estudiante; QA manual; QA automation o SDET; desarrollo de software; datos o inteligencia artificial; ciberseguridad, DevOps o nube; análisis de negocio o producto; otra área tecnológica; y persona en transición hacia tecnología.

Pregunta seis, casillas de verificación opcionales: ¿Qué deseas aprovechar principalmente en Testing 4 All 2026? Las opciones pueden ser: aprender sobre QA y testing; automatización; inteligencia artificial; conocer empresas; buscar orientación profesional; conectar con colegas; explorar oportunidades laborales; y otra.

Pregunta siete, opción sí o no y obligatoria: ¿Confirmas que puedes trasladarte a la Universidad APEC y asistir presencialmente el 11 de julio de 2026, de 9:00 de la mañana a 5:00 de la tarde? Para participar debe seleccionar sí.

Pregunta ocho, opción sí o no y obligatoria: Si resultas seleccionado o seleccionada, ¿te comprometes a responder el DM de confirmación y retirar personalmente tu boleta en la entrada del evento? Para participar debe seleccionar sí.

Pregunta nueve, casilla obligatoria de aceptación: He leído y acepto las bases del sorteo; confirmo que mis datos son correctos, que tengo disponibilidad real para asistir y que autorizo el contacto por correo o mensaje privado exclusivamente para gestionar este sorteo.

Pregunta diez, respuesta breve opcional: ¿Qué esperas aprender o lograr durante el evento? Esta respuesta no debe utilizarse para favorecer a una persona, porque la selección será aleatoria; servirá para conocer las expectativas de los participantes.

En la descripción del formulario conviene informar que se seleccionarán seis personas al azar entre las respuestas válidas, que la boleta no se enviará digitalmente y que será entregada personalmente en la entrada. También debe indicarse la fecha de cierre del formulario cuando se haya decidido.""",
    ORANGE,
    ("VALIDACIÓN", "Solo participan respuestas completas que confirmen disponibilidad real y aceptación de las bases.", "Selección aleatoria de 6 personas."),
    "No solicites cédula, dirección exacta ni información que no sea necesaria para el sorteo.")
raffle_form_slide = prs.slides[-1]
raffle_form_button = rect(raffle_form_slide, 1.05, 6.20, 11.22, 0.48, ORANGE, True, ORANGE)
raffle_form_button.click_action.hyperlink.address = RAFFLE_FORM_URL
raffle_form_link = text(
    raffle_form_slide,
    "ABRIR FORMULARIO DEL SORTEO · 6 BOLETAS  ↗",
    1.28, 6.31, 10.76, 0.24, 11.8, NAVY, True, PP_ALIGN.CENTER,
)
raffle_form_link.click_action.hyperlink.address = RAFFLE_FORM_URL
raffle_form_notes = raffle_form_slide.notes_slide.notes_text_frame
if raffle_form_notes is not None:
    raffle_form_notes.text = raffle_form_notes.text.rstrip() + (
        "\n\nEl botón visible en la diapositiva abre directamente el formulario oficial del sorteo: "
        + RAFFLE_FORM_URL
    )

# 50
slide = prs.slides.add_slide(blank)
rect(slide, 0, 0, 13.333, 7.5, NAVY)
rect(slide, 0, 0, 13.333, 0.12, CYAN)
text(slide, "Cierre", 0.78, 0.72, 2.0, 0.35, 14, CYAN, True)
text(slide, "No estudies todo.", 0.78, 1.42, 8.60, 0.72, 34, WHITE, True)
text(slide, "Estudia lo que exige el trabajo que quieres hacer.", 0.78, 2.25, 10.75, 0.82, 29, CYAN, True)
text(slide, "Practica · demuestra · solicita feedback · aplica · actualiza tu ruta",
     0.82, 3.62, 10.90, 0.52, 18, PALE, True)
rect(slide, 0.82, 4.48, 10.90, 0.07, CYAN)
text(slide, "Preguntas y conversación", 0.82, 5.10, 6.15, 0.45, 22, WHITE, True)
text(slide, "linkedin.com/in/marlenis-judith-c-561117a3/", 0.82, 6.10, 6.65, 0.32, 12, CYAN)
text(slide, "Gracias", 9.35, 5.55, 2.55, 0.62, 28, WHITE, True, PP_ALIGN.CENTER)
text(slide, "Mtra. Marlenis Judith Concepción Cuevas", 0.82, 6.42, 7.15, 0.24, 10.2, WHITE, True)
text(slide, "Senior Software Development Engineer in Test (SDET) · Docente e investigadora",
     0.82, 6.68, 8.10, 0.20, 8.4, CYAN, True)
text(slide, "Escuela de Informática · Facultad de Ciencias · UASD",
     0.82, 6.89, 7.35, 0.18, 8.0, PALE)
if UASD_LOGO.exists():
    slide.shapes.add_picture(str(UASD_LOGO), Inches(11.45), Inches(0.46), width=Inches(1.15))
footer(slide, 46, "La calidad se construye con criterio, evidencia y aprendizaje continuo")
add_notes(slide, """[Tiempo sugerido: 5 minutos]
Cierre retomando la idea principal: el mercado cambia, pero podemos observarlo y responder con intención. No es necesario perseguir todas las herramientas. Se debe elegir un rol, analizar vacantes, priorizar competencias, estudiar mediante práctica, producir evidencia y comenzar a aplicar.

Invite a formular preguntas. Si queda poco tiempo, pida que escriban en el chat su compromiso de los próximos siete días. Agradezca la participación y recuerde revisar su plan de 30 días.""")


# 51 · Tarea para realizar en casa
add_cards_slide(51, "Tarea para casa: Automation Practice",
    "Convierte el repositorio en evidencia para tu portafolio; no te limites a copiar o ejecutar.",
    [
        ("1. Preparar", "Lee el README completo, identifica requisitos y realiza fork o clonación del repositorio.", CYAN),
        ("2. Comprender", "Instala lo indicado, ejecuta el proyecto y explica la estructura y las pruebas existentes.", BLUE),
        ("3. Ampliar", "Agrega o mejora un escenario de prueba y verifica resultados positivos y negativos.", PURPLE),
        ("4. Refactorizar", "Aplica POO, DRY, KISS, SOLID o un patrón cuando aporte claridad y mantenimiento.", ORANGE),
    ],
    """[Tiempo sugerido: 7 minutos]
Texto para leer:
Como práctica para realizar en casa, les propongo trabajar con el repositorio vamcodeAutomationPractice. El enlace al README aparece en esta diapositiva. Primero lean el documento completo para identificar las herramientas, requisitos de instalación, estructura y forma de ejecución. No comiencen copiando comandos sin comprender qué hacen.

Después realicen un fork o clonen el repositorio según su nivel y preferencia. Instalen las dependencias indicadas y ejecuten las pruebas existentes. Si aparece un error, documenten el mensaje, la posible causa y la solución utilizada. Esa bitácora también constituye evidencia de aprendizaje.

El siguiente paso es seleccionar un flujo o comportamiento. Pueden agregar un caso positivo, un caso negativo, un valor límite o mejorar una validación existente. Luego revisen el código e identifiquen una mejora concreta: eliminar duplicación con DRY, simplificar mediante KISS, encapsular una pantalla con POO, separar responsabilidades aplicando SOLID o utilizar Page Object, Builder, Factory o Strategy cuando exista una necesidad real.

Como entregable personal, conserven el enlace a su repositorio o fork, el reporte de ejecución, una captura de evidencia, los commits realizados y una reflexión breve: qué comprendieron, qué dificultad encontraron, qué principio aplicaron y qué mejorarían después. Este ejercicio puede convertirse en una pieza del portafolio profesional.""",
    "Ejercicio sugerido; cada participante puede adaptarlo a su nivel.")
homework_slide = prs.slides[-1]
homework_button = rect(homework_slide, 1.05, 6.38, 11.22, 0.50, NAVY, True, CYAN)
homework_button.click_action.hyperlink.address = (
    "https://github.com/marlenis-concepcion/vamcodeAutomationPractice/blob/main/README.md"
)
homework_link = text(
    homework_slide,
    "ABRIR REPOSITORIO Y GUÍA DE PRÁCTICA EN GITHUB  ↗",
    1.25, 6.50, 10.82, 0.24, 11.5, WHITE, True, PP_ALIGN.CENTER,
)
homework_link.click_action.hyperlink.address = (
    "https://github.com/marlenis-concepcion/vamcodeAutomationPractice/blob/main/README.md"
)

# 52 · Habilidades blandas para el mundo tecnológico
add_cards_slide(52, "Las habilidades blandas convierten conocimiento en impacto",
    "La capacidad técnica abre puertas; la manera de colaborar, decidir y responder sostiene la carrera.",
    [
        ("Comunicación clara", "Explicar riesgos, decisiones, avances y bloqueos según la audiencia.", CYAN),
        ("Escucha y preguntas", "Comprender el problema antes de defender una solución o ejecutar una tarea.", BLUE),
        ("Adaptabilidad", "Aprender, ajustar prioridades y responder al cambio sin trabajar desde el caos.", PURPLE),
        ("Ownership", "Asumir compromisos, comunicar límites y responder por resultados sin buscar culpables.", ORANGE),
        ("Feedback y conflicto", "Dar y recibir observaciones sobre el trabajo con respeto, evidencia y apertura.", GREEN),
        ("Negociación e influencia", "Acordar alcance, tiempo, calidad y expectativas sin imponer ni ceder a ciegas.", TEAL),
    ],
    """[Tiempo sugerido: 4 minutos]
Texto para leer:
Ser muy inteligente o dominar muchas herramientas no garantiza que una persona pueda liderar, colaborar o generar confianza. En tecnología trabajamos con requisitos incompletos, prioridades que compiten, usuarios con necesidades distintas y equipos que dependen unos de otros. Las habilidades blandas convierten el conocimiento técnico en resultados que otras personas pueden comprender y utilizar.

Comunicar bien no significa hablar mucho. Significa explicar el problema, la evidencia, el impacto, las opciones y la recomendación de acuerdo con la audiencia. Un defecto crítico se explica de una manera al desarrollador, de otra al Product Owner y de otra a una persona directiva.

Escuchar y preguntar evita resolver el problema equivocado. La adaptabilidad permite responder al cambio, pero no significa aceptar desorganización permanente. Ownership tampoco significa cargar con todo; implica asumir compromisos, advertir riesgos a tiempo, pedir ayuda y cerrar el ciclo.

El feedback profesional se dirige al trabajo y a sus consecuencias, no al valor personal de quien lo realizó. La negociación permite acordar alcance, fechas, calidad, responsabilidades y prioridades. Liderar no comienza cuando aparece un cargo de manager: comienza cuando una persona aporta claridad, facilita decisiones, reconoce contribuciones y ayuda al equipo a aprender."""
)

# 53 · Comentarios positivos de cultura y liderazgo
add_split_slide(53, "Comentarios positivos y matizados en Glassdoor",
    "Fragmentos breves de reseñas anónimas; las empresas se omiten y cada experiencia requiere contexto.",
    "Comentarios publicados", [
        "“People genuinely care — about the product, about the customer, and about each other.” (P1)",
        "“Urgent requests happen often.” (P2)",
        "“There is more to learn around every corner.” (P3)",
        "“Team-first.” (P4)",
    ],
    "Lectura profesional", [
        "P1: compromiso con producto, clientes y compañeros",
        "P2: incluso una reseña positiva puede revelar urgencias y silos",
        "P3: aprendizaje continuo y posibilidad de crecimiento",
        "P4: colaboración y disposición para ayudar",
    ],
    """[Tiempo sugerido: 4 minutos]
Texto para leer:
En las reseñas positivas compartidas para este taller aparecen varios patrones: sentido de responsabilidad, apoyo entre compañeros, apertura al feedback, interés genuino por el cliente, contratación cuidadosa, aprendizaje continuo y liderazgo cercano. Algunas personas describen ambientes de crecimiento rápido donde se espera esfuerzo, pero también existe apoyo para aprender y evitar el agotamiento.

No debemos confundir una cultura saludable con una cultura sin presión. Una empresa puede moverse rápido y mantener claridad, respeto y prioridades comprensibles. También puede exigir un proceso de entrevista riguroso si explica lo que evalúa y trata a los candidatos con profesionalidad.

Durante una entrevista, la persona candidata también investiga. Pida ejemplos concretos. Si pregunta cómo se manejan los errores y solamente recibe una respuesta genérica, solicite una situación reciente. Pregunte cómo se comunican cambios de prioridad, cómo se reconoce el trabajo, cómo se comparte conocimiento y qué significa crecer en ese puesto.

Una respuesta madura suele incluir procesos, ejemplos y límites. Las respuestas perfectas pero vagas no permiten evaluar la realidad. Observe además cómo lo tratan durante el proceso: puntualidad, claridad, respeto por su tiempo, coherencia entre entrevistadores y disposición para responder preguntas.""",
    "Glassdoor: P1, 4 jun. 2026 (5/5); P2, 18 mar. 2026 (5/5); P3, 11 nov. 2021 (5/5); P4, 24 feb. 2023 (5/5).")

# 54 · Comentarios críticos de cultura y liderazgo
add_split_slide(54, "Comentarios críticos: cuando el problema no es técnico",
    "Los comentarios describen percepciones individuales; los patrones repetidos ayudan a formular mejores preguntas.",
    "Comentarios publicados", [
        "“Leadership frequently changes direction.” (N1)",
        "“Important conversations are often moved to private DMs.” (N2)",
        "“Projects often lacked clear direction.” (N3)",
        "“You have to wear ‘many hats’ working here.” (N4)",
    ],
    "Lectura profesional", [
        "N1: inestabilidad, retrabajo y dificultad para planificar",
        "N2: menor transparencia y riesgo de narrativas selectivas",
        "N3: ambigüedad, promesas incumplidas y crecimiento limitado",
        "N4: sobrecarga, reuniones y responsabilidades poco definidas",
    ],
    """[Tiempo sugerido: 4 minutos]
Texto para leer:
Las reseñas críticas analizadas describen cambios constantes de dirección, reorganizaciones, requisitos inestables, favoritismo, comunicación deficiente, conocimiento aislado, exceso de reuniones y pocas oportunidades de crecimiento. También mencionan ambientes donde la visibilidad o las relaciones parecen pesar más que la evidencia y los resultados.

Estas situaciones no se resuelven solamente siendo mejor programador o tester. Cuando no existe claridad, incluso una persona técnicamente brillante puede trabajar mucho en algo que dejará de ser prioridad. Cuando el conocimiento permanece en la cabeza de pocas personas, el equipo depende de héroes y repite errores. Cuando el crédito no se asigna de manera justa, disminuyen la confianza y la colaboración.

Una reseña individual representa una experiencia, no una sentencia definitiva sobre toda la organización. Debemos buscar patrones entre fechas, puestos y equipos, contrastar opiniones positivas y negativas y hacer preguntas durante la entrevista. También conviene observar si la empresa responde a las críticas con información concreta o solo con declaraciones generales.

El objetivo de leer Glassdoor no es buscar una empresa perfecta. Es identificar riesgos, formular preguntas y decidir cuáles condiciones podemos aceptar según nuestra etapa profesional, salud, responsabilidades y objetivos de crecimiento.""",
    "Glassdoor: N1, 25 jun. 2026 (1/5); N2, 23 ene. 2026 (2/5); N3, 8 ago. 2025 (1/5); N4, 2 feb. 2023 (2/5).")

# 55 · Respuesta profesional ante mala planificación o liderazgo difícil
add_process_slide(55, "Cómo responder profesionalmente en un entorno difícil",
    "No controlas todas las decisiones; sí puedes aportar claridad, proteger evidencia y elegir tus límites.",
    [
        ("Aclarar", "Pregunta por resultado, prioridad, responsable, fecha y criterio de terminado.", CYAN),
        ("Documentar", "Resume acuerdos y cambios en canales compartidos, sin acusaciones.", BLUE),
        ("Negociar", "Si entra una urgencia, pregunta qué alcance, fecha o tarea debe moverse.", PURPLE),
        ("Reducir riesgo", "Comparte conocimiento, automatiza, prueba, registra decisiones y evita héroes.", ORANGE),
        ("Escalar o salir", "Presenta hechos, impacto y opciones; define límites y conserva un plan alternativo.", GREEN),
    ],
    """[Tiempo sugerido: 4 minutos]
Texto para leer:
Cuando la planificación es débil, la primera habilidad es convertir ambigüedad en preguntas. Podemos decir: «Para confirmar, ¿cuál resultado tiene prioridad, quién toma la decisión y cómo sabremos que terminamos?». Si aparece una urgencia, no prometamos todo al mismo tiempo. Una respuesta profesional sería: «Puedo atender esta solicitud hoy; para hacerlo necesito mover la tarea A o reducir el alcance B. ¿Cuál opción prefieren?».

Después de una reunión, documente brevemente las decisiones en un canal compartido: prioridad, responsable, fecha, riesgos y asuntos pendientes. No se trata de crear evidencia para atacar a alguien, sino de reducir interpretaciones distintas y proteger la memoria del equipo. Cuando una conversación importante ocurre por mensaje privado, lleve el acuerdo relevante al espacio común con respeto.

Frente al favoritismo o la apropiación del trabajo, mantenga un registro profesional de contribuciones: tickets, pull requests, documentos, demostraciones y resultados. Reconozca públicamente a quienes participaron y pida el mismo estándar de trazabilidad para todos. Evite responder con chismes o ataques personales; describa comportamientos observables y su impacto.

Un profesional también ayuda a disminuir la dependencia de personas específicas mediante documentación, revisión entre pares, sesiones de transferencia y automatización. Sin embargo, las habilidades blandas no significan soportar cualquier ambiente. Si después de comunicar, negociar y escalar con evidencia persisten amenazas, discriminación, prácticas poco éticas o daño sostenido a la salud, es válido activar el plan B y buscar otro entorno.

Reflexión complementaria para leer:

Como profesionales debemos desarrollar control y autocontrol sobre la dirección de nuestra carrera. Así como las empresas defienden sus intereses, ajustan presupuestos, reorganizan equipos y toman decisiones según sus necesidades, nosotros también debemos proteger nuestros intereses profesionales. Eso implica mantener actualizadas las capacidades técnicas, fortalecer las habilidades blandas y evaluar con serenidad si el lugar donde estamos todavía contribuye a nuestro crecimiento.

Una frase que quiero que recuerden es esta: «Tu reputación profesional no se construye únicamente con lo que sabes, sino con la manera en que respondes cuando las cosas se complican». La reputación se fortalece cuando cumplimos, comunicamos a tiempo, reconocemos errores, respetamos a los demás, protegemos información y actuamos con coherencia. Cambiar de empleo no destruye una reputación cuando la salida se gestiona con profesionalidad; lo que puede dañarla es abandonar compromisos sin comunicar, reaccionar impulsivamente o convertir cada desacuerdo en un conflicto personal.

También quiero dejarles otra idea: «Tu paz mental forma parte de tu compensación profesional». Un salario puede ser atractivo, pero no siempre compensa un ambiente que deteriora constantemente la salud, la confianza o la capacidad de aprender. La experiencia no garantiza que una persona nunca sentirá estrés ni que siempre mantendrá bajo el cortisol. Lo que sí aporta la experiencia es mayor capacidad para reconocer patrones, organizar el trabajo, hacer preguntas, establecer límites y distinguir entre una dificultad normal de aprendizaje y un problema estructural del ambiente.

Cuando una tarea nos genera mucha tensión, debemos preguntarnos con honestidad: ¿me estreso porque todavía me falta una competencia técnica?, ¿porque no comprendí el dominio o necesito más práctica?, ¿porque no pedí ayuda a tiempo?, ¿o porque los requisitos cambian constantemente, nadie define prioridades y la planificación hace imposible cumplir? Si el problema es una brecha propia, corresponde estudiar, practicar, solicitar mentoría y mejorar. Si poseemos las competencias, pero el entorno produce caos permanente, contradicciones, favoritismo o desgaste, entonces debemos reconocer que quizá no necesitamos esforzarnos más dentro del mismo sistema; quizá necesitamos movernos.

Saber cuándo «saltar de barco» también es una habilidad profesional. No significa renunciar ante la primera dificultad. Significa observar evidencia, intentar comunicar, proponer soluciones, negociar expectativas, establecer límites y evaluar los resultados. Si nada cambia y permanecer comienza a causar un daño sostenido, moverse puede ser una decisión madura. No debemos esperar a estar completamente agotados, resentidos o profesionalmente estancados para preparar una alternativa.

Las empresas hablan de transformación, reestructuración y adaptación al mercado. Una persona también tiene derecho a transformar y reconstruir su carrera. Puede decidir profundizar en su especialidad, moverse horizontalmente hacia otra función, buscar mayores ingresos, emprender, trabajar como independiente o avanzar hacia liderazgo. Permanecer diez años en una ruta horizontal puede ser una decisión excelente si coincide con sus intereses. Buscar liderazgo también puede serlo. Ningún camino es exitoso por definición; el éxito depende de lo que la persona valora, de las responsabilidades que desea asumir y del costo que está dispuesta a aceptar.

Una relación laboral comienza con una negociación. La empresa presenta una oferta y unas necesidades; la persona decide si ofrece sus servicios, conocimientos y tiempo bajo esas condiciones. Después ambas partes deben adaptarse y cumplir lo acordado. Dentro del contrato y de la legislación aplicable, la empresa puede prescindir de los servicios cuando cambian sus necesidades, y el profesional también puede renunciar cuando entiende que debe continuar su carrera en otro lugar. No conviene interpretar automáticamente cada salida como una derrota personal. Existen emociones y consecuencias humanas, pero también existen decisiones de negocio.

Negociar bien no significa pedir la cifra más alta ni aceptar la más baja. Significa comprender el valor del puesto, las responsabilidades, el aprendizaje, el riesgo, los beneficios, el horario, la modalidad y las oportunidades futuras. La madurez se observa en la actitud con la que reaccionamos ante los requisitos: preguntar antes de asumir, reconocer lo que sabemos y lo que no sabemos, proponer opciones y no prometer lo imposible.

Por eso, mantengan siempre un plan de desarrollo y un plan alternativo. Actualicen el CV, cuiden las relaciones, documenten resultados, continúen aprendiendo y observen el mercado incluso cuando tengan empleo. No para vivir con miedo, sino para conservar capacidad de decisión. La manera en que interpreten el mercado influirá en su carrera: pueden verlo como algo que controla completamente su destino o como un entorno que deben aprender a leer, negociar y navegar con inteligencia.

Tres frases útiles son: «Veo este riesgo y propongo estas dos opciones»; «Confirmo que la prioridad cambió de A hacia B y que la fecha de A será revisada»; y «No tengo todavía la respuesta, pero investigaré, validaré y regresaré a esta hora». Esa combinación de honestidad, claridad y compromiso suele comunicar más madurez que fingir seguridad.""",
    "Aplicación práctica derivada de los patrones de las reseñas de Glassdoor citadas; documentar hechos, no rumores.")

# 56 · Referencias concretas de los posts de Glassdoor compartidos
add_matrix_slide(56, "Posts de Glassdoor utilizados como casos de reflexión",
    "Las empresas se omiten deliberadamente; se conservan fecha, puntuación, perfil y patrón analizado.",
    ["Caso", "Fecha y puntuación", "Perfil publicado", "Patrón utilizado en el taller"],
    [
        ("P1", "4 jun. 2026 · 5/5", "Empleado actual · <1 año", "Ownership, apoyo, feedback y contratación intencional"),
        ("P2", "18 mar. 2026 · 5/5", "Empleado actual · >1 año", "Mejor ritmo operativo; urgencias y conocimiento aislado"),
        ("P3", "11 nov. 2021 · 5/5", "Lead QA Engineer · actual", "Aprendizaje continuo y oportunidades para crear una trayectoria"),
        ("P4", "24 feb. 2023 · 5/5", "Revenue Operations · actual", "Cultura amable, ayuda entre compañeros y enfoque de equipo"),
        ("N1", "25 jun. 2026 · 1/5", "Software Engineer · actual", "Dirección reactiva, poco crecimiento y silos técnicos"),
        ("N2", "23 ene. 2026 · 2/5", "Senior Software Engineer · actual", "Favoritismo, crédito desigual y comunicación poco transparente"),
        ("N3", "8 ago. 2025 · 1/5", "Exempleado anónimo", "Reorganizaciones, rotación y progresión profesional limitada"),
        ("N4", "2 feb. 2023 · 2/5", "Exempleado anónimo", "Exceso de reuniones, múltiples roles y dificultad para aprender"),
    ],
    """[Tiempo sugerido: 3 minutos]
Texto para leer:
Esta diapositiva identifica los posts de Glassdoor utilizados para construir los ejemplos anteriores. Las empresas se omiten deliberadamente porque el propósito no es atacar una organización, sino aprender a reconocer patrones de cultura, liderazgo y trabajo.

El caso P1 corresponde al post titulado “Exceptional culture with ownership mentality and intentional hiring”, publicado el 4 de junio de 2026 con una puntuación de cinco sobre cinco por una persona empleada durante menos de un año. Se utilizaron sus referencias a responsabilidad compartida, apoyo, apertura al feedback, orientación al cliente y contratación cuidadosa.

El caso P2 corresponde a “Figuring things out”, del 18 de marzo de 2026, también con cinco sobre cinco. Se tomó el contraste entre una mejora en los ritmos operativos y dos desafíos: solicitudes urgentes frecuentes y conocimiento concentrado en pocas personas.

El caso P3 corresponde a una reseña de un Lead QA Engineer publicada el 11 de noviembre de 2021 con cinco sobre cinco. Se utilizó su comentario sobre la posibilidad de aprender continuamente y construir un espacio profesional alineado con intereses y talentos.

El caso P4 corresponde a una reseña de Revenue Operations del 24 de febrero de 2023 con cinco sobre cinco. Se utilizó su énfasis en una cultura amable, colaborativa y dispuesta a ayudar.

El caso N1 es “Reactive environment with little room to grow as an engineer”, del 25 de junio de 2026, con una puntuación de uno sobre cinco. Sus temas principales fueron cambios reactivos de dirección, trabajo dominado por urgencias, poco espacio para crecimiento técnico y conocimiento aislado entre equipos.

El caso N2 es “Good pay, but culture rewards self-promotion, favoritism, and narrative control over real work”, del 23 de enero de 2026, con dos sobre cinco. Se utilizaron sus señalamientos sobre inestabilidad, favoritismo, reconocimiento desigual y conversaciones importantes trasladadas a canales privados.

El caso N3 es “Unstable Environment with Limited Growth Opportunities”, del 8 de agosto de 2025, con uno sobre cinco. Se consideraron las reorganizaciones frecuentes, la rotación, el apoyo limitado y la falta de progresión profesional.

El caso N4 es “Decent For Some But Not Everyone”, del 2 de febrero de 2023, con dos sobre cinco. Se tomaron como temas el exceso de reuniones, la dificultad para obtener ayuda, el poco tiempo para aprender y la expectativa de asumir múltiples funciones.

Estas reseñas expresan experiencias individuales y no prueban por sí solas cómo funciona toda una organización. Deben analizarse junto con otras publicaciones, fechas, puestos, entrevistas y fuentes institucionales. Aquí se utilizan como casos pedagógicos para aprender a preguntar, contrastar y decidir.""",
    "Fuente: posts de Glassdoor suministrados por la expositora; síntesis y anonimización para fines educativos.",
    col_widths=[0.72, 2.12, 2.82, 6.46])

# 57 · Tradición liberal y orientada al mercado
add_cards_slide(57, "Frases · tradición liberal y orientada al mercado",
    "Agrupación orientativa solicitada en rojo: incentivos, precios, beneficios y costos de coordinación.",
    [
        ("Adam Smith · 1776", "«No es por la benevolencia del carnicero, el cervecero o el panadero que esperamos nuestra cena». El intercambio surge cuando cada parte ofrece algo que la otra valora. Fuente: La riqueza de las naciones, I.II.", RED),
        ("Milton Friedman · 1962", "«Usar sus recursos [...] para aumentar sus beneficios, mientras permanezca dentro de las reglas del juego». La búsqueda de rentabilidad está limitada por ley, contratos y competencia sin fraude. Fuente: Capitalism and Freedom, p. 133.", RED),
        ("F. A. Hayek · 1945", "«Lo más significativo de este sistema es la economía de conocimiento con la que opera». Los precios transmiten señales sobre escasez y cambios que ninguna persona conoce por completo. Fuente: AER, sección VI.", RED),
        ("Ronald Coase · 1937", "«La principal razón para establecer una empresa [...] es que usar el mecanismo de precios tiene un costo». Buscar, negociar, contratar y supervisar también consumen recursos. Fuente: Economica, pp. 388–392.", RED),
    ],
    """[Tiempo sugerido: 3 minutos]
Texto para leer:
Desde una perspectiva de mercado, una empresa intenta producir bienes o servicios utilizando una combinación de recursos cuyo costo sea sostenible frente al valor que el cliente está dispuesto a pagar. Esos recursos incluyen dinero, tiempo, trabajo humano, tecnología, instalaciones, información, proveedores y riesgo. Por eso la meta económica no es necesariamente utilizar la menor cantidad de personas, sino organizar los recursos de la manera que permita cumplir con calidad, continuidad y rentabilidad.

Adam Smith explicó que el intercambio no depende únicamente de la benevolencia, sino del interés de cada parte en ofrecer algo que la otra valora. La cita presentada procede de La riqueza de las naciones, libro I, capítulo II. En una relación laboral también existe intercambio: la organización ofrece compensación y condiciones; el profesional ofrece tiempo, conocimiento, ejecución y responsabilidad.

Milton Friedman sostuvo que, en una economía libre, la empresa utiliza sus recursos para aumentar beneficios mientras permanece dentro de las reglas del juego. La parte final es indispensable: legalidad, competencia abierta, contratos y ausencia de engaño o fraude. La cita procede de Capitalism and Freedom, edición conmemorativa de 2002, página 133, publicada originalmente en 1962.

Friedrich A. Hayek explicó que el sistema de precios comunica información dispersa sobre escasez y cambios. Ninguna persona posee todo el conocimiento del mercado. La cita se basa en la sección VI de “The Use of Knowledge in Society”, publicado en 1945. Esto ayuda a comprender por qué salarios, demanda de habilidades, tecnologías y vacantes cambian cuando cambian las condiciones.

Ronald Coase mostró que utilizar el mercado también tiene costos: buscar, negociar, contratar, coordinar y supervisar. Su planteamiento aparece en “The Nature of the Firm”, de 1937. Una empresa decide qué realizar internamente, qué automatizar y qué contratar externamente comparando costos, control, riesgo y capacidad.

Esta lógica no significa que la persona con el salario más bajo sea siempre la alternativa menos costosa. Una contratación deficiente puede producir errores, retrabajo, incidentes, rotación, conocimiento perdido y daño reputacional. Tampoco significa que la empresa deba sacrificar calidad, seguridad o dignidad para reducir gastos. La eficiencia sostenible considera el costo total y las consecuencias de largo plazo.

Para el profesional, la conclusión es clara: no compita solamente ofreciendo cobrar menos. Demuestre que puede resolver problemas, aprender, reducir riesgos, colaborar y generar resultados. La empresa defenderá sus intereses económicos y usted debe defender responsablemente el valor de su trabajo, su desarrollo, su reputación y sus límites. Una buena negociación busca una relación en la que ambas partes comprendan qué ofrecen, qué reciben y bajo cuáles condiciones.""",
    "Fuentes: Smith, The Wealth of Nations, I.II (1776); Friedman, Capitalism and Freedom, p. 133 (1962/2002); Hayek, AER, VI (1945); Coase, Economica, pp. 388–392 (1937).")
for quote_shape in prs.slides[-1].shapes:
    if not getattr(quote_shape, "has_text_frame", False) or "«" not in quote_shape.text:
        continue
    for paragraph in quote_shape.text_frame.paragraphs:
        paragraph.font.size = Pt(10.5)
        for run in paragraph.runs:
            run.font.size = Pt(10.5)

# 59 · Perspectivas socialistas, comunistas y de centroizquierda
add_cards_slide(59, "Frases · perspectivas socialistas y de centroizquierda",
    "Agrupación orientativa solicitada en azul: trabajo, empleo, desigualdad, dignidad y libertad real.",
    [
        ("Karl Marx · 1844", "«El trabajador solo se siente fuera de sí en su trabajo, y consigo mismo fuera del trabajo». Marx critica la alienación cuando la actividad laboral pierde autonomía y sentido. Fuente: Manuscritos económico-filosóficos, “Trabajo enajenado”.", BLUE),
        ("J. M. Keynes · 1936", "«Las fallas sobresalientes [...] son no proporcionar pleno empleo y distribuir arbitraria y desigualmente la riqueza y los ingresos». Keynes propone corregir fallas del capitalismo sin abolirlo. Fuente: Teoría general, cap. 24.", BLUE),
        ("Amartya Sen · 1999", "«El desarrollo puede verse [...] como un proceso de expansión de las libertades reales que disfrutan las personas». El ingreso importa, pero también la capacidad efectiva para elegir y vivir dignamente. Fuente: Development as Freedom, p. 3.", BLUE),
        ("Joseph Stiglitz · 2013", "«La desigualdad es una elección». Las reglas, instituciones y políticas influyen en cómo se distribuyen oportunidades, riesgos y resultados; no todo es consecuencia inevitable del mercado. Fuente: “Inequality Is a Choice”.", BLUE),
    ],
    """[Tiempo sugerido: 3 minutos]
Texto para leer:
Esta segunda diapositiva presenta perspectivas críticas o reformistas. La clasificación por colores es orientativa y responde al diseño solicitado. Karl Marx pertenece a la tradición comunista y socialista. John Maynard Keynes fue un reformista del capitalismo, no un comunista. Amartya Sen y Joseph Stiglitz suelen relacionarse con enfoques progresistas o de centroizquierda, pero sus obras son más amplias que una etiqueta partidaria.

Marx analiza la alienación: el trabajo puede convertirse en una actividad externa, controlada por otros y separada del sentido personal. Esta idea ayuda a reflexionar sobre ambientes donde la persona ejecuta sin comprender, carece de autonomía o siente que su identidad queda reducida a producir resultados.

Keynes identifica como fallas centrales la ausencia de pleno empleo y la distribución desigual de riqueza e ingresos. Su respuesta no consiste en eliminar los mercados, sino en reconocer que pueden existir periodos prolongados de desempleo y que las instituciones públicas tienen un papel estabilizador.

Amartya Sen amplía el concepto de desarrollo. No basta con medir salario, producción o crecimiento; también importa si las personas poseen libertades reales para educarse, cuidar su salud, participar, elegir y construir el tipo de vida que valoran. Aplicado a la carrera, un empleo puede pagar bien y aun así limitar aprendizaje, autonomía o bienestar.

Joseph Stiglitz sostiene que la desigualdad está influida por decisiones institucionales y políticas. Las reglas sobre competencia, educación, impuestos, negociación y protección social cambian la distribución de oportunidades. Esto recuerda que el resultado del mercado no surge en un vacío: depende de normas y relaciones de poder.

Estas perspectivas no deben utilizarse como consignas para aceptar o rechazar automáticamente una empresa. Sirven para ampliar las preguntas: además de eficiencia y beneficio, ¿cómo se distribuyen riesgos y recompensas?, ¿qué autonomía existe?, ¿qué oportunidades de aprendizaje ofrece el trabajo?, ¿cómo se protege la dignidad y qué libertad real conserva la persona?

Al comparar ambas diapositivas, el objetivo no es declarar un ganador ideológico. Es comprender que la empresa observa costos, precios e incentivos, mientras el profesional también debe observar dignidad, crecimiento, seguridad, libertad y poder de negociación. Una decisión madura integra ambas dimensiones.""",
    "Fuentes: Marx, Manuscritos de 1844; Keynes, Teoría general, cap. 24 (1936); Sen, Development as Freedom, p. 3 (1999); Stiglitz, “Inequality Is a Choice” (2013).")
for quote_shape in prs.slides[-1].shapes:
    if not getattr(quote_shape, "has_text_frame", False) or "«" not in quote_shape.text:
        continue
    for paragraph in quote_shape.text_frame.paragraphs:
        paragraph.font.size = Pt(10.5)
        for run in paragraph.runs:
            run.font.size = Pt(10.5)

# 58 · Expectativa educativa frente a la realidad profesional
add_bullet_slide(58, "El título prepara, pero no garantiza un empleo",
    "Graduarse es un logro y una base; convertirlo en carrera exige adaptación, evidencia y lectura del entorno.",
    [
        "Aprendimos que completar la ingeniería o carrera técnica conduciría a un trabajo seguro.",
        "El mercado crea vacantes cuando las organizaciones tienen una necesidad, no cuando terminamos de estudiar.",
        "El trabajo incluye culturas distintas, políticas internas, egos, negociación y conflictos.",
        "Siempre existirán brechas técnicas porque las herramientas y las funciones continúan cambiando.",
        "El idioma puede limitar documentación, entrevistas, colaboración internacional y oportunidades remotas.",
        "La respuesta es gestionar la carrera: aprender, demostrar, comunicar, relacionarse y conservar alternativas.",
    ],
    """[Tiempo sugerido: 4 minutos]
Texto para leer:
A muchos se nos educó con una secuencia aparentemente segura: graduarnos de la preparatoria, completar cuatro o cinco años de ingeniería, licenciatura o formación técnica y después obtener un trabajo estable. Esa expectativa dio estructura y motivación, pero estaba incompleta. La formación académica es valiosa; sin embargo, un título no funciona como un contrato anticipado con el mercado laboral.

No siempre nos explicaron que, además del conocimiento de la carrera, tendríamos que interpretar culturas organizacionales, políticas internas, estilos de liderazgo, egos, conflictos, prioridades cambiantes y relaciones de poder. Tampoco que necesitaríamos negociar expectativas, comunicar desacuerdos, pedir ayuda, establecer límites y reconocer cuándo un problema es técnico y cuándo pertenece al sistema de trabajo.

La universidad puede ofrecer fundamentos, disciplina, pensamiento lógico y oportunidades para aprender. Pero ninguna institución puede enseñarnos de antemano todas las herramientas, industrias y situaciones humanas que encontraremos. Siempre aparecerán lagunas técnicas. Lo profesional no es saberlo todo; es identificar la brecha, investigar, practicar, validar y pedir colaboración sin fingir dominio.

También existen barreras de idioma. En tecnología, una parte importante de la documentación, las comunidades, los mensajes de error, las entrevistas y las vacantes internacionales utiliza inglés. No dominarlo todavía no invalida nuestras capacidades, pero sí puede limitar el acceso. Por eso debe tratarse como una competencia profesional que se desarrolla progresivamente: lectura técnica, vocabulario del puesto, escucha, escritura y conversación.

Las vacantes no aparecen cuando una persona se siente completamente preparada. Aparecen cuando una empresa tiene presupuesto, urgencia y una necesidad concreta. Si no podemos responder en ese momento, otra persona lo hará. Esto no debe producir miedo, sino una actitud de preparación continua: observar el mercado, mantener evidencias, fortalecer redes, practicar entrevistas y actualizar el perfil antes de necesitarlo.

El mensaje no es que estudiar no sirve. El mensaje es que la educación formal es el inicio de una carrera que debemos aprender a gestionar. Al título debemos sumar experiencia demostrable, habilidades blandas, criterio, idioma, networking, negociación y capacidad para movernos cuando cambian nuestras metas o el entorno.

La seguridad profesional no proviene únicamente de permanecer muchos años en una empresa. También proviene de saber que podemos aprender, demostrar valor, construir relaciones y crear alternativas. La carrera deja de ser una promesa externa y se convierte en una responsabilidad que dirigimos con información y autocontrol.""",
    PURPLE,
    ("NUEVA ECUACIÓN", "Título + evidencia + habilidades blandas + idioma + adaptación", "La carrera se construye después de graduarse."),
    "La educación formal aporta fundamentos; la empleabilidad requiere actualización y acción continuas.")

# 60 · Empresas de etapa y gratitud profesional saludable
add_split_slide(60, "Empresas de etapa: aprender, aportar y cerrar ciclos",
    "Un buen comienzo profesional puede ser valioso aunque el proyecto, el contrato o la relación laboral no sean permanentes.",
    "Comentarios publicados", [
        "“Nice place to work, great opportunity to improve your skills, great team work.” (C1)",
        "“The work environment with my colleagues is exceptionally positive.” (C2)",
        "“Very friendly and welcoming team.” (C3)",
        "“Collaborative team. Flexible hours.” (C4)",
    ],
    "Lectura profesional", [
        "Algunas empresas son una etapa de aprendizaje, no un destino permanente",
        "Un proyecto puede terminar o las necesidades de la organización pueden cambiar",
        "«Y todo lo que hagáis, hacedlo de corazón, como para el Señor y no para los hombres». — Colosenses 3:23",
        "Agradece la oportunidad sin convertirla en una deuda infinita",
        "Mantén CV, portafolio, red profesional y finanzas personales preparados",
        "Si no eres accionista, tu vínculo es profesional: aporta valor y conserva autonomía",
    ],
    """[Tiempo sugerido: 4 minutos]
Texto para leer:
Estas reseñas de Santo Domingo muestran que una empresa puede ser un buen espacio para comenzar, aprender tecnologías, recibir apoyo de compañeros y adquirir experiencia. Al mismo tiempo, algunas personas describen desalineación entre el título y las funciones, estancamiento, múltiples proyectos, micromanagement o tareas adicionales producidas por el bajo desempeño de otros compañeros.

El caso C1 corresponde a una reseña de DevOps Engineer publicada el 5 de noviembre de 2024 con cuatro sobre cinco. El caso C2 fue publicado por un Junior Software Engineer el 24 de agosto de 2024 con tres sobre cinco. El caso C3 pertenece a un antiguo Quality Assurance Engineer y fue publicado el 3 de mayo de 2023 con cinco sobre cinco. El caso C4 corresponde a un antiguo Software QA Engineer y fue publicado el 9 de febrero de 2023 con cuatro sobre cinco. Todas identifican Santo Domingo como ubicación y se presentan sin mencionar la empresa.

Hay organizaciones en las que entraremos durante una etapa corta. Puede terminar el proyecto, cambiar el presupuesto, desaparecer la necesidad o la empresa decidir prescindir de nuestros servicios conforme al contrato y la legislación aplicable. Nosotros también podemos decidir cerrar el ciclo cuando el puesto deja de aportar aprendizaje, estabilidad, bienestar o coherencia con nuestros objetivos.

Trabajar cada día como si pudiera ser el último no significa vivir con miedo. Significa mantener el trabajo documentado, entregar calidad, comunicar el estado real, evitar conocimiento escondido y dejar una transferencia que proteja al equipo. También significa conservar actualizado el CV, el portafolio, los contactos profesionales y, cuando sea posible, un fondo financiero para responder a cambios inesperados.

Para quienes conectan con una perspectiva de fe, Colosenses 3:23 expresa: «Y todo lo que hagáis, hacedlo de corazón, como para el Señor y no para los hombres». Esta reflexión invita a trabajar con excelencia por convicción y valores, no únicamente para recibir aprobación de un jefe. En un público diverso puede traducirse como actuar con integridad incluso cuando nadie está observando.

Debemos agradecer sanamente a quien nos ofrece una oportunidad. Agradecer no significa aceptar una deuda infinita ni pensar que debemos permanecer para siempre porque una empresa nos contrató cuando comenzábamos. La organización recibió nuestro trabajo y nosotros recibimos compensación, experiencia y aprendizaje. Fue un intercambio profesional que puede ser valioso sin convertirse en un apego enfermizo.

Si no somos accionistas, la empresa no nos pertenece y tampoco controlamos todas sus decisiones. Esto no reduce la importancia de nuestro trabajo; aclara la relación. Nuestra responsabilidad es aportar valor, respetar acuerdos, cuidar la reputación y preparar el siguiente paso. La lealtad profesional saludable es compatible con reconocer cuándo un ciclo terminó.

No conviertan una salida en traición ni una permanencia prolongada en la única prueba de compromiso. Evalúen resultados, aprendizaje, condiciones, salud, ingresos y dirección profesional. Entrar, aportar, aprender y cerrar bien un ciclo también puede ser una historia de éxito.""",
    "Glassdoor, reseñas anónimas de Santo Domingo: C1, 5 nov. 2024 (4/5); C2, 24 ago. 2024 (3/5); C3, 3 may. 2023 (5/5); C4, 9 feb. 2023 (4/5).")
for stage_shape in prs.slides[-1].shapes:
    if not getattr(stage_shape, "has_text_frame", False) or not stage_shape.text.startswith("•"):
        continue
    for paragraph in stage_shape.text_frame.paragraphs:
        paragraph.font.size = Pt(13.8)
        paragraph.space_after = Pt(6)
        for run in paragraph.runs:
            run.font.size = Pt(13.8)


# Ampliaciones del guion. Estas notas están diseñadas para poder leerse de forma
# natural y también incluyen ejemplos, preguntas y transiciones para la sesión.
NOTE_EXPANSIONS = {
    1: """Texto ampliado para leer:
Cuando hablamos de transformación del QA, no hablamos únicamente de aprender una herramienta nueva. Hablamos de un cambio en la manera de producir software, organizar equipos y demostrar valor. Antes, muchas organizaciones esperaban que una persona probara el sistema cuando el desarrollo había terminado. Hoy se espera que la calidad se trabaje desde las preguntas iniciales, continúe durante la construcción y se observe después de liberar el producto.

La inteligencia artificial puede acelerar análisis, documentación y automatización, pero también crea riesgos nuevos. Por eso el profesional de calidad no pierde importancia: necesita comprender el producto, cuestionar los resultados, proteger la información y comunicar decisiones. Durante la sesión relacionaremos cada concepto con una pregunta práctica: ¿cómo me ayuda esto a conseguir trabajo, crecer o mantenerme vigente?""",

    2: """Texto ampliado para leer:
Quiero que observen que los cuatro resultados de esta diapositiva forman un sistema. La dirección profesional indica hacia dónde vamos. La ruta de aprendizaje define qué estudiar. Las evidencias permiten demostrarlo. El plan de acción convierte la intención en fechas y actividades.

Una persona puede completar diez cursos y continuar sin saber a qué empleo aplicar. Otra puede estudiar menos contenidos, pero seleccionar los que aparecen en sus vacantes objetivo, construir un proyecto relacionado y explicar sus decisiones. Generalmente, la segunda persona presenta una historia profesional más clara.

Mientras avancemos, no se pregunten solamente «¿entiendo este concepto?». Pregúntense también «¿qué podría construir con esto?», «¿cómo lo explicaría en una entrevista?» y «¿qué evidencia mostraría en mi CV o portafolio?».""",

    4: """Texto ampliado para leer:
Esta actividad también nos ayuda a reconocer que el público puede encontrarse en momentos distintos. Quien está iniciando necesita fundamentos y práctica guiada. Quien trabaja en QA manual puede fortalecer API, SQL, análisis y automatización. Quien ya automatiza puede avanzar hacia arquitectura, integración continua, rendimiento, seguridad, inteligencia artificial o liderazgo.

No comparen su capítulo uno con el capítulo diez de otra persona. Lo importante es identificar un siguiente paso que sea exigente pero alcanzable. Si alguien desea cambiar de área, también debe reconocer cuáles competencias ya son transferibles. Por ejemplo, una persona de soporte suele conocer incidentes y usuarios; alguien de desarrollo comprende código; una persona de negocio comprende procesos. Todas esas experiencias pueden convertirse en ventajas dentro de QA.""",

    6: """Texto ampliado para leer:
El informe del Foro Económico Mundial muestra una combinación importante: crecen las competencias tecnológicas relacionadas con IA y datos, pero el pensamiento analítico, la resiliencia, la colaboración y la adaptación continúan siendo centrales. Para QA esto tiene mucho sentido. Las herramientas pueden ejecutar rápidamente, pero el profesional debe decidir qué merece atención y cómo interpretar un resultado.

LinkedIn también identifica la alfabetización en inteligencia artificial entre las habilidades de crecimiento acelerado. Alfabetización no significa aceptar cualquier respuesta generada. Significa saber formular una tarea, entregar contexto, reconocer limitaciones, verificar resultados y utilizar la herramienta con responsabilidad.

La pregunta para el público es: si una parte de las competencias cambiará, ¿qué capacidad nos permite seguir siendo relevantes? La respuesta es aprender de forma continua y demostrar que podemos transferir fundamentos a nuevas herramientas.""",

    7: """Texto ampliado para leer:
Pensemos en dos profesionales. El primero ejecuta siempre el mismo conjunto de pasos y copia el resultado. El segundo comprende el riesgo, selecciona pruebas, investiga comportamientos inesperados y explica el impacto. La automatización puede sustituir buena parte del trabajo repetitivo del primero, pero amplifica la productividad del segundo.

La oportunidad está en desplazarse hacia actividades de mayor valor: análisis de requisitos, diseño basado en riesgo, pruebas exploratorias, calidad de datos, observabilidad, seguridad y comunicación con negocio. Esto no significa abandonar la ejecución manual. Significa utilizarla donde aporta conocimiento y automatizar donde aporta velocidad y repetibilidad.

Pregunte al grupo: ¿qué parte de su trabajo actual se repite y podría automatizarse? Luego pregunte: si recuperaran ese tiempo, ¿qué riesgo o problema podrían analizar mejor?""",

    8: """Texto ampliado para leer:
En la etapa de descubrir, QA pregunta quién usará el producto y qué daño produciría una falla. Al definir, convierte reglas vagas en criterios verificables. Durante la construcción, puede revisar ejemplos, contratos de API, manejo de errores, accesibilidad y pruebas del código. En la validación, combina pruebas manuales, automáticas y no funcionales. Finalmente, al observar producción, aprende de métricas, incidentes y comportamiento real.

Este ciclo abre especialidades distintas. Una persona interesada en negocio puede profundizar en requisitos. Otra interesada en programación puede trabajar automatización o SDET. Quien disfruta analizar sistemas puede avanzar hacia rendimiento, seguridad u observabilidad. Quien se interesa por usuarios puede especializarse en accesibilidad, usabilidad o experiencia.

La calidad no es una única actividad; es una perspectiva que puede aplicarse en todo el ciclo.""",

    9: """Texto ampliado para leer:
Dentro de QA existen muchas más especialidades de las que caben en la diapositiva. Podemos mencionar QA funcional, test analyst, QA de API, tester de bases de datos, mobile tester, automation engineer, SDET, performance engineer, security tester, accessibility tester, localization tester, data quality engineer, ETL tester, AI o machine learning tester, reliability engineer, QA Lead, Test Manager y Quality Coach.

QA funcional profundiza en requisitos, negocio y comportamiento. API y datos se concentran en servicios, contratos e integridad. Automatización y SDET requieren programación y diseño técnico. Rendimiento estudia carga, capacidad y estabilidad. Seguridad analiza vulnerabilidades y riesgos. Accesibilidad valida que el producto pueda utilizarse por personas con diferentes capacidades. Testing de IA trabaja con datos, modelos, sesgo, variabilidad y drift.

No deben aprenderlas todas. Utilicen esta lista para explorar y luego revisen cuáles aparecen en el mercado que les interesa.""",

    10: """Texto ampliado para leer:
También hay profesiones cercanas a las que una persona de QA puede migrar. En análisis de negocio puede aprovechar su experiencia aclarando requisitos y procesos. En producto puede utilizar su comprensión del usuario, los riesgos y el valor. En DevOps o release management puede aplicar automatización, ambientes, pipelines y monitoreo. En ciberseguridad puede avanzar hacia application security, pruebas de penetración o DevSecOps.

En datos puede trabajar calidad de datos, validación de ETL, gobierno o pruebas de analítica. En inteligencia artificial puede especializarse en evaluación de modelos y sistemas generativos. También puede evolucionar hacia desarrollo, gestión de proyectos, Scrum Master, UX research, accesibilidad, soporte avanzado, documentación técnica, consultoría, docencia o liderazgo de calidad.

La clave es reconocer la competencia transferible y la brecha adicional. QA ofrece una base, pero cada transición exige práctica específica y evidencia nueva.""",

    11: """Texto ampliado para leer:
Si una persona todavía no conoce el nombre exacto del puesto, puede comenzar describiendo el tipo de trabajo que disfruta. ¿Le gusta hablar con usuarios y aclarar procesos? Puede explorar análisis funcional o negocio. ¿Disfruta programar y resolver problemas técnicos? Puede revisar automation o SDET. ¿Le interesan métricas, estabilidad y sistemas? Performance, observabilidad o DevOps pueden ser opciones. ¿Le preocupa la protección de información? Seguridad y DevSecOps merecen explorarse.

Después de elegir una hipótesis, busque al menos diez vacantes. La elección inicial puede cambiar. Eso no es fracaso; es aprendizaje basado en evidencia. El objetivo de esta actividad es reducir el campo y comenzar una investigación profesional concreta.""",

    13: """Texto ampliado para leer:
La IA puede utilizarse como asistente en cada etapa, pero nunca debe convertirse en una caja negra que toma decisiones sin revisión. Al analizar requisitos, puede proponer preguntas que no habíamos considerado. Al diseñar pruebas, puede sugerir particiones, valores límite y combinaciones. Durante la ejecución, puede ayudar a interpretar mensajes de error y comparar grandes volúmenes de texto. En documentación, puede mejorar claridad y estructura.

Para aprender, puede explicar código, generar ejercicios, simular entrevistas y adaptar ejemplos. Sin embargo, la persona debe intentar resolver antes de solicitar una solución completa. Una técnica útil es pedir pistas progresivas: primero una pregunta, luego una orientación y finalmente un ejemplo. Así la IA acompaña el razonamiento en lugar de reemplazarlo.""",

    14: """Texto ampliado para leer:
Vamos a imaginar que recibimos una historia que dice: «Como cliente quiero transferir dinero para pagar rápidamente». Esa frase no aclara límites, monedas, cuentas bloqueadas, saldo insuficiente, doble envío, permisos, horarios ni confirmación. Un prompt con contexto puede ayudarnos a preparar una conversación con negocio.

Observe que la salida solicitada no es solamente una lista de casos. Primero pedimos preguntas, luego riesgos, después escenarios y datos. Ese orden obliga a comprender antes de ejecutar. Cuando llegue la respuesta, clasifique cada elemento como confirmado, supuesto o pendiente.

Una buena práctica es conservar el requisito original, la salida generada y la versión validada. Esa trazabilidad permite explicar qué aportó la IA y qué decidió el profesional.""",

    16: """Texto ampliado para leer:
Un error frecuente es compartir información real porque parece más cómodo. Antes de utilizar IA, sustituya nombres, correos, cuentas, tokens, direcciones y datos de clientes. Si el código pertenece a una organización, confirme qué herramientas están autorizadas. La ausencia de una prohibición explícita no equivale a permiso.

También debemos revisar si el contenido generado reproduce sesgos o crea una falsa sensación de seguridad. Una suite generada automáticamente puede ejecutar correctamente y aun así probar lo equivocado. Un reporte muy bien escrito puede describir una conclusión falsa.

Como regla práctica: nada producido por IA debe convertirse en requisito, evidencia, código o decisión sin una forma de verificación. Pregunte siempre: ¿cómo sé que esto es correcto?, ¿qué fuente o ejecución lo demuestra?""",

    17: """Texto ampliado para leer:
Testing de inteligencia artificial es una especialidad con crecimiento propio. En sistemas convencionales solemos esperar que la misma entrada produzca una salida definida. En modelos probabilísticos puede existir variación y no siempre hay un oráculo sencillo. Esto obliga a trabajar con métricas, conjuntos de evaluación, tolerancias y revisión experta.

En datos se valida completitud, representatividad, calidad de etiquetas y privacidad. En modelos se estudian precisión, robustez, ataques y degradación. En sistemas generativos se evalúan alucinaciones, contenido dañino, prompt injection, filtración de información y resistencia al abuso.

Una persona interesada en esta ruta puede combinar fundamentos de QA con datos, estadística básica, machine learning, ética, seguridad y técnicas de evaluación. No necesita dominar todo desde el primer día, pero sí construir una base progresiva.""",

    20: """Texto ampliado para leer:
Para decidir qué automatizar, podemos utilizar cuatro preguntas. Primera: ¿qué riesgo cubre esta prueba? Segunda: ¿con qué frecuencia se ejecutará? Tercera: ¿el comportamiento es suficientemente estable? Cuarta: ¿podemos determinar automáticamente si pasó o falló?

Un flujo de pago crítico que se ejecuta en cada versión suele ser buen candidato. Una interfaz que cambia diariamente quizá no. Una revisión de claridad visual o facilidad de uso puede requerir juicio humano. Algunas pruebas se automatizan parcialmente: preparar datos, navegar y capturar evidencia, dejando la valoración final a una persona.

La automatización profesional también incluye mantenimiento, documentación, ambientes, datos, reportes y análisis de fallos. El script es solo una parte del sistema.""",

    21: """Texto ampliado para leer:
La ruta puede adaptarse según el puesto. Para QA manual junior, la prioridad puede ser requisitos, casos, defectos, Jira, SQL y API básica. Para automation, se añade un lenguaje, Git, estructura de proyecto, localizadores, aserciones, patrones y CI. Para SDET se profundiza en ingeniería de software, arquitectura, servicios, código limpio y herramientas internas.

Para performance se requieren conceptos de carga, concurrencia, métricas y observabilidad. Para seguridad, redes, autenticación, amenazas y OWASP. Para datos, SQL más profundo, ETL, calidad y pipelines. Para mobile, plataformas, dispositivos, permisos y redes.

La pregunta no es cuál ruta es superior. La pregunta es cuál coincide con las tareas que usted desea realizar y con las oportunidades que encuentra.""",

    22: """Texto ampliado para leer:
Observe que la tabla no obliga a utilizar una marca específica. Jira puede reemplazarse por otra herramienta de gestión; Playwright puede sustituirse por Cypress o Selenium; GitHub Actions por GitLab CI, Jenkins o Azure DevOps. Lo transferible es comprender la función que cumple cada herramienta.

Cuando una vacante mencione una tecnología que no conoce, primero identifique la categoría. Si domina otra herramienta equivalente, explique esa transferencia y prepare una evidencia de adaptación. Por ejemplo, quien comprende HTTP, aserciones y ambientes en Postman puede aprender otra plataforma de API con mayor rapidez.

En el CV no escriba solamente nombres. Relacione cada tecnología con una tarea y un resultado que pueda defender.""",

    23: """Texto ampliado para leer:
En una entrevista pueden preguntarle por qué seleccionó esos localizadores. La respuesta no debe ser «porque lo vi en un tutorial». Puede explicar que los roles y etiquetas se acercan al comportamiento del usuario y suelen ser más resistentes que selectores dependientes de estilos internos.

También pueden preguntarle qué faltaría. Aquí faltan manejo seguro de credenciales, preparación de datos, limpieza, casos negativos, configuración de ambientes y quizá una abstracción para reutilizar acciones. Reconocer limitaciones demuestra criterio.

Si utiliza IA para producir una prueba, asegúrese de poder explicar cada línea, ejecutarla, modificarla y depurarla. El código que no puede defender no debe presentarse como dominio propio.""",

    24: """Texto ampliado para leer:
El pipeline representa una conversación automática entre el cambio y la calidad. Cuando alguien envía código, el sistema prepara un ambiente, ejecuta controles y publica una señal. Esa señal debe ser rápida y confiable. Si las pruebas fallan constantemente por razones aleatorias, el equipo comienza a ignorarlas.

Esto abre otra especialización: calidad en CI/CD, release engineering, DevOps y reliability. Una persona de QA que disfruta ambientes, automatización y análisis de fallos puede crecer hacia estas áreas. Necesitará aprender sistemas operativos, contenedores, nube, infraestructura como código, monitoreo y seguridad.

Para un portafolio junior basta comenzar con un workflow sencillo que ejecute la suite y conserve el reporte.""",

    25: """Texto ampliado para leer:
Este proyecto también puede adaptarse a la especialidad elegida. Para accesibilidad, incluya revisión con teclado, contraste, etiquetas y herramientas automáticas. Para rendimiento, agregue un escenario de carga, métricas y análisis. Para seguridad, documente un modelo de amenazas y pruebas autorizadas. Para datos, valide calidad, transformaciones y consistencia. Para IA, construya un pequeño conjunto de evaluación y criterios.

El proyecto debe ser honesto sobre su alcance. Explique qué cubre, qué no cubre y qué mejoraría con más tiempo. Esa transparencia comunica madurez.

Antes de publicarlo, pida a otra persona que siga el README. Si no puede ejecutar el proyecto, todavía falta documentación o reproducibilidad.""",

    29: """Texto ampliado para leer:
Esta estrategia evita dos extremos. El primero es estudiar sin mirar el mercado. El segundo es obedecer ciegamente cada requisito de una vacante. Analizar varias ofertas permite identificar patrones y distinguir entre requisitos esenciales, deseables y particulares de una empresa.

Busque vacantes del nivel actual y del nivel siguiente. Las primeras ayudan a aplicar ahora; las segundas muestran hacia dónde crecer. También compare vacantes locales, remotas y en inglés, porque pueden pedir combinaciones diferentes.

No espere cumplir el cien por ciento. Si posee una base sólida y puede cubrir una parte importante de las funciones, considere aplicar. La propia búsqueda ofrece retroalimentación para ajustar el estudio.""",

    30: """Texto ampliado para leer:
El método VACANTE puede repetirse cada trimestre. Visualizar evita estudiar sin destino. Analizar convierte las vacantes en fuente de información. Clasificar permite agrupar fundamentos, herramientas, habilidades humanas, dominio e idioma. Autoevaluar exige pruebas y no percepciones vagas.

La última fase, nivelar, trabajar y evidenciar, ocurre simultáneamente. No es necesario terminar toda la ruta antes de aplicar. Mientras se estudia una brecha, se puede mejorar el portafolio, practicar entrevistas y solicitar posiciones adecuadas.

Si después de varias semanas el mercado no responde, revise la hipótesis: quizá el nivel, el título, el CV, la evidencia o los canales necesitan ajuste. La autonomía incluye observar resultados y cambiar la estrategia.""",

    31: """Texto ampliado para leer:
Para construir esta matriz, copie únicamente información útil, no toda la descripción. Separe tareas de herramientas. Una tarea puede ser «diseñar pruebas de API» y la herramienta puede ser Postman. La tarea es más estable que la marca.

Utilice una escala de dominio basada en desempeño: cero, no lo conozco; uno, comprendo conceptos; dos, realizo ejercicios con guía; tres, ejecuto de manera autónoma; cuatro, puedo diseñar, explicar y ayudar a otra persona. Luego añada un enlace a la evidencia.

La prioridad puede calcularse combinando frecuencia, importancia para el puesto, nivel actual y tiempo requerido. Empiece por brechas que desbloqueen varias tareas, como fundamentos, API, SQL, Git o comunicación.""",

    33: """Texto ampliado para leer:
Una guía autónoma debe ser específica. Para cada competencia escriba un resultado de aprendizaje: «al finalizar podré diseñar una colección de API con autenticación y aserciones». Después seleccione una fuente principal, preferiblemente oficial, y una fuente de apoyo.

Diseñe una práctica que no sea copiar el ejemplo. Cambie datos, agregue casos negativos, provoque errores y explique los resultados. La evidencia puede ser código, reporte, video breve, documentación o una demostración.

Finalmente defina el criterio de terminado. Por ejemplo: ejecutar sin ayuda, explicar decisiones, resolver tres variaciones y recibir revisión. Sin criterio de dominio, estudiar puede convertirse en consumir contenido indefinidamente.""",

    34: """Texto ampliado para leer:
Esta tabla puede convertirse en el centro del plan semanal. La columna «brecha» proviene de las vacantes. «Resultado esperado» expresa una capacidad laboral. «Práctica» define lo que hará la persona. «Evidencia» indica qué quedará disponible y «fecha» crea un límite.

Añada dos columnas si lo necesita: fuente principal y estado. Los estados pueden ser pendiente, en aprendizaje, en práctica, demostrable y revisado. No marque una competencia como completa porque terminó un video. Márquela como demostrable cuando pueda realizar una tarea y explicar su trabajo.

Cada viernes revise la guía. Mantenga lo que funciona, reduzca actividades demasiado grandes y reemplace recursos que no aportan claridad.""",

    35: """Texto ampliado para leer:
Veamos un ejemplo de práctica deliberada. Si falla al escribir aserciones de API, no necesita repetir un curso completo. Puede tomar cinco respuestas distintas y practicar solamente códigos, encabezados, esquemas y campos obligatorios. Después integra esa habilidad al proyecto.

La recuperación activa puede hacerse explicando un concepto sin mirar notas, resolviendo una tarea en blanco o respondiendo preguntas. La bitácora debe registrar errores, porque allí ocurre gran parte del aprendizaje.

Ser autónomo también significa saber cuándo pedir ayuda. Primero describa lo que intentó, el resultado obtenido, el resultado esperado y la evidencia. Esa forma de solicitar apoyo mejora la comunicación profesional y facilita recibir respuestas útiles.""",

    37: """Texto ampliado para leer:
Utilice búsquedas en español e inglés porque los títulos cambian. Pruebe combinaciones como «QA junior», «software tester», «quality analyst», «manual tester», «QA automation», «test automation engineer», «SDET», «API tester» y «quality engineer». Añada ubicación, remoto, internship o junior.

No dependa únicamente del botón de aplicación rápida. Visite las páginas de las empresas, siga reclutadores especializados, participe en comunidades y solicite referidos de manera respetuosa. Un referido no significa pedir que alguien mienta; significa que una persona conozca su trabajo y pueda señalar una oportunidad.

Registre resultados semanalmente. La búsqueda también se prueba, se mide y se mejora como cualquier proceso de calidad.""",

    38: """Texto ampliado para leer:
El CV debe responder rápidamente tres preguntas: qué puesto busca, qué puede hacer y dónde se observa la evidencia. Utilice palabras de la vacante cuando sean verdaderas. Si la oferta habla de pruebas de API y usted posee un proyecto, ubíquelo de manera visible.

Para quienes cambian de carrera, traduzca la experiencia previa. Atención al cliente puede demostrar análisis de incidentes y comunicación. Docencia puede mostrar planificación, documentación y facilitación. Desarrollo demuestra código y depuración. Administración puede aportar procesos y detalle.

No oculte que un proyecto es académico o personal. Preséntelo con profesionalidad: problema, alcance, acciones, herramientas y resultados. Revise ortografía, enlaces y consistencia de fechas.""",

    39: """Texto ampliado para leer:
Un buen titular de LinkedIn puede decir «QA Analyst | Pruebas funcionales, API y SQL | En formación en automatización». Es más informativo que escribir únicamente «buscando oportunidades». La sección Acerca de puede explicar el tipo de producto que le interesa, sus capacidades y la dirección de crecimiento.

En GitHub, fije los repositorios más relevantes. Un reclutador no debería adivinar cuál revisar. Agregue capturas, reportes o una demostración breve cuando ayuden, pero no publique secretos ni datos reales.

La actividad pública también cuenta: documentación, issues, pull requests y colaboración muestran cómo trabaja con otras personas. Calidad no es solamente escribir scripts; es contribuir de forma clara y confiable.""",

    40: """Texto ampliado para leer:
Prepare una presentación profesional de uno o dos minutos: quién es, qué experiencia trae, qué rol busca, qué proyecto demuestra su preparación y por qué le interesa la empresa. Debe sonar natural, no memorizada mecánicamente.

Practique preguntas como: ¿cómo probaría un inicio de sesión?, ¿qué diferencia existe entre severidad y prioridad?, ¿cómo selecciona una regresión?, ¿cómo investigaría una falla intermitente?, ¿qué automatizaría primero? y ¿cómo reportaría un defecto rechazado por desarrollo?

Para especialidades técnicas, espere ejercicios relacionados: una consulta SQL, una colección de API, revisión de código o diseño de framework. Piense en voz alta, declare supuestos y priorice. El proceso de razonamiento suele importar tanto como la respuesta final.""",

    41: """Texto ampliado para leer:
Construya un radar con pocas fuentes de alta calidad. Para herramientas, siga documentación y notas de versiones. Para fundamentos, utilice programas reconocidos y libros. Para seguridad, revise OWASP. Para tendencias laborales, analice vacantes y reportes. Para comunidad, participe en eventos y conversaciones profesionales.

No cambie de herramienta cada semana. Cuando aparezca una novedad, pregunte: ¿resuelve un problema de mi rol?, ¿aparece en mis vacantes?, ¿mejora mi proyecto?, ¿puedo probarla en poco tiempo? Si la respuesta es no, guárdela en una lista futura.

Mantenerse actualizado es seleccionar, experimentar y decidir; no consumir todas las noticias. Su ruta profesional necesita estabilidad suficiente para profundizar y flexibilidad suficiente para adaptarse.""",

    42: """Texto ampliado para leer:
El plan de treinta días debe combinar aprendizaje y empleabilidad. No coloque cuatro semanas únicamente de cursos. Incluya análisis de vacantes, práctica, publicación de evidencia, adaptación del CV, simulación de entrevista y aplicaciones.

Un plan posible sería: semana uno, analizar vacantes y completar API; semana dos, SQL y documentación; semana tres, automatizar un flujo; semana cuatro, integrar CI, mejorar README y realizar aplicaciones. Cada viernes se revisan resultados.

Defina una versión mínima. Si dispone de cinco horas semanales, reduzca el alcance, no abandone el sistema. La meta no es terminar una lista perfecta; es producir progreso visible y recibir señales del mercado que permitan mejorar la siguiente versión.""",

    50: """Texto ampliado para leer:
Antes de despedirnos, recuerden que una carrera tecnológica no se construye intentando predecir cada cambio. Se construye desarrollando fundamentos, observando el entorno y aprendiendo a responder. La persona que sabe investigar una herramienta, practicarla, verificar resultados y comunicar lo aprendido puede adaptarse mejor que quien depende de una lista fija.

QA puede conducir a muchas rutas: análisis, automatización, SDET, rendimiento, seguridad, accesibilidad, datos, IA, DevOps, producto, negocio y liderazgo. No tienen que decidir toda su carrera hoy. Solo necesitan seleccionar un próximo paso, construir evidencia y revisarlo con honestidad.

Mi invitación final es sencilla: esta semana busquen vacantes, actualicen su matriz y comiencen una evidencia pequeña. La empleabilidad se fortalece cuando el aprendizaje se convierte en trabajo visible.""",
}

for slide_number, expansion in NOTE_EXPANSIONS.items():
    notes_frame = prs.slides[slide_number - 1].notes_slide.notes_text_frame
    if notes_frame is not None:
        notes_frame.text = notes_frame.text.rstrip() + "\n\n" + expansion.strip()


def replace_exact_text(slide, old_value, new_value):
    """Replace a one-run text box while preserving its formatting."""
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        if shape.text.strip() != old_value:
            continue
        paragraph = shape.text_frame.paragraphs[0]
        if paragraph.runs:
            paragraph.runs[0].text = new_value
        else:
            paragraph.text = new_value
        return


def note_body(slide):
    frame = slide.notes_slide.notes_text_frame
    if frame is None:
        return ""
    return re.sub(r"^\[Tiempo sugerido: \d+ minutos\]\s*", "", frame.text.strip())


# Unir el contenido de diapositivas complementarias antes de compactar.
replace_exact_text(prs.slides[1], "¿Qué te llevarás de este taller?", "Resultados y ruta de las tres horas")
replace_exact_text(prs.slides[22], "Ejemplo mínimo con Playwright", "Playwright y CI: del código al pipeline")
replace_exact_text(prs.slides[37], "CV adaptado al puesto, no biografía completa", "CV, LinkedIn y GitHub alineados con la vacante")
replace_exact_text(prs.slides[40], "Sistema para mantenerte actualizado", "Actualización continua y plan de 30 días")

MERGED_NOTES = {
    2: [3],       # resultados + agenda
    23: [24],     # Playwright + integración continua
    38: [39],     # CV + LinkedIn y GitHub
    41: [42],     # actualización + plan de 30 días
}
for target_number, source_numbers in MERGED_NOTES.items():
    target_frame = prs.slides[target_number - 1].notes_slide.notes_text_frame
    if target_frame is None:
        continue
    for source_number in source_numbers:
        target_frame.text = (
            target_frame.text.rstrip()
            + "\n\nContenido integrado de la versión extensa:\n"
            + note_body(prs.slides[source_number - 1])
        )


# Orden final compacto: 40 diapositivas. Los números corresponden a la versión
# extensa antes de compactar; 48-50 contienen sorteo, formulario y cierre.
KEEP_ORDER = [
    1, 2, 4, 6, 7, 9, 10, 46, 11, 13,
    14, 16, 17, 20, 21, 44, 45, 23, 51, 27,
    29, 31, 33, 35, 37, 57, 58, 59, 52, 53, 54, 56, 60, 55,
    38,
    40, 41, 48, 49, 50,
]
slide_id_list = prs.slides._sldIdLst
all_slide_ids = list(slide_id_list)
kept_ids = [all_slide_ids[index - 1] for index in KEEP_ORDER]
kept_set = set(KEEP_ORDER)
for index, slide_id in enumerate(all_slide_ids, start=1):
    slide_id_list.remove(slide_id)
    if index not in kept_set:
        prs.part.drop_rel(slide_id.rId)
for slide_id in kept_ids:
    slide_id_list.append(slide_id)


# Renumerar los indicadores visuales después de reducir y reordenar.
for new_number, slide in enumerate(prs.slides, start=1):
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        value = shape.text.strip()
        if not value.isdigit():
            continue
        x = shape.left / 914400
        y = shape.top / 914400
        new_value = None
        if x < 1.20 and y < 0.85:
            new_value = f"{new_number:02d}"
        elif x > 11.40 and y > 6.75:
            new_value = f"{new_number:02d}"
        if new_value is None:
            continue
        paragraph = shape.text_frame.paragraphs[0]
        if paragraph.runs:
            paragraph.runs[0].text = new_value
        else:
            paragraph.text = new_value


# Distribución exacta de los 180 minutos, incluida una pausa de 10 minutos.
TIMINGS = [
    3, 4, 4, 4, 4, 5, 4, 4, 4, 4,
    4, 4, 4, 4, 4, 4, 4, 4, 4, 10,
    5, 7, 7, 7, 6, 3, 3, 4, 4, 4,
    4, 3, 4, 4, 7,
    6, 4, 4, 4, 4,
]
assert len(TIMINGS) == len(prs.slides)
assert sum(TIMINGS) == 180
for slide, minutes in zip(prs.slides, TIMINGS):
    notes_frame = slide.notes_slide.notes_text_frame
    if notes_frame is not None:
        notes_frame.text = re.sub(
            r"\[Tiempo sugerido: \d+ minutos\]",
            f"[Tiempo sugerido: {minutes} minutos]",
            notes_frame.text,
            count=1,
        )

prs.core_properties.title = "La transformación del Software QA Engineering: IA, automatización y empleabilidad"
prs.core_properties.subject = "Charla-taller virtual del 2 de julio de 2026"
prs.core_properties.author = "Marlenis Judith Concepción Cuevas"
prs.core_properties.keywords = "QA, IA, automatización, empleabilidad, aprendizaje autónomo, UASD"
prs.core_properties.comments = "Incluye notas del orador en cada diapositiva."

prs.save(OUT)
print(f"Presentación creada: {OUT}")
print(f"Diapositivas: {len(prs.slides)}")
